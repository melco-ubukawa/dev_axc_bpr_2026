# utils/file_processor.py

import io
import os
import re
import json
import base64
import logging
import subprocess
import magic # python-magic
import pptx # python-pptx
import docx # python-docx
import chardet # chardet
import pandas as pd
import fitz # PyMuPDF
from PIL import Image, ImageStat
from bs4 import BeautifulSoup
from google.genai import types
from extract_msg import Message
from .gemini import generate_with_gemini, ensure_parent_log_blob_name
from google.genai import types
import config

# --- ロガー設定 ---
logger = logging.getLogger(__name__)

# --- 定数 ---
MAX_FILE_SIZE_BYTES_SERVER = 100 * 1024 * 1024  # 100MB

# --- ソースコード最適化関数 ---
def _optimize_source_code(text: str, filename: str) -> str:
    """
    ソースコードやテキストファイルから、LLMの理解に不要かつトークン消費が激しい要素
    （埋め込みBase64データ、巨大な1行データ、ソースマップ、SVGパス、ライセンス条文など）
    を安全に削除・短縮する。
    """
    if not text:
        return ""

    # 安全対策: プレーンテキストやMarkdown、CSVなどは、誤爆リスクが高いため
    # 構造的な削除（Base64削除など）以外はスキップする
    ext = os.path.splitext(filename)[1].lower()
    is_code_file = ext in ['.js', '.jsx', '.ts', '.tsx', '.py', '.java', '.c', '.cpp', '.h', '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.html', '.css', '.xml', '.vue', '.json']

    original_len = len(text)
    
    # ---------------------------------------------------------
    # 1. Base64 / Hex データの削除 (全ファイル共通)
    # ---------------------------------------------------------
    # data URI スキーム形式のもの（明確なので安全に削除）
    text = re.sub(
        r'data:[a-z-]+/[a-z0-9-.+]+;base64,[a-zA-Z0-9+/=\s]{20,}',
        r'[...Base64 Data Omitted...]',
        text
    )

    # 汎用的な長い英数字の羅列 (誤爆防止のため閾値を1000文字に上げ、引用符で囲まれているものに限定)
    # コードファイルの場合のみ適用（txtファイル等での誤爆防止）
    if is_code_file:
        text = re.sub(
            r'([\'"])[a-zA-Z0-9+/=]{1000,}([\'"])', 
            r'\1[...Large Data String Omitted...]\2', 
            text
        )

    # ---------------------------------------------------------
    # 2. SVG / Path データの削除 (HTML/XML/JS内)
    # ---------------------------------------------------------
    # if ext in ['.html', '.xml', '.js', '.jsx', '.ts', '.tsx', '.vue', '.svg']:
    #     # 修正: \b (単語境界) を追加して、id="..." や data-id="..." への誤爆を防止
    #     # また、属性値が非常に長い場合のみ対象とする (100文字以上)
    #     text = re.sub(
    #         r'\b(d|points)=([\'"])([\s\S]{100,}?)([\'"])', 
    #         r'\1=\2[...SVG Path Omitted...]\4', 
    #         text
    #     )

    # ---------------------------------------------------------
    # 3. ソースマップと巨大コメントの削除 (コードファイルのみ)
    # ---------------------------------------------------------
    if is_code_file:
        # Source Map
        text = re.sub(r'/\*# sourceMappingURL=.*?\*/', '', text)
        text = re.sub(r'//# sourceMappingURL=.*', '', text)
        
        # ライセンスヘッダーなどの巨大なブロックコメント削除
        # 修正: 正規表現で一括置換すると、"/* code */ code /* License */" のように
        # コメント間のコードを巻き込んで消してしまうリスクがあるため、
        # ブロックコメントを一つずつ抽出して判定する方式に変更。
        
        def comment_replacer(match):
            comment_content = match.group(0)
            # 200文字以上 かつ "License" や "Copyright" が含まれる場合のみ省略
            if len(comment_content) > 200 and re.search(r'(License|Copyright)', comment_content, re.IGNORECASE):
                return '/* ...License/Copyright Info Omitted... */'
            return comment_content

        # 非貪欲マッチでブロックコメントを抽出
        text = re.sub(r'/\*[\s\S]*?\*/', comment_replacer, text)

    # ---------------------------------------------------------
    # 4. 巨大な数値/データ配列の短縮 (行単位処理)
    # ---------------------------------------------------------
    # 行単位で処理を行い、構文破壊を防ぐ
    
    lines = text.split('\n')
    optimized_lines = []
    
    for line in lines:
        stripped = line.strip()
        line_len = len(line)

        # A. 極端に長い行 (Minified Code対策)
        # 構文を壊さないよう、中抜きではなく「行ごと省略」または「データ配列のみ省略」にする
        if line_len > 1000: # 閾値を500->1000へ安全側に変更
            
            # ケース1: データ配列 ([1,2,3,...] や "data", "data"...)
            # 数字、引用符、カンマ、スペースだけで構成されている割合が高い場合
            if re.match(r'^[\d\s,.\-\'\"\[\]{}]+$', stripped) or (stripped.count(',') > 50 and stripped.count(' ') < stripped.count(',') * 2):
                 optimized_lines.append(f" ... [Long Data Row Omitted: {line_len} chars] ... ")
                 continue
            
            # ケース2: Minified Code (スペースが極端に少ない)
            space_ratio = line.count(' ') / line_len
            if space_ratio < 0.02: # 2%未満はほぼ確実にMinified
                # Minifiedコードを中抜きすると構文エラーになるため、行全体を省略する方が安全
                optimized_lines.append(f"// ... [Minified Code Line Omitted: {line_len} chars] ... ")
                continue

        optimized_lines.append(line)

    text = '\n'.join(optimized_lines)

    # ---------------------------------------------------------
    # 5. 余分な空白・改行の圧縮
    # ---------------------------------------------------------
    # 3つ以上の連続する改行を2つに
    text = re.sub(r'\n{3,}', '\n\n', text)

    if len(text) < original_len:
        logger.info(f"Optimized source code '{filename}': {original_len} -> {len(text)} chars ({(len(text)/original_len)*100:.1f}%)")
    else:
        logger.debug(f"No optimization applied to '{filename}'")

    return text

# --- ファイル内容抽出ヘルパー関数 ---
def extract_file_content_to_part(uploaded_file) -> types.Part | None:
    """
    アップロードされたファイルの内容を読み取り、MIMEタイプと拡張子に応じて
    google.generativeai.types.Part オブジェクト、またはエラー時にNoneを返す。
    テキストデータの場合は、AIがファイル名を認識できるようにヘッダーを付与する。
    """
    if not uploaded_file:
        return None

    filename = uploaded_file.filename
    logger.info(f"Processing file: {filename}")

    # テキストパート生成用の内部ヘルパー関数
    def create_text_part_with_filename(text_content, fname):

        optimized_text = _optimize_source_code(text_content, fname)

        # ファイル名を明記してAIにコンテキストを伝える
        contextualized_text = f"--- File: {fname} ---\n{optimized_text}\n--- End of File ---"
        return types.Part.from_text(text=contextualized_text)

    try:
        file_content = uploaded_file.read()
        if len(file_content) > MAX_FILE_SIZE_BYTES_SERVER:
            logger.error(f"File size exceeds limit: {filename}")
            return None
        uploaded_file.seek(0)

        try:
            mime_type = magic.from_buffer(file_content, mime=True)
        except Exception as magic_err:
            logger.warning(f"python-magic failed: {magic_err}. Skipping.")
            mime_type = "application/octet-stream" 

        logger.info(f"Detected MIME type for {filename}: {mime_type}")

        if mime_type.startswith("video/"):
            #logger.error(f"Video file upload is not supported here: {filename}")
            #return None
            logger.info(f"Processed '{filename}' as an video file.")
            # 動画ファイルはバイナリなのでファイル名付与はプロンプト側で行う必要があるが、
            # ここではPart生成の責務としてバイナリを返す
            return types.Part.from_bytes(data=file_content, mime_type=mime_type)

        if mime_type.startswith("audio/"):
            logger.info(f"Processed '{filename}' as an audio file.")
            # 音声ファイルはバイナリなのでファイル名付与はプロンプト側で行う必要があるが、
            # ここではPart生成の責務としてバイナリを返す
            return types.Part.from_bytes(data=file_content, mime_type=mime_type)

        if filename.lower().endswith('.msg'):
            try:
                msg = Message(file_content)
                body_text = BeautifulSoup(msg.htmlBody, 'html.parser').get_text(separator='\n', strip=True) if msg.htmlBody else msg.body
                extracted_text = f"From: {getattr(msg, 'sender', 'N/A')}\nTo: {getattr(msg, 'to', 'N/A')}\nCC: {getattr(msg, 'cc', 'N/A')}\nSubject: {getattr(msg, 'subject', 'N/A')}\nDate: {getattr(msg, 'date', 'N/A')}\n--- Body ---\n{body_text.strip()}"
                # .msgは内部ですでに構造化しているが、統一フォーマットでラップする
                return create_text_part_with_filename(extracted_text, filename)
            except Exception as msg_err:
                logger.error(f"Error reading MSG file '{filename}': {msg_err}")
                return None

        if filename.lower().endswith('.pdf'):
            return types.Part.from_bytes(data=file_content, mime_type="application/pdf")

        if filename.lower().endswith('.docx'):
            try:
                doc = docx.Document(io.BytesIO(file_content))
                all_text = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        all_text.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        all_text.append(" | ".join(row_text))
                extracted_text = "\n".join(all_text)
                
                MAX_CHARS_FROM_DOC = 700000 
                if len(extracted_text) > MAX_CHARS_FROM_DOC:
                    logger.warning(f"Extracted text from DOCX '{filename}' is too long ({len(extracted_text)} chars). Truncating.")
                    extracted_text = extracted_text[:MAX_CHARS_FROM_DOC] + "\n\n[...文書が長すぎるため、ここで内容を切り捨てました...]"
                
                return create_text_part_with_filename(extracted_text, filename)
            except Exception as docx_err:
                logger.error(f"Error reading DOCX file '{filename}': {docx_err}")
                return None
        
        if filename.lower().endswith('.pptx'):
            try:
                presentation = pptx.Presentation(io.BytesIO(file_content))
                full_text = [shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                return create_text_part_with_filename("\n".join(full_text), filename)
            except Exception as pptx_err:
                logger.error(f"Error reading PPTX file {filename}: {pptx_err}")
                return None

        if filename.lower().endswith(('.xlsx', '.xls')):
            try:
                xls_dfs = pd.read_excel(io.BytesIO(file_content), sheet_name=None)
                full_text_content = []
                for sheet_name, df_sheet in xls_dfs.items():
                    full_text_content.append(f"--- Sheet: {sheet_name} ---\n")
                    full_text_content.append(df_sheet.to_csv(index=False))
                    full_text_content.append("\n")
                
                extracted_text = "\n".join(full_text_content)
                MAX_CHARS_FROM_TABLE = 700000 
                if len(extracted_text) > MAX_CHARS_FROM_TABLE:
                    logger.warning(f"Extracted text from Excel '{filename}' is too long. Truncating.")
                    extracted_text = extracted_text[:MAX_CHARS_FROM_TABLE] + "\n\n[...内容が長すぎるため、ここで切り捨てました...]"
                
                return create_text_part_with_filename(extracted_text, filename)
            except Exception as excel_err:
                logger.error(f"Error reading Excel file '{filename}': {excel_err}")
                return None

        if filename.lower().endswith('.csv'):
            potential_encodings = ['utf-8', 'utf-8-sig', 'cp932', 'euc-jp']
            try:
                detected_encoding_info = chardet.detect(file_content)
                detected_encoding = detected_encoding_info.get('encoding')
                if detected_encoding and detected_encoding.lower() not in [e.lower() for e in potential_encodings]:
                    potential_encodings.insert(0, detected_encoding)
            except Exception as e_chardet:
                logger.warning(f"Chardet failed for '{filename}': {e_chardet}")
            
            for encoding in potential_encodings:
                try:
                    df_csv = pd.read_csv(io.BytesIO(file_content), encoding=encoding, sep=None, engine='python')
                    extracted_text = df_csv.to_csv(index=False)
                    MAX_CHARS_FROM_TABLE = 700000
                    if len(extracted_text) > MAX_CHARS_FROM_TABLE:
                        logger.warning(f"Extracted text from CSV '{filename}' is too long. Truncating.")
                        extracted_text = extracted_text[:MAX_CHARS_FROM_TABLE] + "\n\n[...内容が長すぎるため、ここで切り捨てました...]"
                    return create_text_part_with_filename(extracted_text, filename)
                except (UnicodeDecodeError, pd.errors.ParserError):
                    continue
                except Exception as e:
                    logger.error(f"Unexpected error reading CSV '{filename}' with encoding '{encoding}': {e}", exc_info=True)
                    break 
            logger.error(f"Failed to read CSV file '{filename}' with all attempted encodings.")
            return None
            
        # SVGファイルを特別にテキストとして扱う
        if filename.lower().endswith('.svg') or mime_type == "image/svg+xml":
            try:
                encoding = chardet.detect(file_content).get("encoding") or "utf-8"
                extracted_text = file_content.decode(encoding, errors="replace")
                logger.info(f"Processed '{filename}' as a text-based SVG file.")
                return create_text_part_with_filename(extracted_text, filename)
            except Exception as svg_err:
                logger.error(f"Error reading SVG file '{filename}' as text: {svg_err}")
        
        # 画像ファイル (SVG以外)
        if mime_type.startswith("image/"):
            try:
                with Image.open(io.BytesIO(file_content)):
                    return types.Part.from_bytes(data=file_content, mime_type=mime_type)
            except Exception as img_err:
                logger.warning(f"File {filename} (MIME: {mime_type}) could not be opened as an image, will be treated as text. Error: {img_err}")

        # MIMEタイプに基づく汎用的なテキストファイルの処理 (Pythonファイルなどはここに来る可能性が高い)
        if mime_type.startswith("text/") or mime_type in ["application/json", "application/xml", "application/javascript", "application/x-python-code", "text/x-python"]:
            try:
                encoding = chardet.detect(file_content).get("encoding") or "utf-8"
                logger.info(f"Processed '{filename}' as a generic text file with MIME type {mime_type}.")
                extracted_text = file_content.decode(encoding, errors="replace")
                return create_text_part_with_filename(extracted_text, filename)
            except (UnicodeDecodeError, LookupError) as decode_err:
                 logger.warning(f"Could not decode '{filename}' with detected encoding, will use final fallback. Error: {decode_err}")

        # 最終フォールバック: 上記のいずれにも当てはまらないファイルをテキストとしてデコード試行
        # (python-magicが text/x-python などを返さず application/octet-stream を返した場合など)
        try:
            encoding = chardet.detect(file_content).get("encoding") or "utf-8"
            logger.info(f"Processed '{filename}' using final fallback text handler.")
            extracted_text = file_content.decode(encoding, errors="replace")
            return create_text_part_with_filename(extracted_text, filename)
        except Exception as final_err:
             logger.error(f"Failed to decode '{filename}' with any method: {final_err}")
             return None

    except Exception as e:
        logger.error(f"General error processing file {filename}: {e}", exc_info=True)
        return None

def extract_images_from_pdf(pdf_stream) -> list[str]:
    """PDFファイルストリームから画像を抽出し、Base64エンコードされた文字列のリストを返す。"""
    extracted_images_base64 = []
    try:
        pdf_stream.seek(0)
        doc = fitz.open(stream=pdf_stream.read(), filetype="pdf")
        img_count = 0
        MAX_IMAGES_TO_EXTRACT = 100 
        
        MIN_EMBEDDED_IMAGE_WIDTH = 20
        MIN_EMBEDDED_IMAGE_HEIGHT = 20
        MIN_EMBEDDED_IMAGE_BYTES = 128

        MIN_RASTERIZED_PAGE_WIDTH = 50 
        MIN_RASTERIZED_PAGE_HEIGHT = 50
        MIN_RASTERIZED_PAGE_BYTES = 1536 
        RASTERIZE_DPI = 150
        
        TEXT_HEAVY_THRESHOLD = 500 
        SIGNIFICANT_DRAWING_PATH_COUNT_THRESHOLD = 3 
        IMAGE_COMPLEXITY_STDDEV_THRESHOLD = 4.0
        TABLE_PAGE_COMPLEXITY_STDDEV_THRESHOLD = 1.5

        logger.info(f"Starting image extraction from PDF...")

        for page_num in range(len(doc)):
            if img_count >= MAX_IMAGES_TO_EXTRACT:
                logger.warning(f"Reached max image extraction limit ({MAX_IMAGES_TO_EXTRACT}). Stopping.")
                break
            
            page = doc.load_page(page_num)
            page_has_actual_content_for_rasterize = False 
            
            image_list = page.get_images(full=True)
            if image_list:
                for img_info in image_list:
                    if img_count >= MAX_IMAGES_TO_EXTRACT: break
                    xref = img_info[0]
                    try:
                        base_image_info = doc.extract_image(xref)
                        if not base_image_info: continue
                        image_bytes, image_ext, width, height = base_image_info["image"], base_image_info["ext"].lower(), base_image_info.get("width",0), base_image_info.get("height",0)

                        if len(image_bytes) < MIN_EMBEDDED_IMAGE_BYTES or width < MIN_EMBEDDED_IMAGE_WIDTH or height < MIN_EMBEDDED_IMAGE_HEIGHT:
                            continue
                        
                        supported_extensions = ['png', 'jpeg', 'jpg', 'gif', 'bmp', 'tiff', 'webp', 'jp2', 'jpx']
                        if image_ext in ['jpeg', 'jpg']:
                            try:
                                img_pil = Image.open(io.BytesIO(image_bytes))
                                if img_pil.mode == 'CMYK':
                                    rgb_img = img_pil.convert('RGB'); output_buffer = io.BytesIO(); rgb_img.save(output_buffer, format="PNG"); image_bytes = output_buffer.getvalue(); image_ext = "png"
                            except: pass
                        if image_ext not in supported_extensions: continue

                        mime_type = f"image/{image_ext}"
                        if image_ext in ["jp2", "jpx"]: mime_type = "image/jp2"
                        elif image_ext in ["jpeg", "jpg"]: mime_type = "image/jpeg"
                        elif image_ext == "tiff": mime_type = "image/tiff"
                        
                        extracted_images_base64.append(f"data:{mime_type};base64,{base64.b64encode(image_bytes).decode('utf-8')}")
                        img_count += 1
                        page_has_actual_content_for_rasterize = True 
                        logger.info(f"  Extracted embedded image {img_count} from page {page_num + 1}.")
                    except Exception: pass
            
            page_text_content = page.get_text("text", sort=True).strip()
            page_text_len = len(page_text_content)
            drawing_paths = page.get_drawings()
            num_drawing_paths = len(drawing_paths) if drawing_paths else 0

            contains_table_keyword = False
            if re.search(r'\\b表\\s*\\d*[\\d\\.]*\\d*|\\btable\\b', page_text_content, re.IGNORECASE):
                contains_table_keyword = True
                page_has_actual_content_for_rasterize = True
            
            is_likely_text_only_page = page_text_len > TEXT_HEAVY_THRESHOLD and num_drawing_paths < SIGNIFICANT_DRAWING_PATH_COUNT_THRESHOLD and not page_has_actual_content_for_rasterize and not contains_table_keyword

            should_rasterize_page = False
            if img_count < MAX_IMAGES_TO_EXTRACT and not is_likely_text_only_page:
                if contains_table_keyword:
                    should_rasterize_page = True
                elif page_text_len < (TEXT_HEAVY_THRESHOLD / 2) and num_drawing_paths > 0:
                    should_rasterize_page = True
                    page_has_actual_content_for_rasterize = True
                elif page_text_len < TEXT_HEAVY_THRESHOLD and num_drawing_paths >= (SIGNIFICANT_DRAWING_PATH_COUNT_THRESHOLD * 2): 
                    should_rasterize_page = True
                    page_has_actual_content_for_rasterize = True
                elif not page_has_actual_content_for_rasterize and page_text_len < (TEXT_HEAVY_THRESHOLD / 5) and num_drawing_paths < SIGNIFICANT_DRAWING_PATH_COUNT_THRESHOLD and num_drawing_paths > 0:
                    should_rasterize_page = True
                    page_has_actual_content_for_rasterize = True
            
            if should_rasterize_page and page_has_actual_content_for_rasterize:
                try:
                    zoom_matrix = fitz.Matrix(RASTERIZE_DPI/72, RASTERIZE_DPI/72)
                    pix = page.get_pixmap(matrix=zoom_matrix, alpha=False) 
                    
                    if pix.width < MIN_RASTERIZED_PAGE_WIDTH or pix.height < MIN_RASTERIZED_PAGE_HEIGHT or pix.n < 3:
                        continue 

                    image_bytes_raster = pix.tobytes("png")
                    if len(image_bytes_raster) < MIN_RASTERIZED_PAGE_BYTES:
                         continue

                    current_complexity_threshold = TABLE_PAGE_COMPLEXITY_STDDEV_THRESHOLD if contains_table_keyword else IMAGE_COMPLEXITY_STDDEV_THRESHOLD
                    
                    try:
                        pil_img_raster = Image.open(io.BytesIO(image_bytes_raster)).convert('L')
                        stat = ImageStat.Stat(pil_img_raster)
                        stddev = stat.stddev[0]
                        if stddev < current_complexity_threshold:
                            continue
                    except Exception as e_pil_stat:
                        logger.warning(f"    Could not assess complexity for rasterized page {page_num + 1}: {e_pil_stat}. Proceeding.")

                    extracted_images_base64.append(f"data:image/png;base64,{base64.b64encode(image_bytes_raster).decode('utf-8')}")
                    img_count += 1
                    logger.info(f"  Rendered page {page_num + 1} as image {img_count}.")
                except Exception as e_render:
                    logger.error(f"  Error rendering page {page_num + 1}: {e_render}", exc_info=False)
        
        logger.info(f"Finished PDF processing. Total {len(extracted_images_base64)} images extracted.")
    except fitz.fitz.FitzError as fe:
        logger.error(f"PyMuPDF (Fitz) error: {fe}", exc_info=True)
    except Exception as e_doc:
        logger.error(f"General PDF processing error: {e_doc}", exc_info=True)
    return extracted_images_base64

def _read_dataframe_with_multiple_encodings(file_input, filename: str = None) -> pd.DataFrame | None:
    """
    指定された入力（パスまたはBytesIO）から、複数のエンコーディングを試行してDataFrameを読み込む。
    
    :param file_input: ファイルパス(str) または io.BytesIOオブジェクト
    :param filename: ファイル名（file_inputがBytesIOの場合に形式判別のために使用）
    """
    is_stream = isinstance(file_input, (io.BytesIO, io.IOBase))
    
    # 1. 形式判別用の名前を特定
    target_name = filename if is_stream else file_input
    if not target_name:
        logger.error("DataFrame読み込みエラー: ファイル名またはパスが指定されていません。")
        return None

    # 2. 存在チェックと生データの取得（エンコーディング判定用）
    try:
        if is_stream:
            file_input.seek(0)
            raw_data = file_input.read()
            file_input.seek(0) # ポインタを戻す
        else:
            if not os.path.exists(file_input):
                logger.error(f"DataFrame読み込みエラー: ファイルが存在しません - {file_input}")
                return None
            with open(file_input, 'rb') as f:
                raw_data = f.read()
    except Exception as e:
        logger.error(f"データの読み込みに失敗しました: {e}")
        return None

    # 3. エンコーディングの推定
    potential_encodings = ['utf-8', 'utf-8-sig', 'cp932', 'euc-jp']
    try:
        detected_encoding_info = chardet.detect(raw_data)
        detected_encoding = detected_encoding_info.get('encoding')
        if detected_encoding and detected_encoding.lower() not in [e.lower() for e in potential_encodings]:
            potential_encodings.insert(0, detected_encoding)
            logger.info(f"Chardetが '{detected_encoding}' を推定しました。")
    except Exception as e_chardet:
        logger.warning(f"Chardetでのエンコーディング推定に失敗: {e_chardet}。")

    # 4. 読み込み試行
    last_error = None
    for encoding in potential_encodings:
        try:
            # ストリームの場合は毎回ポインタを先頭に戻す
            if is_stream:
                file_input.seek(0)
            
            # 読み込み対象（パスまたはストリーム）
            read_source = file_input

            if target_name.lower().endswith('.csv'):
                df = pd.read_csv(read_source, encoding=encoding, sep=None, engine='python')
            elif target_name.lower().endswith(('.xls', '.xlsx')):
                df = pd.read_excel(read_source, engine='openpyxl')
            else:
                logger.warning(f"サポートされていないファイル形式: {target_name}")
                return None
            
            logger.info(f"DataFrameの読み込みに成功: '{os.path.basename(target_name) if not is_stream else target_name}' (Encoding: {encoding})")
            return df
        except (UnicodeDecodeError, pd.errors.ParserError) as e:
            last_error = e
            continue
        except Exception as e:
            logger.error(f"DataFrame読み込み中に予期しないエラー: {e}", exc_info=True)
            last_error = e
            break

    logger.error(f"すべてのエンコーディングで読み込みに失敗: {target_name}. 最終エラー: {last_error}")
    return None

def _get_text_from_file_for_diff(uploaded_file):
    """ファイル比較機能専用のテキスト抽出ヘルパー。"""
    filename = uploaded_file.filename.lower()
    file_bytes = uploaded_file.read()
    uploaded_file.seek(0)

    try:
        if filename.endswith('.pdf'):
            with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                text = "".join(page.get_text() for page in doc)
            return text
        elif filename.endswith('.docx'):
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join([p.text for p in doc.paragraphs])
        elif filename.endswith('.pptx'):
            prs = pptx.Presentation(io.BytesIO(file_bytes))
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)
        elif filename.endswith(('.xls', '.xlsx')):
            df_dict = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None, header=0)
            all_text = []
            sheets = df_dict if isinstance(df_dict, dict) else {"Sheet1": df_dict}
            for sheet_name, sheet_df in sheets.items():
                all_text.append(f"--- Sheet: {sheet_name} ---")
                sheet_df = sheet_df.reset_index(drop=True)
                sheet_df.index = range(1, len(sheet_df) + 1)
                sheet_df.index.name = "行番号"
                all_text.append(sheet_df.to_string(index=True))
            return "\n".join(all_text)
        elif filename.endswith('.csv'):
            encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'
            lines = file_bytes.decode(encoding, errors='replace').splitlines()
            # ヘッダー行は行番号なし、データ行は1始まりで行番号を付与
            numbered = []
            for i, line in enumerate(lines):
                if i == 0:
                    numbered.append(f"[ヘッダー]: {line}")
                else:
                    numbered.append(f"行{i}: {line}")
            return "\n".join(numbered)
        else:
            encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'
            return file_bytes.decode(encoding, errors='replace')
    except Exception as e:
        logger.error(f"Error extracting text from {filename} for diff: {e}", exc_info=True)
        return None

def _get_video_mime_type(filename: str) -> str:
    """ファイル名から動画のMIMEタイプを堅牢に推定する。"""
    import mimetypes
    mime_type, _ = mimetypes.guess_type(filename)
    if mime_type and mime_type.startswith("video/"):
        return mime_type

    file_extension = os.path.splitext(filename)[1].lower().strip('.')
    ext_map = {
        "mp4": "video/mp4", "mov": "video/quicktime", "webm": "video/webm",
        "mkv": "video/x-matroska", "mpg": "video/mpeg", "mpeg": "video/mpeg",
        "avi": "video/x-msvideo", "wmv": "video/x-ms-wmv",
    }
    fallback_mime = ext_map.get(file_extension, "video/mp4")
    logger.warning(f"mimetypes failed for '{filename}'. Falling back to MIME type '{fallback_mime}'.")
    return fallback_mime

def _split_video_ffmpeg(input_path: str, output_dir: str, num_parts: int = 2) -> list[str] | None:
    """FFmpegを使用して動画ファイルを指定された数に分割する。"""
    import shutil
    
    # FFmpegとffprobeの実行可能ファイルを検索
    ffmpeg_executable = None
    ffprobe_executable = None
    
    # まずPATHから検索
    if shutil.which("ffmpeg") and shutil.which("ffprobe"):
        ffmpeg_executable = "ffmpeg"
        ffprobe_executable = "ffprobe"
    elif os.name == 'nt':  # Windows
        # Windowsでの代替パスを試す
        alternative_paths = [
            (r'C:\FFmpeg\bin\ffmpeg.exe', r'C:\FFmpeg\bin\ffprobe.exe'),
            (r'C:\ffmpeg\bin\ffmpeg.exe', r'C:\ffmpeg\bin\ffprobe.exe'),
            ('ffmpeg.exe', 'ffprobe.exe')
        ]
        
        for ffmpeg_path, ffprobe_path in alternative_paths:
            try:
                result = subprocess.run([ffmpeg_path, '-version'], capture_output=True, text=True, timeout=5)
                if result.returncode == 0:
                    result2 = subprocess.run([ffprobe_path, '-version'], capture_output=True, text=True, timeout=5)
                    if result2.returncode == 0:
                        ffmpeg_executable = ffmpeg_path
                        ffprobe_executable = ffprobe_path
                        logger.info(f"FFmpegとffprobeを発見: {ffmpeg_path}, {ffprobe_path}")
                        break
            except (subprocess.TimeoutExpired, subprocess.CalledProcessError, FileNotFoundError):
                continue
    
    if not ffmpeg_executable or not ffprobe_executable:
        logger.error("ffmpeg or ffprobe not found. Please install FFmpeg.")
        return None
    if not os.path.exists(input_path):
        logger.error(f"Input file not found for ffmpeg: {input_path}")
        return None

    try:
        probe_cmd = [ffprobe_executable, "-v", "error", "-show_entries", "format=duration", "-of", "default=noprint_wrappers=1:nokey=1", input_path]
        result = subprocess.run(probe_cmd, capture_output=True, text=True, check=True)
        duration = float(result.stdout.strip())
        
        segment_duration = duration / num_parts
        output_files = []
        base, ext = os.path.splitext(os.path.basename(input_path))

        for i in range(num_parts):
            start_time = i * segment_duration
            output_filename = os.path.join(output_dir, f"{base}_part{i+1}{ext}")
            
            split_cmd = [ffmpeg_executable, "-y", "-ss", str(start_time), "-i", input_path, "-t", str(segment_duration), "-c", "copy", output_filename]
            
            subprocess.run(split_cmd, check=True, capture_output=True, text=True)
            output_files.append(output_filename)

        logger.info(f"Video split completed. Output files: {output_files}")
        return output_files

    except (subprocess.CalledProcessError, FileNotFoundError, ValueError) as e:
        stderr_output = e.stderr if hasattr(e, 'stderr') else "(details not available)"
        logger.error(f"Video splitting failed: {e}\n{stderr_output}", exc_info=True)
        return None
    
def parse_vtt_content(vtt_content: str) -> str:
    """
    VTT形式のトランスクリプトから不要な情報（タイムスタンプなど）を除去し、
    話者と発言内容のテキストを抽出します。
    """
    lines = vtt_content.splitlines()
    text_lines = []
    speaker = "不明な話者"
    speech_buffer = []

    for line in lines:
        # ヘッダー、タイムスタンプ、空行は無視
        if '-->' in line or line.strip().upper() == 'WEBVTT' or not line.strip():
            if speech_buffer:
                text_lines.append(f"{speaker}: {' '.join(speech_buffer)}")
                speech_buffer = []
            continue
        
        # 話者情報（例: <v John Doe>）を抽出
        speaker_match = re.match(r'<v\s+([^>]+)>', line)
        if speaker_match:
            if speech_buffer:
                text_lines.append(f"{speaker}: {' '.join(speech_buffer)}")
                speech_buffer = []
            
            speaker = speaker_match.group(1).strip()
            line_content = re.sub(r'<v\s+[^>]+>', '', line).strip()
            if line_content:
                speech_buffer.append(line_content)
        else:
            speech_buffer.append(line.strip())

    if speech_buffer:
        text_lines.append(f"{speaker}: {' '.join(speech_buffer)}")

    return "\n".join(text_lines)

def parse_docx_transcript(file_bytes: bytes) -> str:
    """
    .docx形式のTeamsトランスクリプトファイルを解析し、話者と発言内容を抽出する。
    """
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        
        # Word文書内の表からトランスクリプトを抽出することを想定
        # (Teamsのdocxトランスクリプトは通常、表形式で保存される)
        transcript_lines = []
        for table in doc.tables:
            # 2列のテーブル（1列目: 話者, 2列目: 発言）を想定
            if len(table.columns) >= 2:
                for row in table.rows:
                    speaker = row.cells.text.strip()
                    speech = row.cells.text.strip()
                    if speaker and speech:
                        transcript_lines.append(f"{speaker}: {speech}")
        
        # 表が見つからない場合は、通常の段落からテキストを抽出
        if not transcript_lines:
            for para in doc.paragraphs:
                if para.text.strip():
                    transcript_lines.append(para.text.strip())

        return "\n".join(transcript_lines)
    except Exception as e:
        logger.error(f"DOCXトランスクリプトの解析に失敗: {e}")
        return "(DOCXファイルの解析に失敗しました)"

def _extract_text_from_file_with_ai(uploaded_file, user_info: dict, original_filename: str) -> str | None:
    """
    生成AIを使用して、アップロードされたファイルからテキスト情報を抽出する。
    画像やPDFにも対応。
    """
    logger.info(f"AIによるテキスト抽出を開始: {original_filename}")
    
    # 既存の関数を使って、まずPartオブジェクトを作成
    file_part = extract_file_content_to_part(uploaded_file)
    if not file_part:
        logger.error(f"ファイル '{original_filename}' からPartオブジェクトを作成できませんでした。")
        return None
        
    prompt = """
    あなたは、提供されたファイルから検索インデックス登録用の高品質なテキストを生成する専門家です。
    以下のファイルの内容を分析し、その内容を網羅的かつ詳細にテキスト化してください。

    - **テキストベースのファイル (PDF, DOCX, TXTなど) の場合:**
      - テキストを抽出し、重要なキーワードや概念を保持しつつ、冗長な表現を削除して要約してください。
      - 元のドキュメントの構造（見出し、リストなど）がわかるように、適切な改行やインデントを使ってください。
    - **画像ファイル (PNG, JPGなど) の場合:**
      - 画像に写っているオブジェクト、人物、風景、テキスト、グラフ、図表などを詳細に説明してください。
      - 画像の目的や伝えているメッセージを推測し、それもテキストに含めてください。

    出力は、生成されたテキストのみとしてください。前置きや説明は不要です。
    """

    try:
        parent_blob_name = ensure_parent_log_blob_name(
            user_info=user_info,
            feature_category="RAGインデックス作成 (親ログ)",
            model_name=config.DEFAULT_MODEL_NAME,
            additional_log_params={"source": "file_processor", "original_filename": original_filename},
        )

        # generate_with_gemini を呼び出してテキストを生成
        extracted_text, metadata = generate_with_gemini(
            model_name=config.DEFAULT_MODEL_NAME,
            contents=[types.Content(role="user", parts=[file_part, types.Part.from_text(text=prompt)])],
            system_instruction="あなたは、あらゆる種類のファイルからテキスト情報を抽出し、検索しやすいように要約・説明するAIアシスタントです。応答は抽出されたテキストのみとしてください。",
            max_output_tokens=65535, # 長めのテキストも扱えるように
            temperature=0.1, # 決定性を高める
            user_info=user_info,
            log_prompt_text=f"AI Text Extraction for RAG from {original_filename}",
            log_filenames=[original_filename],
            feature_category="RAGインデックス作成 (AIテキスト抽出)",
            skip_cosmos_log=True, # この内部処理は個別にログ記録しない
            parent_log_id=parent_blob_name,
        )

        if extracted_text and not extracted_text.startswith("[エラー:") and not extracted_text.startswith("[APIエラー:"):
            logger.info(f"AIによるテキスト抽出に成功: {original_filename} (文字数: {len(extracted_text)})")
            return extracted_text.strip()
        else:
            logger.error(f"AIによるテキスト抽出に失敗 (APIエラーまたは空の応答): {original_filename}. Response: {extracted_text}")
            return None

    except Exception as e:
        logger.error(f"AIテキスト抽出処理中に予期せぬエラーが発生 ({original_filename}): {e}", exc_info=True)
        return None
    
def _extract_structured_content_with_ai(uploaded_file, user_info: dict, original_filename: str) -> dict | None:
    """
    【改訂版 v5・指示強化版】生成AI(Gemini)を使用して、ファイルから構造化されたコンテンツをJSON形式で抽出する。
    長い入力に対しても指示を遵守させるため、プロンプトの構造を「指示の事前定義 -> タスク実行」の順に変更し、ルールを強化する。
    """
    logger.info(f"AIによる構造化コンテンツ抽出(v5・指示強化版)を開始: {original_filename}")
    
    file_part = extract_file_content_to_part(uploaded_file)
    if not file_part:
        logger.error(f"ファイル '{original_filename}' からPartオブジェクトを作成できませんでした。")
        return None
        
    prompt = """
あなたは、検索インデックス作成のために、あらゆる種類のドキュメントを構造化されたJSON形式に変換する、高度なドキュメント解析AIです。
あなたの唯一のタスクは、提供されたファイルを分析し、その内容を後述する【厳格なJSON出力形式】に完全に準拠した形で出力することです。

### 【最重要】厳格なJSON出力形式 (この形式を絶対に遵守してください)
- **出力は単一のJSONオブジェクトのみ:** 応答は、必ず `{"content_blocks": [...]}` という構造を持つ単一のJSONオブジェクトの文字列でなければなりません。説明文や、```json のようなマークダウンの囲み文字は**絶対に含めないでください**。
- **ブロックオブジェクトの必須キー:** 各ブロックオブジェクトは、`type` キーと、それに紐づく必須キーをすべて含んでいる必要があります。
- **JSONスキーマ定義:**
  - `content_blocks`: (Array) ドキュメントの構成要素ブロックのリスト。
    - `type`: (String) ブロックの種類。`"heading"`, `"paragraph"`, `"table"`, `"image_analysis"` のいずれか。
    - **(type: "heading")**
      - `level`: (Integer) 見出しレベル (例: 1, 2, 3)。
      - `text`: (String) 見出しのテキスト。
    - **(type: "paragraph")**
      - `raw_text`: (String) 元の段落のテキスト。
      - `summary`: (String) その段落の簡潔な要約。
      - `keywords`: (Array of Strings) その段落の主要なキーワード。
    - **(type: "table")**
      - `markdown_representation`: (String) 表をMarkdown形式で表現した文字列。
      - `summary`: (String) 表が示す内容の簡潔な要約。
      - `keywords`: (Array of Strings) 表に関連するキーワード。
    - **(type: "image_analysis")**
      - `description`: (String) 画像の内容に関する詳細な説明。
      - `ocr_text`: (String) 画像から読み取ったテキスト（OCR）。なければ空文字。
      - `summary`: (String) 画像の内容とOCRテキストを統合した要約。
      - `keywords`: (Array of Strings) 画像に関連するキーワード。

### 思考プロセス (このプロセスに従ってタスクを実行してください)
1.  **ブロック単位での分解:** 提供されたファイルを、意味のある構成単位（見出し、段落、表、画像）に一つずつ分解します。
2.  **各ブロックの分析:** 分解した**ブロックごと**に、上記のJSONスキーマで定義された情報を生成します。
3.  **JSONの構築:** 生成した情報を、上記のスキーマに厳密に従って組み立てます。文書全体の包括的な要約は生成しないでください。

---
**それでは、以下のファイルの内容を分析し、上記の指示に厳密に従ってJSONオブジェクトを生成してください。**
"""

    try:
        model_name_for_extraction = os.environ.get("DOC_GEN_MODEL_NAME", config.DEFAULT_MODEL_NAME)

        parent_blob_name = ensure_parent_log_blob_name(
            user_info=user_info,
            feature_category="RAGインデックス作成 (親ログ)",
            model_name=model_name_for_extraction,
            additional_log_params={"source": "file_processor", "original_filename": original_filename},
        )
        
        extracted_json_str, metadata = generate_with_gemini(
            model_name=model_name_for_extraction,
            contents=[types.Content(role="user", parts=[types.Part.from_text(text=prompt), file_part])], # 指示プロンプトとファイルを分離
            max_output_tokens=65535, # モデルの最大値近くまで許容
            temperature=0.0, # 決定性を最大にする
            user_info=user_info,
            log_prompt_text=f"AI Structured Extraction v5 for RAG from {original_filename}",
            log_filenames=[original_filename],
            feature_category="RAGインデックス作成 (AI構造化抽出 v5)",
            skip_cosmos_log=True,
            parent_log_id=parent_blob_name,
            generation_config_override={"response_mime_type": "application/json"}
        )

        if extracted_json_str and not extracted_json_str.startswith("[エラー:"):
            # 制御文字を除去するクリーニング処理は維持
            cleaned_json_str = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', extracted_json_str)

            json_match = re.search(r'\{.*\}', cleaned_json_str, re.DOTALL)
            if not json_match:
                # エラーメッセージをより具体的に
                raise ValueError(f"AIがJSONオブジェクトを生成しませんでした。AI応答(先頭500文字): {cleaned_json_str[:500]}")

            json_to_parse = json_match.group(0)
            parsed_json = json.loads(json_to_parse, strict=False)

            if "content_blocks" in parsed_json and isinstance(parsed_json["content_blocks"], list):
                logger.info(f"AIによる構造化コンテンツ抽出(v5)に成功: {original_filename} ({len(parsed_json['content_blocks'])}ブロック)")
                return parsed_json
            else:
                raise ValueError("JSONの形式が正しくありません。'content_blocks'キーが見つからないか、配列ではありません。")
        else:
            logger.error(f"AIによる構造化コンテンツ抽出(v5)に失敗: {original_filename}. Response: {extracted_json_str}")
            return None

    except (json.JSONDecodeError, ValueError) as e:
        logger.error(f"AI応答のJSONパースに失敗 (v5) ({original_filename}): {e}. Raw response: {extracted_json_str[:500]}...")
        return None
    except Exception as e:
        logger.error(f"AI構造化抽出処理(v5)中に予期せぬエラー ({original_filename}): {e}", exc_info=True)
        return None

def extract_text_for_simple_rag(file_bytes: bytes, filename: str) -> str:
    """
    【新規追加】
    RAG用に、AIを使わずにファイルからプレーンテキストを抽出する関数。
    extract_file_content_to_part のロジックをベースに、文字列のみを返すように簡略化。
    """
    filename_lower = filename.lower()
    
    try:
        # 1. PDF
        if filename_lower.endswith('.pdf'):
            try:
                with fitz.open(stream=file_bytes, filetype="pdf") as doc:
                    text = ""
                    for page in doc:
                        text += page.get_text() + "\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting text from PDF '{filename}': {e}")
                return ""

        # 2. Word (docx)
        elif filename_lower.endswith('.docx'):
            try:
                doc = docx.Document(io.BytesIO(file_bytes))
                return "\n".join([p.text for p in doc.paragraphs])
            except Exception as e:
                logger.error(f"Error extracting text from DOCX '{filename}': {e}")
                return ""

        # 3. PowerPoint (pptx)
        elif filename_lower.endswith('.pptx'):
            try:
                presentation = pptx.Presentation(io.BytesIO(file_bytes))
                text_runs = []
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return "\n".join(text_runs)
            except Exception as e:
                logger.error(f"Error extracting text from PPTX '{filename}': {e}")
                return ""

        # 4. Excel (xlsx, xls)
        elif filename_lower.endswith(('.xlsx', '.xls')):
            try:
                xls_dfs = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
                text_parts = []
                for sheet_name, df in xls_dfs.items():
                    text_parts.append(f"--- Sheet: {sheet_name} ---")
                    text_parts.append(df.to_csv(index=False))
                return "\n".join(text_parts)
            except Exception as e:
                logger.error(f"Error extracting text from Excel '{filename}': {e}")
                return ""

        # 5. CSV
        elif filename_lower.endswith('.csv'):
            try:
                # 文字コード判定
                encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'
                df = pd.read_csv(io.BytesIO(file_bytes), encoding=encoding)
                return df.to_csv(index=False)
            except Exception as e:
                logger.error(f"Error extracting text from CSV '{filename}': {e}")
                return ""
        
        # 6. 一般的なテキストファイル (txt, md, py, js, etc.)
        # python-magic等での判定も良いが、ここではシンプルに拡張子とデコード試行で対応
        else:
            try:
                encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'
                return file_bytes.decode(encoding, errors='replace')
            except Exception as e:
                logger.warning(f"Failed to decode text file '{filename}': {e}")
                return ""

    except Exception as e:
        logger.error(f"Unexpected error in simple text extraction for '{filename}': {e}", exc_info=True)
        return ""

