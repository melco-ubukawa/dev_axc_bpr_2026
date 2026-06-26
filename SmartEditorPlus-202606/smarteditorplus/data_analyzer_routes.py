import os
import uuid
import json
import io
import re
import base64
import logging
import tempfile
import shutil
import config
import requests
import markdown
import copy
import pandas as pd
import numpy as np
import joblib
import matplotlib.pyplot as plt
import seaborn as sns  

from bs4 import BeautifulSoup, NavigableString
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx import Presentation
from pptx.util import Inches

from sklearn.preprocessing import StandardScaler
from sklearn.pipeline import Pipeline
from sklearn.compose import ColumnTransformer
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import train_test_split
from sklearn.metrics import mean_squared_error, f1_score, mean_absolute_percentage_error
from sklearn.ensemble import RandomForestRegressor, RandomForestClassifier
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.impute import SimpleImputer
from statsmodels.tsa.statespace.sarimax import SARIMAX
import shap

from flask import Blueprint, render_template, request, jsonify, send_file
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from .utils.auth import get_entra_user_info
from .utils.cosmos_db import (
    get_tasks_container, create_task_in_cosmos, get_task_from_cosmos, update_task_in_cosmos,
    get_sessions_container, get_projects_by_user, InMemoryContainer, delete_design_project, get_models_by_project_id
)
from .utils.storage import upload_bytes_to_blob, download_blob_to_bytes, delete_blob, get_blob_sas_url, get_gcs_signed_url, _upload_to_gcs, _delete_from_gcs, download_gcs_blob_to_bytes
from .utils.gemini import generate_with_gemini, ensure_parent_log_blob_name
from .utils.file_processor import _read_dataframe_with_multiple_encodings
from .utils.sandbox import execute_python_safely
from google.genai import types

# --- ロガー設定 ---
logger = logging.getLogger(__name__)

# --- Blueprintの作成 ---
data_analyzer_bp = Blueprint(
    'data_analyzer',
    __name__,
    template_folder='../templates',
    static_folder='../static'
)

# --- 非同期タスク実行用のExecutor ---
executor = ThreadPoolExecutor(max_workers=max(4, os.cpu_count() or 1))

# === データベースヘルパー (既存のutils/cosmos_db.pyのラッパー) ===
def get_analyzer_container():
    return get_sessions_container()

def create_analysis_project(project_data):
    container = get_analyzer_container()
    if not container:
        logger.error("Failed to get analyzer container.")
        return None
    try:
        return container.create_item(body=project_data)
    except Exception as e:
        logger.error(f"Failed to create analysis project in Cosmos DB: {e}", exc_info=True)
        return None

def get_analysis_project(project_id, user_id):
    container = get_analyzer_container()
    if not container:
        logger.error("Failed to get analyzer container.")
        return None
    try:
        return container.read_item(item=project_id, partition_key=user_id)
    except Exception:
        return None

def update_analysis_project(project_id, user_id, updates):
    container = get_analyzer_container()
    if not container:
        logger.error("Failed to get analyzer container.")
        return None
    
    # Patch操作の構築
    patch_operations = []
    for key, value in updates.items():
        patch_operations.append({"op": "add", "path": f"/{key}", "value": value})
    
    if not patch_operations:
        return None

    try:
        # patch_item を使用して部分更新を行う (RU消費を劇的に削減)
        return container.patch_item(
            item=project_id,
            partition_key=user_id,
            patch_operations=patch_operations
        )
    except Exception as e:
        logger.error(f"Failed to patch analysis project {project_id}: {e}", exc_info=True)
        # パッチが失敗した場合（ドキュメントが存在しない等）のフォールバックが必要ならここに記述
        return None

def _load_df_from_bytes_with_temp_file(data_bytes, original_filename):
    temp_dir = tempfile.mkdtemp()
    temp_file_path = os.path.join(temp_dir, original_filename)
    try:
        with open(temp_file_path, 'wb') as f:
            f.write(data_bytes)
        df = _read_dataframe_with_multiple_encodings(temp_file_path)
        return df
    finally:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)

# === ページレンダリング用ルート ===
@data_analyzer_bp.route('/upload')
def upload_page():
    return render_template('upload.html')

@data_analyzer_bp.route('/profiling/<project_id>')
def profiling_page(project_id):
    return render_template('profiling.html', project_id=project_id)

@data_analyzer_bp.route('/analysis/<project_id>')
def analysis_page(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    project_data = get_analysis_project(project_id, user_id)
    project_status = project_data.get("status") if project_data else "unknown"
    return render_template('analysis.html', project_id=project_id, project_status=project_status)

@data_analyzer_bp.route('/history')
def history_page():
    """分析履歴一覧ページをレンダリングします。"""
    # このルートに対応する新しいHTMLファイル (history.html) が別途必要になります。
    return render_template('history.html')

@data_analyzer_bp.route('/models')
def models_page():
    """予測モデル管理ページをレンダリングします。"""
    return render_template('models.html')

@data_analyzer_bp.route('/models/<model_id>')
def model_detail_page(model_id):
    """予測モデル詳細ページをレンダリングします。"""
    # model_idをテンプレートに渡して、JSが利用できるようにする
    return render_template('model_detail.html', model_id=model_id)

# === APIエンドポイント ===

@data_analyzer_bp.route('/api/projects', methods=['GET'])
def get_user_projects_api():
    """現在のユーザーに紐づく過去の分析プロジェクト一覧を返します。"""
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        # 新しく追加した、ユーザーIDでプロジェクトを取得する関数を呼び出します
        projects = get_projects_by_user(user_id)
        return jsonify(projects)
    except Exception as e:
        logger.error(f"ユーザー {user_id} のプロジェクト一覧取得中にエラー: {e}", exc_info=True)
        return jsonify({"error": "プロジェクト一覧の取得に失敗しました。"}), 500

@data_analyzer_bp.route('/api/task_status/<task_id>', methods=['GET'])
def get_task_status_api(task_id):
    try:
        task = get_task_from_cosmos(task_id)
        if task:
            return jsonify(task)
        else:
            return jsonify({"error": "Task not found", "status": "NOT_FOUND"}), 404
    except Exception as e:
        logger.error(f"Error getting task status for {task_id}: {e}", exc_info=True)
        return jsonify({"error": "Failed to retrieve task status"}), 500

@data_analyzer_bp.route('/api/upload_data', methods=['POST'])
def upload_data_api():
    user_info = get_entra_user_info(request)
    if 'data_file' not in request.files: return jsonify({"error": "分析対象のデータファイルがありません。"}), 400
    data_file = request.files['data_file']
    reference_files = request.files.getlist('reference_files')
    session_id = f"icp_session_{uuid.uuid4().hex}"
    try:
        data_bytes = data_file.read()
        data_blob_name = f"temp_uploads/{session_id}/{data_file.filename}"
        upload_bytes_to_blob(data_bytes, data_blob_name)
        reference_file_blobs = []
        for ref_file in reference_files:
            ref_bytes = ref_file.read()
            ref_blob_name = f"temp_uploads/{session_id}/references/{ref_file.filename}"
            upload_bytes_to_blob(ref_bytes, ref_blob_name)
            reference_file_blobs.append({"filename": ref_file.filename, "blob_name": ref_blob_name})
        return jsonify({"success": True, "session_id": session_id, "data_file": {"filename": data_file.filename, "blob_name": data_blob_name}, "reference_files": reference_file_blobs}), 200
    except Exception as e:
        logger.error(f"データの一時アップロード中にエラー: {e}", exc_info=True)
        return jsonify({"error": "ファイルの一時アップロードに失敗しました。"}), 500

def _convert_numeric_to_datetime_based_on_role(df: pd.DataFrame, column_roles: dict) -> pd.DataFrame:
    """
    ユーザーが「日付 / 時刻」と指定した数値(int/float)列を、datetime型に変換する。
    YYYYMMDD, YYYYMM, YYYY 形式を自動的に判定しようと試みる。
    """
    for col_name, role in column_roles.items():
        if role == '日付 / 時刻' and col_name in df.columns and pd.api.types.is_numeric_dtype(df[col_name].dtype):
            try:
                series = df[col_name].dropna()
                if series.empty:
                    continue

                # 値の大きさや桁数からフォーマットを推測する
                # 例: 20240101.0 のようなfloat型も考慮してintに変換
                max_val = int(series.max())
                
                format_str = None
                if max_val > 19000000 and max_val < 22000000: # YYYYMMDD (例: 20240515)
                    format_str = '%Y%m%d'
                elif max_val > 190000 and max_val < 220000: # YYYYMM (例: 202405)
                    format_str = '%Y%m'
                elif max_val > 1900 and max_val < 2200: # YYYY (例: 2024)
                    format_str = '%Y'

                if format_str:
                    # to_datetimeは文字列を期待するため、一度文字列に変換してから適用
                    converted_series = pd.to_datetime(df[col_name].astype(str), format=format_str, errors='coerce')
                    
                    # 50%以上が正常に変換できたら採用（元のロジックに合わせる）
                    if converted_series.notna().sum() / df[col_name].notna().sum() > 0.5:
                        df[col_name] = converted_series
                        logger.info(f"役割設定に基づき、数値列 '{col_name}' を日付/時刻型に変換しました (フォーマット: {format_str})。")
            except Exception as e:
                logger.warning(f"役割設定に基づく数値列 '{col_name}' の日付変換に失敗しました: {e}")
    return df

def _try_convert_datetime_columns(df: pd.DataFrame) -> pd.DataFrame:
    """
    DataFrame内のobject型（文字列）の列を、日付/時刻型へ自動変換することを試みます。
    列内の値の50%以上が日付として正常に変換できた場合のみ、その列のデータ型をdatetime型に置き換えます。
    'YYYYMMDD', 'YYYY/MM/DD', 'YYYY-MM' など、多様な形式に対応します。
    """
    for col_name in df.select_dtypes(include=['object']).columns:
        try:
            # `errors='coerce'`は、日付に変換できない値をエラーにせず、NaT(Not a Time)に変換します
            converted_series = pd.to_datetime(df[col_name], errors='coerce')
            
            # 元のデータでNULLでなかった値の数をカウント
            original_not_null_count = df[col_name].notna().sum()
            if original_not_null_count == 0:
                continue

            # 日付変換後にNULLでなくなった値（正常に変換された値）の数をカウント
            converted_not_null_count = converted_series.notna().sum()
            
            # 変換の成功率を計算
            success_rate = converted_not_null_count / original_not_null_count
            
            # 成功率が50%を超え、かつデータ型が実際にdatetimeに変わった場合、列を置き換える
            if success_rate > 0.5 and pd.api.types.is_datetime64_any_dtype(converted_series.dtype):
                df[col_name] = converted_series
                logger.info(f"列 '{col_name}' を日付/時刻型に自動変換しました (成功率: {success_rate:.1%})。")
        except (ValueError, TypeError):
            # 非常に大きな数値や混合型で `to_datetime` が失敗するケースを無視
            continue
        except Exception as e:
            logger.warning(f"列 '{col_name}' の日付変換試行中に予期せぬエラーが発生しました: {e}")
    return df

@data_analyzer_bp.route('/api/setup_suggestions', methods=['POST'])
def get_setup_suggestions_api():
    user_info = get_entra_user_info(request)
    data = request.get_json()
    session_id, data_file_info, analysis_purpose = data.get('session_id'), data.get('data_file'), data.get('purpose')
    if not all([session_id, data_file_info, analysis_purpose]): return jsonify({"error": "必要な情報が不足しています。"}), 400
        
    try:
        data_bytes = download_blob_to_bytes(data_file_info['blob_name'])
        df = _load_df_from_bytes_with_temp_file(data_bytes, data_file_info['filename'])
        if df is None: raise Exception("データファイルの読み込みに失敗しました。")

        if hasattr(config, 'MAX_UPLOAD_ROWS') and len(df) > config.MAX_UPLOAD_ROWS:
            error_message = f"データが上限行数を超えています。上限: {config.MAX_UPLOAD_ROWS: ,}行, 現在の行数: {len(df): ,}行"
            logger.warning(f"Upload rejected for user {user_info.get('userId')}: {error_message}")
            # 400 Bad Request を返すことで、フロントエンドでエラー内容をハンドリングしやすくする
            return jsonify({"error": error_message}), 400

        df = _try_convert_datetime_columns(df)

        columns_with_details = []
        for col_name in df.columns:
            series = df[col_name]
            # 日付型の場合、表示形式を整える
            if pd.api.types.is_datetime64_any_dtype(series.dtype):
                sample_values = series.dropna().head(3).dt.strftime('%Y-%m-%d').tolist()
            else:
                sample_values = series.dropna().head(3).astype(str).tolist()
            
            details = {"column_name": col_name, "type": str(series.dtype), "sample_values": sample_values}
            columns_with_details.append(details)
        
        detailed_info_for_ai = json.dumps(columns_with_details, indent=2, ensure_ascii=False)
        
        suggestion_prompt = f"""あなたはデータ分析プロジェクトの立ち上げを支援するAIアシスタントです。ユーザーの分析目的とデータ列の詳細情報を元に、最適なセットアップを提案してください。日本語で回答してください。

### ユーザーの分析目的
{analysis_purpose}

### データ列の詳細情報```json
{detailed_info_for_ai}
```

### あなたのタスク
上記の情報を元に、以下の3つの項目を**単一のJSONオブジェクト**として提案してください。
1.  **analysis_recipe**: 分析目的に最適な分析レシピ。
2.  **column_roles**: 各列の役割。役割は `ID / キー`, `予測ターゲット`, `層別（セグメント）`, `日付 / 時刻`, `自由記述（テキスト）`, `数値（連続）`, `選択式回答（カテゴリ）`, `分析対象外` から選択。
3.  **setup_suggestions**: 重要指標、参考資料の活用方針、警告ルール設定に関する提案。

### 出力形式 (JSONのみ、キー名は厳守)
```json
{{
  "analysis_recipe": {{ "name": "...", "description": "..." }},
  "column_roles": [
    {{"column_name": "列名1", "suggested_role": "推定役割1"}},
    {{"column_name": "列名2", "suggested_role": "推定役割2"}}
  ],
  "setup_suggestions": {{ "key_metrics": "...", "reference_usage": "...", "alert_rules": "..." }}
}}
```"""
        
        ai_response_str, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME, contents=[types.Part.from_text(text=suggestion_prompt)],
            max_output_tokens=65535, temperature=0.1, user_info=user_info,
            generation_config_override={"response_mime_type": "application/json"},
            feature_category="Insight Canvas",
        )
        if not ai_response_str or ai_response_str.strip().startswith("[エラー") or ai_response_str.strip().startswith("[APIエラー"):
             # エラーが返ってきた場合は、その内容を例外としてスローする
             raise Exception(f"AIによる提案の生成に失敗しました: {ai_response_str}")

        ai_suggestions = json.loads(ai_response_str)
        
        final_suggestions = {
            "analysis_recipe": ai_suggestions.get("analysis_recipe"),
            "setup_suggestions": ai_suggestions.get("setup_suggestions"),
            "column_details": [],
            "data_summary": {
                "rowCount": len(df),
                "columnCount": len(df.columns)
            }
        }
        
        ai_roles_map = {item.get('column_name'): item.get('suggested_role') for item in ai_suggestions.get('column_roles', [])}

        for col_detail in columns_with_details:
            if pd.api.types.is_datetime64_any_dtype(col_detail['type']):
                ai_roles_map[col_detail['column_name']] = "日付 / 時刻"

        for col_detail in columns_with_details:
            col_name = col_detail['column_name']
            final_suggestions["column_details"].append({
                "column_name": col_name,
                "type": col_detail["type"],
                "sample_values": col_detail["sample_values"],
                "suggested_role": ai_roles_map.get(col_name, "分析対象外")
            })

        return jsonify({"success": True, "suggestions": final_suggestions})

    except Exception as e:
        logger.error(f"セットアップ提案の生成中にエラー: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

@data_analyzer_bp.route('/api/project', methods=['PUT'])
def finalize_project_api():
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    session_id, data_file_info, reference_files_info, analysis_purpose, column_roles = data.get('session_id'), data.get('data_file'), data.get('reference_files', []), data.get('purpose', ''), data.get('column_roles', {})
    if not all([session_id, data_file_info, analysis_purpose, column_roles]): return jsonify({"error": "プロジェクトの確定に必要な情報が不足しています。"}), 400
    project_id = f"icp_{uuid.uuid4().hex}"
    try:
        project_data = {"id": project_id, "userId": user_id, "type": "AnalysisProject", "projectName": data_file_info['filename'], "analysisPurpose": analysis_purpose, "columnRoles": column_roles, "dataFile": data_file_info, "referenceFiles": reference_files_info, "status": "profiling_started", "createdAt": datetime.now(timezone.utc).isoformat()}
        create_analysis_project(project_data)
        task_id = f"task_{project_id}_profiling"
        create_task_in_cosmos({"id": task_id, "task_id": task_id, "status": "PENDING", "task_type": "data_profiling", "project_id": project_id, "user_id": user_id})
        executor.submit(_run_profiling_task, task_id, project_id, user_id)
        return jsonify({"project_id": project_id, "task_id": task_id}), 202
    except Exception as e:
        logger.error(f"プロジェクト確定処理中にエラー: {e}", exc_info=True)
        return jsonify({"error": "プロジェクトの作成に失敗しました。"}), 500

def _sanitize_stats_dict(stats: dict) -> dict:
    """pandas.describe()が出力する辞書をJSONシリアライズ可能にする"""
    sanitized = {}
    for key, value in stats.items():
        # pd.isna()はNaNを、isfinite()は無限大(inf)をチェックする
        if pd.isna(value) or (isinstance(value, float) and not np.isfinite(value)):
            sanitized[key] = None  # JSONのnullに変換
        elif isinstance(value, pd.Timestamp):
            sanitized[key] = value.isoformat()
        else:
            sanitized[key] = value
    return sanitized

def _run_profiling_task(task_id, project_id, user_id):
    update_task_in_cosmos(task_id, {"status": "PROCESSING", "status_message": "プロジェクト情報を読み込み中...", "progress_percent": 10})
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data: raise Exception("プロジェクトデータが見つかりません。")
        update_task_in_cosmos(task_id, {"status_message": "データファイルを読み込み中...", "progress_percent": 20})
        data_bytes = download_blob_to_bytes(project_data['dataFile']['blob_name'])
        df = _load_df_from_bytes_with_temp_file(data_bytes, project_data['dataFile']['filename'])
        if df is None: raise Exception("データファイルの読み込みに失敗しました。")

        column_roles_info = project_data.get('columnRoles', {})
        df = _convert_numeric_to_datetime_based_on_role(df, column_roles_info)

        df = _try_convert_datetime_columns(df)
        update_task_in_cosmos(task_id, {"status_message": "データの基本統計を計算中...", "progress_percent": 40})
        profiling_result = {'summary': {'tableName': project_data['projectName'], 'rowCount': len(df), 'columnCount': len(df.columns), 'lastModified': project_data.get('createdAt')}}
        
        column_roles = project_data.get('columnRoles', {})
        column_summaries = []
        for col in df.columns:
            series = df[col]
            summary = {
                'name': col, 
                'type': str(series.dtype), 
                'missing_percentage': series.isnull().sum() / len(df) * 100 if len(df) > 0 else 0,
                'role': column_roles.get(col, '未設定'),
                'memo': ''
            }
            if pd.api.types.is_numeric_dtype(series.dtype):
                summary['stats'] = series.describe().to_dict()
                sample_size = min(1000, len(series.dropna()))
                summary['distribution_data'] = series.dropna().sample(sample_size, random_state=1).tolist() if sample_size > 0 else []
            else:
                summary['stats'] = series.describe().to_dict()
                
            if pd.api.types.is_numeric_dtype(series.dtype):
                stats_dict = series.describe().to_dict()
                summary['stats'] = _sanitize_stats_dict(stats_dict) # ヘルパー関数を適用
                sample_size = min(1000, len(series.dropna()))
                summary['distribution_data'] = series.dropna().sample(sample_size, random_state=1).tolist() if sample_size > 0 else []
            else:
                stats_dict = series.describe().to_dict()
                summary['stats'] = _sanitize_stats_dict(stats_dict) # こちらもヘルパー関数を適用

                if pd.api.types.is_datetime64_any_dtype(series.dtype):
                    vc = series.value_counts(normalize=True).head(5)
                    vc.index = vc.index.strftime('%Y-%m-%d %H:%M:%S')
                    summary['distribution_data'] = vc.to_dict()
                else:
                    summary['distribution_data'] = series.value_counts(normalize=True).head(5).to_dict()
            column_summaries.append(summary)

        profiling_result['column_summaries'] = column_summaries

        serializable_profiling_result = json.loads(json.dumps(profiling_result, default=str))

        update_task_in_cosmos(task_id, {"status_message": "AIがデータ品質を評価中...", "progress_percent": 70})

        profiling_for_ai = copy.deepcopy(serializable_profiling_result)
        if 'summary' in profiling_for_ai:
            # lastModified や tableName はデータの中身の品質とは無関係なので削除
            profiling_for_ai['summary'].pop('lastModified', None)
            profiling_for_ai['summary'].pop('tableName', None)
        
        # JSON文字列化には、フィルタリング済みの辞書を使用
        profiling_text_for_ai = json.dumps(profiling_for_ai, indent=2, ensure_ascii=False)

        ai_quality_prompt = f"以下のデータプロファイル(JSON)を分析し、「品質インジケーター」と「品質アラート」(問題点と推奨アクションのセット)を日本語で生成してください。\n\n### データプロファイル (抜粋)\n```json\n{profiling_text_for_ai[:8000]}\n```\n\n### 出力形式 (JSONのみ)\n```json\n{{\"quality_indicator\": {{\"status\": \"（安定/要確認/危険）\", \"metrics\": [{{\"label\": \"...\", \"value\": \"...\"}}]}},\"quality_alerts\": [{{\"issue\": \"...\", \"recommendation\": \"...\"}}]}}\n```"

        ai_quality_response_str, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME, 
            contents=[types.Part.from_text(text=ai_quality_prompt)], 
            max_output_tokens=65535, 
            temperature=0.1, 
            generation_config_override={"response_mime_type": "application/json"},
            feature_category="Insight Canvas",
        )
        serializable_profiling_result.update(json.loads(ai_quality_response_str))

        # 1. 結果をJSONバイトデータに変換
        result_bytes = json.dumps(serializable_profiling_result, ensure_ascii=False).encode('utf-8')
        
        # 2. Blob Storageにアップロード
        blob_name = f"results/{project_id}/profiling_result.json"
        upload_bytes_to_blob(result_bytes, blob_name)

        # 3. CosmosDBには結果本体ではなく、Blob名とステータスを保存
        update_analysis_project(project_id, user_id, {
            "profilingResultBlobName": blob_name, # 結果本体の代わりにBlob名
            "status": "profiling_completed"
        })

        update_task_in_cosmos(task_id, {"status": "SUCCESS", "status_message": "プロファイリングが完了しました。", "progress_percent": 100})
    except Exception as e:
        logger.error(f"プロファイリングタスク {task_id} でエラー: {e}", exc_info=True)
        update_task_in_cosmos(task_id, {"status": "FAILURE", "error": str(e), "status_message": "プロファイリング中にエラーが発生しました。"})

@data_analyzer_bp.route('/api/profile/<project_id>', methods=['GET'])
def get_profiling_result_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    project_data = get_analysis_project(project_id, user_id)
    if not project_data: return jsonify({"error": "プロジェクトが見つかりません。"}), 404
    
    profiling_result = None
    # 新しいBlob名フィールドを優先してチェック
    if blob_name := project_data.get("profilingResultBlobName"):
        try:
            result_bytes = download_blob_to_bytes(blob_name)
            if result_bytes:
                profiling_result = json.loads(result_bytes.decode('utf-8'))
            else:
                return jsonify({"error": f"プロファイリング結果ファイル({blob_name})のダウンロードに失敗しました。"}), 404
        except Exception as e:
            logger.error(f"Blob '{blob_name}' からのプロファイリング結果の読み込みに失敗: {e}", exc_info=True)
            return jsonify({"error": "プロファイリング結果の読み込み中にエラーが発生しました。"}), 500
            
    # 後方互換性のため、古いフィールドもチェック
    elif "profilingResult" in project_data:
        profiling_result = project_data.get("profilingResult")    

    if not profiling_result: return jsonify({"error": "プロファイリング結果がまだありません。"}), 404
    return jsonify(profiling_result)
    
@data_analyzer_bp.route('/api/profile/<project_id>/memo', methods=['PATCH'])
def update_column_memo_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    column_name = data.get('column_name')
    memo_text = data.get('memo')

    if not column_name or memo_text is None:
        return jsonify({"error": "カラム名とメモの内容は必須です。"}), 400

    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404

        profiling_result = None
        blob_name = project_data.get("profilingResultBlobName")
        
        # Blobからデータを読み込む
        if blob_name:
            try:
                result_bytes = download_blob_to_bytes(blob_name)
                if result_bytes:
                    profiling_result = json.loads(result_bytes.decode('utf-8'))
                else:
                    raise Exception("Blob data is empty.")
            except Exception as e:
                return jsonify({"error": f"プロファイリング結果ファイルの読み込みに失敗しました: {e}"}), 500
        # 後方互換性
        elif "profilingResult" in project_data:
            profiling_result = project_data.get("profilingResult")

        if not profiling_result or 'column_summaries' not in profiling_result:
            return jsonify({"error": "プロファイリング結果が見つかりません。"}), 404

        target_column = next((col for col in profiling_result['column_summaries'] if col['name'] == column_name), None)

        if not target_column:
            return jsonify({"error": f"カラム '{column_name}' が見つかりません。"}), 404

        target_column['memo'] = memo_text
        
        # Blobに書き戻す
        if blob_name:
            updated_bytes = json.dumps(profiling_result, ensure_ascii=False).encode('utf-8')
            upload_bytes_to_blob(updated_bytes, blob_name)
        # 後方互換性: Cosmos DBフィールドを更新
        else:
            update_analysis_project(project_id, user_id, {"profilingResult": profiling_result})
        
        return jsonify({"success": True, "message": "メモを更新しました。"})

    except Exception as e:
        logger.error(f"プロジェクト {project_id} のメモ更新中にエラー: {e}", exc_info=True)
        return jsonify({"error": "メモの更新に失敗しました。"}), 500

def _generate_feature_engineering_code(rules: dict) -> str:
    """保存された特徴量エンジニアリングルール(JSON)を、実行可能なPython(pandas)コードに変換する。"""
    if not rules:
        return ""

    code_lines = ["# --- 自動特徴量エンジニアリング ---", "import numpy as np", "import pandas as pd"]
    
    # 1. 時間ベース特徴量
    time_rules = rules.get('time_based', {})
    for col, features in time_rules.items():
        code_lines.append(f"\n# [{col}] から時間ベースの特徴量を生成")
        code_lines.append(f"if '{col}' in df.columns:")
        code_lines.append(f"    # Ensure column is datetime format, coercing errors")
        code_lines.append(f"    df['{col}'] = pd.to_datetime(df['{col}'], errors='coerce')")
        
        if features.get('year'):
            code_lines.append(f"    if pd.api.types.is_datetime64_any_dtype(df['{col}']): df['{col}_year'] = df['{col}'].dt.year")
        if features.get('month'):
            code_lines.append(f"    if pd.api.types.is_datetime64_any_dtype(df['{col}']): df['{col}_month'] = df['{col}'].dt.month")
        if features.get('weekday'):
            code_lines.append(f"    if pd.api.types.is_datetime64_any_dtype(df['{col}']): df['{col}_weekday'] = df['{col}'].dt.day_name()")
        if features.get('quarter'):
            code_lines.append(f"    if pd.api.types.is_datetime64_any_dtype(df['{col}']): df['{col}_quarter'] = df['{col}'].dt.quarter")
        if features.get('is_weekend'):
            code_lines.append(f"    if pd.api.types.is_datetime64_any_dtype(df['{col}']): df['{col}_is_weekend'] = (df['{col}'].dt.weekday >= 5).astype(int)")

    # 2. 交互作用特徴量
    interaction_rules = rules.get('interactions', [])
    for interaction in interaction_rules:
        cols = interaction.get('columns', [])
        if len(cols) >= 2:
            new_col_name = '_x_'.join(cols)
            # 生成するコード内でカラムの存在を確認するif文を作成
            col_existence_check = " and ".join([f"'{c}' in df.columns" for c in cols])
            
            code_lines.append(f"\n# [{', '.join(cols)}] から交互作用特徴量を生成")
            code_lines.append(f"if {col_existence_check}:")
            # 以下のコードはifブロック内にインデントして生成する
            str_cols_to_join = " + '_' + ".join([f"df['{c}'].astype(str)" for c in cols])
            code_lines.append(f"    df['{new_col_name}'] = {str_cols_to_join}")

    # 3. 数値特徴量の変換
    numeric_rules = rules.get('numeric_transforms', {})
    for col, transform in numeric_rules.items():
        method = transform.get('method')
        code_lines.append(f"\n# [{col}] 列の数値変換: {method}")

        check_str = f"if '{col}' in df.columns and pd.api.types.is_numeric_dtype(df['{col}']):"

        if method == 'log':
            code_lines.append(check_str)
            code_lines.append(f"    # 0や負の値に対応するため np.log1p を使用")
            # SettingWithCopyWarningを回避するため .loc を使用
            code_lines.append(f"    df['{col}_log'] = np.log1p(df.loc[df['{col}'] >= 0, '{col}'])")
        elif method == 'binning':
            bins = transform.get('bins', 4)
            code_lines.append(check_str)
            code_lines.append(f"    df['{col}_binned'] = pd.cut(df['{col}'], bins={bins}, labels=False, include_lowest=True)")

    code_lines.append("\n# --- 特徴量エンジニアリング完了 ---")
    return "\n".join(code_lines)

@data_analyzer_bp.route('/api/profile/<project_id>/interaction_suggestion', methods=['GET'])
def get_interaction_suggestion_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404

        profiling_result = None
        # 新しいBlob名フィールドを優先してチェック
        if blob_name := project_data.get("profilingResultBlobName"):
            try:
                result_bytes = download_blob_to_bytes(blob_name)
                if result_bytes:
                    profiling_result = json.loads(result_bytes.decode('utf-8'))
            except Exception as e:
                logger.error(f"Blob '{blob_name}' からのプロファイリング結果の読み込みに失敗: {e}", exc_info=True)
                # エラーが発生しても、後方互換性のために処理を続ける
        
        # 後方互換性のため、古いフィールドもチェック
        if profiling_result is None and "profilingResult" in project_data:
            profiling_result = project_data.get("profilingResult")

        if not profiling_result:
             return jsonify({"error": "プロファイリング結果が見つかりません。"}), 404

        analysis_purpose = project_data.get('analysisPurpose', '目的は設定されていません')
        
        column_summaries = profiling_result.get("column_summaries", [])
        categorical_columns = [
            col['name'] for col in column_summaries 
            if not pd.api.types.is_numeric_dtype(col.get('type','')) 
            and not pd.api.types.is_datetime64_any_dtype(col.get('type',''))
            and project_data.get('columnRoles', {}).get(col['name']) not in ['ID / キー', '自由記述（テキスト）']
        ]

        if len(categorical_columns) < 2:
            return jsonify({"suggestions": []})

        prompt = f"""あなたは、特徴量エンジニアリングに精通したデータサイエンティストです。
ユーザーの分析目的とデータカラムの情報を基に、クロス集計や分析モデルで特に有効そうな「交互作用特徴量」の組み合わせを提案してください。

### ユーザーの分析目的
{analysis_purpose}

### カテゴリ（文字列）型のカラム一覧
{json.dumps(categorical_columns, ensure_ascii=False)}

### あなたのタスク
上記の目的を達成するために、組み合わせることで新しい洞察が得られそうなカラムの組み合わせを**最大10個**まで提案してください。
**2つのカラムだけでなく、分析に有効であれば3つのカラムの組み合わせも提案に含めてください。**

### 【最重要】出力形式 (解説やマークダウンを含まない、単一のJSON配列のみ)
```json
[
  {{
    "columns": ["地域", "商品カテゴリ"],
    "description": "地域ごとの売れ筋商品を特定し、エリアマーケティング戦略に活用できます。"
  }},
  {{
    "columns": ["年代", "性別", "会員ランク"],
    "description": "詳細なデモグラフィック属性と会員ランクを組み合わせ、ターゲット層を細分化します。"
  }}
]
```"""

        ai_response_str, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME,
            contents=[types.Part.from_text(text=prompt)],
            max_output_tokens=65535,
            temperature=0.2,
            user_info=user_info,
            generation_config_override={"response_mime_type": "application/json"},
            feature_category="Insight Canvas"
        )

        suggestions = json.loads(ai_response_str)
        return jsonify({"suggestions": suggestions})

    except Exception as e:
        logger.error(f"交互作用特徴量の提案生成中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "AIによる提案の生成に失敗しました。"}), 500


@data_analyzer_bp.route('/api/profile/<project_id>/preprocessing_info', methods=['GET'])
def get_preprocessing_info_api(project_id):
    """前処理・特徴量エンジニアリングに必要なカラム情報を返すAPI"""
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404

        profiling_result = None
        # 新しいBlob名フィールドを優先してチェック
        if blob_name := project_data.get("profilingResultBlobName"):
            try:
                result_bytes = download_blob_to_bytes(blob_name)
                if result_bytes:
                    profiling_result = json.loads(result_bytes.decode('utf-8'))
                else:
                    return jsonify({"error": f"プロファイリング結果ファイル({blob_name})のダウンロードに失敗しました。"}), 404
            except Exception as e:
                logger.error(f"Blob '{blob_name}' からのプロファイリング結果の読み込みに失敗: {e}", exc_info=True)
                return jsonify({"error": "プロファイリング結果の読み込み中にエラーが発生しました。"}), 500
                
        # 後方互換性のため、古いフィールドもチェック
        elif "profilingResult" in project_data:
            profiling_result = project_data.get("profilingResult")    

        if not profiling_result:
            return jsonify({"error": "プロファイリング結果がまだ見つかりません。"}), 404

        column_summaries = profiling_result.get("column_summaries", [])
        column_roles = project_data.get("columnRoles", {})
        
        missing_cols = []
        categorical_cols = []
        numeric_cols = []
        text_cols = []
        datetime_cols = [] # <<< 追加

        for col in column_summaries:
            col_name = col.get("name")
            col_type = col.get("type", "")
            role = column_roles.get(col_name)

            if col.get("missing_percentage", 0) > 0:
                missing_cols.append({"name": col_name, "type": col_type})
            
            # 型判定のロジックをより厳密化
            if pd.api.types.is_datetime64_any_dtype(col_type):
                datetime_cols.append({"name": col_name, "type": col_type})
            elif pd.api.types.is_numeric_dtype(col_type):
                numeric_cols.append({"name": col_name, "type": col_type})
            elif role == '自由記述（テキスト）':
                text_cols.append({"name": col_name, "type": col_type})
            else:
                categorical_cols.append({"name": col_name, "type": col_type})

        return jsonify({
            "missing_value_columns": missing_cols,
            "categorical_columns": categorical_cols,
            "numeric_columns": numeric_cols,
            "text_columns": text_cols,
            "datetime_columns": datetime_cols, # <<< 追加
            "existing_rules": project_data.get("preprocessingRules", {})
        })

    except Exception as e:
        logger.error(f"前処理情報取得中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "情報の取得に失敗しました。"}), 500


@data_analyzer_bp.route('/api/project/<project_id>/preprocessing_rules', methods=['PATCH'])
def update_preprocessing_rules_api(project_id):
    """ユーザーが設定した前処理ルールを保存するAPI"""
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    rules = request.get_json()

    if not rules:
        return jsonify({"error": "ルールが指定されていません。"}), 400

    try:
        # 既存のルールを取得して、新しいルールで更新する
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404
            
        existing_rules = project_data.get("preprocessingRules", {})
        existing_rules.update(rules)
        
        update_analysis_project(project_id, user_id, {"preprocessingRules": existing_rules})
        
        return jsonify({"success": True, "message": "前処理ルールを保存しました。"})

    except Exception as e:
        logger.error(f"前処理ルールの保存中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "ルールの保存に失敗しました。"}), 500

@data_analyzer_bp.route('/api/profile/<project_id>/preprocessing_suggestion', methods=['GET'])
def get_preprocessing_suggestion_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404

        profiling_result = None
        # 新しいBlob名フィールドを優先してチェック (Blob Storage対応の追加)
        if blob_name := project_data.get("profilingResultBlobName"):
            try:
                result_bytes = download_blob_to_bytes(blob_name)
                if result_bytes:
                    profiling_result = json.loads(result_bytes.decode('utf-8'))
                else:
                    return jsonify({"error": f"プロファイリング結果ファイル({blob_name})のダウンロードに失敗しました。"}), 404
            except Exception as e:
                logger.error(f"Blob '{blob_name}' からのプロファイリング結果の読み込みに失敗: {e}", exc_info=True)
                return jsonify({"error": "プロファイリング結果の読み込み中にエラーが発生しました。"}), 500
                
        # 後方互換性のため、古いフィールドもチェック
        elif "profilingResult" in project_data:
            profiling_result = project_data.get("profilingResult")

        if not profiling_result or 'quality_alerts' not in profiling_result:
            return jsonify({"error": "品質アラートが見つかりません。"}), 404

        quality_alerts = profiling_result['quality_alerts']
        
        # AIへの指示をさらに厳格で安全なものに修正
        prompt = f"""あなたは優秀なデータサイエンティストです。以下のデータ品質に関するアラートを解決するための、安全なPython（pandas）コードを生成してください。

### 品質アラート
```json
{json.dumps(quality_alerts, indent=2, ensure_ascii=False)}
```

### 厳守すべきルール
1.  **出力は、`df`という既存のDataFrameを直接操作する一連のPythonコードのみ**とします。
2.  **サンプルデータの作成 (`pd.DataFrame(...)`)、関数定義 (`def ...`), `print()`文、`import`文は絶対に出力に含めないでください。**
3.  指摘された問題点のみを修正し、他のデータには一切触れないでください。
4.  各処理には、なぜその処理を行うのかを説明するコメントを必ず付けてください。
5.  元のDataFrame `df` を直接変更してください。コピーを作成する必要はありません。

### 良いコードの例（このような形式で出力してください）
```python
# 'age'列に欠損値があれば、その中央値で補完します。
if 'age' in df.columns and df['age'].isnull().any():
    median_age = df['age'].median()
    df['age'] = df['age'].fillna(median_age)

# 'sales'列の外れ値を99パーセンタイル値でクリッピングします。
if 'sales' in df.columns and not df['sales'].isnull().all():
    upper_bound = df['sales'].quantile(0.99)
    df['sales'] = df['sales'].clip(upper=upper_bound)
```

### あなたのタスク
上記のアラートを修正するための、**データ操作を行うPythonコードだけ**を生成してください。
"""
        
        generated_code, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME,
            contents=[types.Part.from_text(text=prompt)],
            max_output_tokens=65535,
            temperature=0.0,
            user_info=user_info,
            feature_category="Insight Canvas",
        )
        
        cleaned_code = generated_code.strip().lstrip('```python').rstrip('```').strip()

        return jsonify({"code": cleaned_code})

    except Exception as e:
        logger.error(f"前処理提案の生成中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "AIによるコード提案の生成に失敗しました。"}), 500


@data_analyzer_bp.route('/api/run_analysis/<project_id>', methods=['POST'])
def run_analysis_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    request_data = request.get_json()
    recipe = request_data.get('recipe', 'auto')
    insight_count = request_data.get('insight_count', 10) # デフォルトは10件

    task_id = f"task_{project_id}_analysis"
    task_data = {"id": task_id, "task_id": task_id, "status": "PENDING", "task_type": "data_analysis", "project_id": project_id, "user_id": user_id, "recipe": recipe}

    try:

        # 新しい分析を開始する前に、前回の分析結果をクリアする
        update_analysis_project(project_id, user_id, {
            "status": "analysis_started",
            "preprocessingSummary": None, # 前処理サマリーをクリア
            "preprocessingLog": [],       # 前処理ログをクリア
            "insights": {},               # インサイトをクリア
            "summaryReport": "",          # サマリーレポートをクリア
            "clusterAnalysisResult": None,
            "regressionAnalysisResult": None,
            "timeseriesForecastResult": None
        })

        container = get_tasks_container()
        if container:
            container.upsert_item(body=task_data)
        else:
            raise Exception("Tasks container is not available.")
    except Exception as e:
        logger.error(f"分析タスク {task_id} の作成/更新に失敗: {e}", exc_info=True)
        return jsonify({"error": "タスクの開始に失敗しました。"}), 500

    executor.submit(_run_analysis_task, task_id, project_id, user_id, recipe, insight_count)
    return jsonify({"task_id": task_id}), 202

def _process_text_columns_in_batch(df: pd.DataFrame, rules: dict, batch_size: int = 1000, log_callback=None):
    if not rules:
        return df

    def _log(message):
        if log_callback:
            log_callback(message)
        else:
            logger.info(message)

    df_processed = df.copy()

    for col_name, col_rules in rules.items():
        if col_name not in df_processed.columns:
            logger.warning(f"テキスト処理ルールに指定された列 '{col_name}' がデータに存在しません。")
            continue
        
        _log(f"列 '{col_name}' のテキスト処理を開始します。")
        
        categorize_rule = next((r for r in col_rules if r['type'] == 'categorize'), None)
        score_rule = next((r for r in col_rules if r['type'] == 'score'), None)

        if not categorize_rule and not score_rule:
            continue
            
        if categorize_rule:
            if not categorize_rule.get('categories'):
                try:
                    _log(f"列 '{col_name}' のカテゴリをAIが自動発見しています...")
                    sample_texts = df_processed[col_name].dropna().sample(n=min(300, len(df_processed[col_name].dropna())), random_state=1).tolist()
                    
                    discovery_prompt = f"""あなたは、大量のテキストから主要なテーマを抽出し、簡潔なカテゴリ名を作成する専門家です。
### 分析対象データ
これは「{col_name}」という項目に対する回答のサンプルです。
{json.dumps(sample_texts, ensure_ascii=False, indent=2)}

### あなたのタスク
上記のサンプル全体を読み解き、これらの回答を分類するための最適なカテゴリを**4〜6個**考案してください。
- 各カテゴリは、内容を的確に表す簡潔な日本語のラベルにしてください。
- 必ず「その他」のカテゴリを含めてください。
### 【最重要】出力形式
考案したカテゴリ名を、**カンマ区切りの文字列のみ**で回答してください。解説や番号付けは不要です。
### 出力例
UI/UX, 機能への言及, セキュリティ, 料金, その他"""

                    discovered_categories_str, _ = generate_with_gemini(
                        model_name=config.DATA_ANALYSIS_MODEL_NAME,
                        max_output_tokens=65535,
                        temperature=0.1,
                        contents=[types.Part.from_text(text=discovery_prompt)],
                        feature_category="Insight Canvas"
                    )
                    
                    if discovered_categories_str and not discovered_categories_str.startswith("[エラー"):
                        categorize_rule['categories'] = discovered_categories_str.strip()
                        _log(f"AIがカテゴリを発見しました: '{discovered_categories_str.strip()}'")
                    else:
                        raise Exception("AIによるカテゴリ発見に失敗しました。")

                except Exception as e_discover:
                    _log(f"列 '{col_name}' のカテゴリ発見中にエラーが発生しました。汎用カテゴリを使用します。")
                    categorize_rule['categories'] = "意見・感想, 質問, 要望, その他"
        
        temp_df = df_processed[[col_name]].copy()
        temp_df['__id'] = range(len(temp_df))
        texts_to_process = temp_df[temp_df[col_name].notna() & (temp_df[col_name] != '')]
        all_results = {}

        for i in range(0, len(texts_to_process), batch_size):
            chunk = texts_to_process.iloc[i:i + batch_size]
            input_json_for_ai = chunk.apply(lambda row: {"id": row['__id'], "text": str(row[col_name])}, axis=1).tolist()
            
            _log(f"列 '{col_name}' のテキストをバッチ処理中... ({i+1}件目～)")
            
            prompt_parts = ["あなたは、複数のテキストデータを一度に処理する、高度に最適化されたテキスト解析AIです。"]
            output_schema_parts = ['"id": (Integer)']
            
            if categorize_rule and categorize_rule.get('categories'):
                prompt_parts.append(f"\n1. 【カテゴリ分類】: 各テキストを以下のカテゴリから1つ選択してください。\n    カテゴリ: {categorize_rule['categories']}")
                output_schema_parts.append(f'"category": (String, 選択肢: {categorize_rule["categories"]})')
            if score_rule:
                prompt_parts.append(f"\n2. 【スコアリング】: 各テキストを以下の基準で評価してください。\n    基準: {score_rule['criteria']}")
                output_schema_parts.append('"score": (Number)')

            prompt_parts.append(f"\n### 入力データ (JSON)\n{json.dumps(input_json_for_ai, ensure_ascii=False, indent=2)}")
            prompt_parts.append("\n### 【最重要】出力形式\nあなたの応答は、入力データの各IDに対応する結果を含む、単一のJSON配列の文字列**のみ**でなければなりません。")
            prompt_parts.append("- テキストが空、または分類不能な場合は、categoryやscoreを `null` としてください。")
            prompt_parts.append("- 解説やマークダウンの囲みは**絶対に含めないでください。**")
            prompt_parts.append(f"- 出力JSONのスキーマ: `{{ {', '.join(output_schema_parts)} }}`")

            final_prompt = "\n".join(prompt_parts)
            
            try:
                response_str, _ = generate_with_gemini(
                    model_name=config.DATA_ANALYSIS_MODEL_NAME,
                    max_output_tokens=65535,
                    temperature=0.0,
                    contents=[types.Part.from_text(text=final_prompt)],
                    generation_config_override={"response_mime_type": "application/json"},
                    feature_category="Insight Canvas"
                )
                json_match = re.search(r'\[.*\]', response_str, re.DOTALL)
                if json_match:
                    batch_results = json.loads(json_match.group(0))
                    for res in batch_results:
                        all_results[res['id']] = res
                else:
                    logger.warning(f"列 '{col_name}' のバッチ処理でAIが有効なJSON配列を返しませんでした。")
            except Exception as e:
                _log(f"列 '{col_name}' のバッチ処理中にエラーが発生しました。")
                continue

        if categorize_rule and categorize_rule.get('new_column'):
            new_col_cat = categorize_rule['new_column']
            df_processed[new_col_cat] = df_processed.index.map(lambda idx: all_results.get(idx, {}).get('category'))
        
        if score_rule and score_rule.get('new_column'):
            new_col_score = score_rule['new_column']
            df_processed[new_col_score] = df_processed.index.map(lambda idx: all_results.get(idx, {}).get('score'))
        
        _log(f"列 '{col_name}' のテキスト処理が完了しました。")

    return df_processed


def _get_newly_created_columns(df_before: pd.DataFrame, df_after: pd.DataFrame) -> list:
    """前処理の前後でDataFrameを比較し、新しく追加された列名のリストを返す"""
    return [col for col in df_after.columns if col not in df_before.columns]

def _run_prediction_task(df: pd.DataFrame, target_variable: str, model_type: str, project_id: str, n_factors: int = 5) -> tuple[dict | None, list]:
    """
    データフレーム、目的変数、モデルタイプを受け取り、予測モデルを構築・評価する。
    (改訂版 v2) 再現コード生成位置の安全性確保。
    """
    logger.info(f"プロジェクト {project_id} のための予測モデル構築タスクを開始します。")
    logger.info(f"  - 目的変数: {target_variable}")
    logger.info(f"  - モデルタイプ: {model_type}")

    try:
        # --- 1. データ準備 ---
        if target_variable not in df.columns:
            raise ValueError(f"目的変数 '{target_variable}' がデータに存在しません。")

        # 日付列を除外
        df_for_prediction = df.select_dtypes(exclude=['datetime64[ns]', 'datetime64[ns, UTC]'])
        # 目的変数の欠損を除外
        df_for_prediction = df_for_prediction.dropna(subset=[target_variable])
        
        if df_for_prediction.empty:
            raise ValueError("目的変数の欠損値を除外した結果、データが0件になりました。")
        
        if target_variable not in df_for_prediction.columns:
            raise ValueError(f"目的変数 '{target_variable}' がデータから失われました。")

        X = df_for_prediction.drop(columns=[target_variable])
        y = df_for_prediction[target_variable]

        numeric_features = X.select_dtypes(include=[np.number, 'bool']).columns.tolist()
        categorical_features = X.select_dtypes(include=['object', 'category']).columns.tolist()
        
        model_original_features = X.columns.tolist()

        # --- 2. 前処理パイプラインの構築 ---
        numeric_transformer = Pipeline(steps=[
            ('imputer', SimpleImputer(strategy='median')),
            ('scaler', StandardScaler())
        ])
        categorical_transformer = Pipeline(steps=[
            ('imputer', SimpleImputer(strategy='most_frequent')),
            ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
        ])
        
        preprocessor = ColumnTransformer(
            transformers=[
                ('num', numeric_transformer, numeric_features),
                ('cat', categorical_transformer, categorical_features)
            ],
            remainder='drop' 
        )

        # --- 3. モデルの選択 ---
        if model_type == 'regression':
            model = RandomForestRegressor(random_state=42, n_jobs=1)
            algorithm_name = "Random Forest Regressor"
        else:
            model = RandomForestClassifier(random_state=42, n_jobs=1)
            algorithm_name = "Random Forest Classifier"
        
        pipeline = Pipeline(steps=[('preprocessor', preprocessor), ('classifier', model)])

        # --- 4. 学習の実行 ---
        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        
        # UIフォーム生成用の特徴量詳細情報を作成
        feature_details = {}
        MAX_CATEGORIES_FOR_SELECT = 50
        
        for feature in model_original_features:
            if feature not in X_train.columns: continue

            col_type = X_train[feature].dtype
            if pd.api.types.is_numeric_dtype(col_type):
                feature_details[feature] = {
                    "type": "numeric", 
                    "min": float(X_train[feature].min()),
                    "max": float(X_train[feature].max()), 
                    "median": float(X_train[feature].median())
                }
            else:
                unique_values = X_train[feature].dropna().unique().tolist()
                mode_value = X_train[feature].mode().iloc[0] if not X_train[feature].mode().empty else ""
                
                if len(unique_values) > 50: 
                    feature_details[feature] = {"type": "text", "mode": str(mode_value)}
                else:
                    feature_details[feature] = {
                        "type": "categorical", 
                        "categories": sorted([str(v) for v in unique_values]),
                        "mode": str(mode_value)
                    }
        
        pipeline.fit(X_train, y_train)

        # --- 5. 評価指標の計算 ---
        y_pred = pipeline.predict(X_test)
        if model_type == 'regression':
            accuracy_score_val = pipeline.score(X_test, y_test)
            error_metric_val = mean_absolute_percentage_error(y_test, y_pred)
        else:
            accuracy_score_val = pipeline.score(X_test, y_test)
            error_metric_val = f1_score(y_test, y_pred, average='weighted')

        # --- 6. 特徴量重要度の計算 ---
        top_factors = []
        original_feature_importances = {}
        try:
            transformed_feature_names = pipeline.named_steps['preprocessor'].get_feature_names_out()
            importances = pipeline.named_steps['classifier'].feature_importances_
            
            for name, imp in zip(transformed_feature_names, importances):
                original_name = re.sub(r'^(num__|cat__|remainder__)', '', name)
                if any(cat_feat in name for cat_feat in categorical_features):
                    original_name = re.sub(r'_[^_]*$', '', original_name)
                
                original_feature_importances.setdefault(original_name, 0)
                original_feature_importances[original_name] += imp

            sorted_original_factors = sorted(original_feature_importances.items(), key=lambda item: item[1], reverse=True)
            top_factors = [name for name, _ in sorted_original_factors[:n_factors]]

        except Exception as e_feat:
            logger.error(f"特徴量重要度の計算に失敗: {e_feat}", exc_info=True)

        # --- 7. グラフ生成 ---
        feature_importance_chart_url, chart_blob_name = None, None
        if original_feature_importances:
            try:
                df_for_plot = pd.DataFrame(
                    list(original_feature_importances.items()),
                    columns=['feature', 'importance']
                ).sort_values('importance', ascending=False).head(20)

                plt.figure(figsize=(10, 8))
                sns.barplot(x='importance', y='feature', data=df_for_plot)
                plt.title(f'Feature Importance for {target_variable} (Top 20)')
                plt.tight_layout()
                
                buf = io.BytesIO()
                plt.savefig(buf, format='png')
                plt.close()
                buf.seek(0)
                
                chart_blob_name = f"charts/{project_id}/feature_importance_{uuid.uuid4().hex}.png"
                upload_bytes_to_blob(buf.getvalue(), chart_blob_name)
            except Exception as e_chart:
                logger.error(f"グラフ生成に失敗: {e_chart}", exc_info=True)

        # --- 8. 解説生成 ---
        top_feature_name = top_factors[0] if top_factors else "N/A"
        if model_type == 'regression':
            commentary = f"AIによる自動生成された解説: 回帰モデルを構築し、テストデータで評価した結果、精度(R²)は **{accuracy_score_val:.2f}** となりました。また、予測値が実際の値から平均して **{error_metric_val:.2%}** 程度乖離していることを示すMAPE（平均絶対パーセント誤差）が算出されました。予測に最も影響を与えているのは「**{top_feature_name}**」です。"
        else:
            commentary = f"AIによる自動生成された解説: 分類モデルを構築し、テストデータで評価した結果、精度(正解率)は **{accuracy_score_val:.2f}** となりました。予測に最も影響を与えているのは「**{top_feature_name}**」です。"

        # --- 9. 保存 ---
        pipeline_blob_name = None
        try:
            pipeline_buffer = io.BytesIO()
            joblib.dump(pipeline, pipeline_buffer)
            pipeline_buffer.seek(0)
            pipeline_blob_name = f"models/{project_id}/pipeline_v1_{uuid.uuid4().hex[:8]}.joblib"
            upload_bytes_to_blob(pipeline_buffer.getvalue(), pipeline_blob_name)
        except Exception as e_save:
            logger.error(f"パイプラインの保存に失敗: {e_save}", exc_info=True)

        # --- 10. 再現用コードの生成 (安全な位置) ---
        reproduction_code = f"""# --- 予測モデル構築コード（再現用） ---
# 注意: このコードはシステム内部で実行された処理を再現するためのものです。
# 実際にはサーバー上の環境で実行されました。

import pandas as pd
from sklearn.model_selection import train_test_split
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.ensemble import RandomForest{'Regressor' if model_type=='regression' else 'Classifier'}
from sklearn.metrics import {'mean_absolute_percentage_error' if model_type=='regression' else 'f1_score'}

# 1. データの準備
target = '{target_variable}'
X = df.drop(columns=[target])
y = df[target]

# 数値・カテゴリ列の定義
numeric_features = X.select_dtypes(include=['number']).columns.tolist()
categorical_features = X.select_dtypes(include=['object', 'category']).columns.tolist()

# 2. 前処理パイプライン
numeric_transformer = Pipeline(steps=[
    ('imputer', SimpleImputer(strategy='median')),
    ('scaler', StandardScaler())
])

categorical_transformer = Pipeline(steps=[
    ('imputer', SimpleImputer(strategy='most_frequent')),
    ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
])

preprocessor = ColumnTransformer(
    transformers=[
        ('num', numeric_transformer, numeric_features),
        ('cat', categorical_transformer, categorical_features)
    ]
)

# 3. モデル定義
model = RandomForest{'Regressor' if model_type=='regression' else 'Classifier'}(random_state=42, n_jobs=-1)

clf = Pipeline(steps=[('preprocessor', preprocessor), ('classifier', model)])

# 4. 学習と評価
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
clf.fit(X_train, y_train)

print(f"Model Score: {{clf.score(X_test, y_test):.4f}}")
"""

        model_id = f"model_{project_id}_{uuid.uuid4().hex[:8]}"
        result = {
            "id": model_id, 
            "model_id": model_id, 
            "project_id": project_id,
            "projectName": df.attrs.get('name', project_id), 
            "target_variable": target_variable,
            "model_type": model_type, 
            "algorithm": algorithm_name, 
            "accuracy_score": float(accuracy_score_val),
            "error_metric": float(error_metric_val), 
            "feature_names": model_original_features,
            "feature_details": feature_details, 
            "commentary": commentary,
            "feature_importance_chart_blob_name": chart_blob_name,
            "pipeline_blob_name": pipeline_blob_name,
            "createdAt": datetime.now(timezone.utc).isoformat(), 
            "version": 1, 
            "error": None,
            "code": reproduction_code, # ★コードを含める
            "type": "PredictionModel",
        }
        
        return result, top_factors

    except Exception as e:
        logger.error(f"予測モデルタスク {project_id} でエラー: {e}", exc_info=True)
        return {"error": str(e)}, []

# === インフォグラフィック生成ヘルパー関数 ===
def _generate_and_save_infographic(summary_text, project_id, user_id):
    """
    分析サマリーを分割し、2枚のインフォグラフィック画像を生成してBlobに保存する。
    戻り値: (blob_name_1, blob_name_2) のタプル
    """
    try:
        logger.info(f"[{project_id}] インフォグラフィックの生成を開始します。")

        # 画像生成(サブ呼び出し)のusageを1つの親ログに集約する
        parent_log_blob_name = ensure_parent_log_blob_name(
            user_info={"userId": user_id},
            feature_category="Insight Canvas (Infographic)",
            model_name=config.IMAGE_GENERATION_MODEL_NAME,
        )
        
        # レポートをセクションに分割
        parsed_report = _parse_summary_report(summary_text)
        
        # パート1: サマリーとストーリー
        part1_text = f"## エグゼクティブサマリー\n{parsed_report.get('エグゼクティブサマリー', '')}\n\n## 分析のストーリー\n{parsed_report.get('分析のストーリー', '')}"
        
        # パート2: アクションプラン
        part2_text = f"## 推奨される次のアクション\n{parsed_report.get('推奨される次のアクション', '')}"
        
        blob_name_1 = None
        blob_name_2 = None

        # --- 画像1の生成 ---
        if part1_text.strip():
            prompt1 = f"""
            SYSTEM COMMAND: Create a high-quality, professional infographic image (Part 1 of 2).
            
            **OBJECTIVE:** Visualize the "Executive Summary" and "Analysis Story" sections of a business report.
            **STYLE:** Modern, clean, corporate data visualization, flat design, white background.
            **ASPECT RATIO:** 16:9 (Landscape).
            
            **CONTENT:**
            {part1_text[:1500]} 
            
            **RULES:** Use charts, icons, and layout to structure the information. Text MUST be in Japanese.
            """
            blob_name_1 = _generate_image_from_prompt(
                prompt1,
                project_id,
                user_id,
                "infographic_part1.png",
                parent_log_id=parent_log_blob_name,
            )

        # --- 画像2の生成 ---
        if part2_text.strip():
            prompt2 = f"""
            SYSTEM COMMAND: Create a high-quality, professional infographic image (Part 2 of 2).
            
            **OBJECTIVE:** Visualize the "Recommended Actions" (Action Plan) section of a business report.
            **STYLE:** Modern, clean, corporate data visualization, flat design, white background. Use a table-like or card layout for actions.
            **ASPECT RATIO:** 16:9 (Landscape).
            
            **CONTENT:**
            {part2_text[:1500]}
            
            **RULES:** Clearly visualize the action plan table/list. Text MUST be in Japanese.
            """
            blob_name_2 = _generate_image_from_prompt(
                prompt2,
                project_id,
                user_id,
                "infographic_part2.png",
                parent_log_id=parent_log_blob_name,
            )

        return blob_name_1, blob_name_2

    except Exception as e:
        logger.error(f"[{project_id}] インフォグラフィック生成中にエラー: {e}", exc_info=True)
        return None, None

def _generate_image_from_prompt(prompt, project_id, user_id, filename, parent_log_id=None):
    """(内部ヘルパー) プロンプトから画像を生成しBlobに保存"""
    try:
        part = types.Part.from_text(text=prompt)
        content = types.Content(role="user", parts=[part])

        # 親が確保できている場合は、ログBlobは作らずusageのみ親へ加算する
        skip_cosmos_log = bool(parent_log_id)

        response_iterator = generate_with_gemini(
            model_name=config.IMAGE_GENERATION_MODEL_NAME,
            contents=[content], # 修正: [Part] ではなく [Content] を渡す
            stream=True,
            user_info={"userId": user_id},
            feature_category="Insight Canvas (Infographic)",
            temperature=0.8,
            aspect_ratio="16:9",
            skip_cosmos_log=skip_cosmos_log,
            parent_log_id=parent_log_id,
        )
        full_base64_data = ""
        for chunk in response_iterator:
            if hasattr(chunk, 'image_data_url') and chunk.image_data_url:
                _, base64_data = chunk.image_data_url.split(",", 1)
                full_base64_data += base64_data
        
        if full_base64_data:
            image_bytes = base64.b64decode(full_base64_data)
            blob_name = f"results/{project_id}/{filename}"
            upload_bytes_to_blob(image_bytes, blob_name)
            logger.info(f"[{project_id}] 画像を保存しました: {blob_name}")
            return blob_name
    except Exception as e:
        logger.error(f"画像生成エラー ({filename}): {e}")
    return None

def _run_analysis_task(task_id, project_id, user_id, recipe, insight_count):
    try:
        update_task_in_cosmos(task_id, {"status": "PROCESSING", "status_message": "分析準備中...", "progress_percent": 5})
    except Exception as e:
        logger.error(f"タスク {task_id} の開始ステータス更新に失敗: {e}", exc_info=True)
        update_task_in_cosmos(task_id, {"status": "FAILURE", "error": "タスクの初期化に失敗しました。"})
        return

    # ログ蓄積用リスト
    preprocessing_logs = []

    # 内部関数: ログ保存処理 (Blob Storageへ保存し、Cosmos DBにはパスのみ記録)
    def _save_logs_to_blob():
        if not preprocessing_logs:
            return None
        try:
            log_bytes = json.dumps(preprocessing_logs, indent=2, ensure_ascii=False).encode('utf-8')
            blob_name = f"logs/{project_id}/preprocessing_log.json"
            upload_bytes_to_blob(log_bytes, blob_name)
            return blob_name
        except Exception as e:
            logger.error(f"ログのBlob保存中にエラー: {e}", exc_info=True)
            return None

    # 内部関数: メモリ上へのログ追加のみを行う (RU消費対策 + 上限キャップ)
    def add_log_and_update(message):
        nonlocal preprocessing_logs
        MAX_LOG_ENTRIES = 5000 
        
        if len(preprocessing_logs) >= MAX_LOG_ENTRIES:
            if len(preprocessing_logs) == MAX_LOG_ENTRIES:
                logger.warning(f"[{project_id}] ログ件数が上限({MAX_LOG_ENTRIES})に達したため、これ以上のメモリ記録を停止します。")
                preprocessing_logs.append({
                    "message": "--- ログ件数が上限に達しました。以降のログは省略されます ---", 
                    "timestamp": datetime.now(timezone.utc).isoformat()
                })
            logger.info(f"[{project_id}] (Omitted) {message}")
            return

        timestamp = datetime.now(timezone.utc).isoformat()
        log_entry = {"message": message, "timestamp": timestamp}
        preprocessing_logs.append(log_entry)
        logger.info(f"[{project_id}] {message}")

    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data: raise Exception("Project data not found.")
        
        add_log_and_update("分析コンテキストの準備を開始しました。")

        # まずDataFrameを読み込み、元の列名と小文字化後のマッピングを作成
        add_log_and_update("データファイルを読み込んでいます。")
        data_bytes = download_blob_to_bytes(project_data['dataFile']['blob_name'])
        df = _load_df_from_bytes_with_temp_file(data_bytes, project_data['dataFile']['filename'])
        if df is None: raise Exception("Failed to load data for analysis.")
        
        # オリジナルの列名と小文字化後のマッピングを作成
        column_map = {col: col.lower() for col in df.columns}
        # データフレームの列名をすべて小文字に統一
        df.columns = [col.lower() for col in df.columns]
        add_log_and_update("データカラム名を小文字に正規化しました。")

        analysis_purpose = project_data.get('analysisPurpose', 'データ全体の傾向を把握する')
        
        # 元の役割情報を小文字のキーに変換
        original_column_roles_info = project_data.get('columnRoles', {})
        column_roles_info = {column_map.get(k, k).lower(): v for k, v in original_column_roles_info.items()}

        df = _convert_numeric_to_datetime_based_on_role(df, column_roles_info)

        # 小文字化した後で日付変換を実行
        df = _try_convert_datetime_columns(df)

        # プロファイル結果も小文字化 (コピーを作成して安全に変更)
        original_profiling_result = project_data.get('profilingResult', {})
        profiling_result = copy.deepcopy(original_profiling_result)
        if 'column_summaries' in profiling_result:
            for col_summary in profiling_result['column_summaries']:
                if 'name' in col_summary:
                    col_summary['name'] = col_summary['name'].lower()
        
        # 前処理ルールの正規化
        original_preprocessing_rules = project_data.get('preprocessingRules', {})
        def normalize_rules_recursively(obj):
            if isinstance(obj, dict):
                new_dict = {}
                for k, v in obj.items():
                    new_key = column_map.get(k, k).lower()
                    new_dict[new_key] = normalize_rules_recursively(v)
                return new_dict
            elif isinstance(obj, list):
                return [normalize_rules_recursively(item) for item in obj]
            elif isinstance(obj, str):
                return column_map.get(obj, obj).lower()
            else:
                return obj
        preprocessing_rules = normalize_rules_recursively(original_preprocessing_rules)

        feature_engineering_rules = preprocessing_rules.get('feature_engineering', {})
        reference_files = project_data.get('referenceFiles', [])
        setup_suggestions = project_data.get('setupSuggestions', {})
        
        reference_summary = "なし"
        if reference_files:
            file_names = ", ".join([f.get('filename', '不明なファイル') for f in reference_files])
            reference_summary = f"以下の参考資料が提供されています: {file_names}"

        key_metrics_info = setup_suggestions.get('key_metrics', '未設定')
        alert_rules_info = setup_suggestions.get('alert_rules', '未設定')

        column_summary_with_memos = []
        for col_summary in profiling_result.get('column_summaries', []):
            memo = col_summary.get('memo', '').strip()
            if memo:
                column_summary_with_memos.append(f"- {col_summary.get('name')}: {memo}")
        column_memos_summary = "\n".join(column_summary_with_memos) if column_summary_with_memos else "なし"

        context_summary_for_ai = f"""
## プロジェクト全体のコンテキスト
- **分析の最終目的**: {analysis_purpose}
- **選択された分析レシピ**: {recipe}
- **特に注目すべき重要指標**: {key_metrics_info}
- **設定された警告ルール**: {alert_rules_info}
- **関連する参考資料**: {reference_summary}
- **ユーザーによるカラムへのメモ**:
{column_memos_summary}
- **ユーザーが設定した前処理ルール**: 
{json.dumps(preprocessing_rules, indent=2, ensure_ascii=False)}
- **各列の役割定義**:
{json.dumps(column_roles_info, indent=2, ensure_ascii=False)}
"""

        if feature_engineering_rules:
            update_task_in_cosmos(task_id, {"status_message": "特徴量を生成中...", "progress_percent": 8})
            add_log_and_update("ルールベースの特徴量エンジニアリングを実行しています...")
            feature_code = _generate_feature_engineering_code(feature_engineering_rules)
            if feature_code:
                exec_result_feat = execute_python_safely(feature_code, session_id=f"{project_id}_feat_eng", global_vars={'df': df.copy()})
                if exec_result_feat.get('error'):
                    logger.error(f"特徴量エンジニアリングでエラーが発生: {exec_result_feat.get('error')}")
                    add_log_and_update("特徴量エンジニアリングでエラーが発生したため、スキップしました。")
                else:
                    df_featured = exec_result_feat.get('final_scope', {}).get('df')
                    if df_featured is not None and isinstance(df_featured, pd.DataFrame):
                        df = df_featured
                        add_log_and_update("特徴量エンジニアリングが正常に完了しました。")
                    else:
                        add_log_and_update("特徴量エンジニアリングは実行されましたが、結果は適用されませんでした。")

        text_processing_rules = preprocessing_rules.get('text_processing', {})
        free_text_columns = [col for col, role in column_roles_info.items() if role == '自由記述（テキスト）']
        
        for col_name in free_text_columns:
            if col_name not in text_processing_rules:
                text_processing_rules[col_name] = [
                    {"type": "categorize", "categories": "", "new_column": f"{col_name}_自動カテゴリ"},
                    {"type": "score", "criteria": "文章の感情を-1(否定的)から1(肯定的)の範囲でスコアリング", "new_column": f"{col_name}_自動スコア"}
                ]
        
        if text_processing_rules:
            update_task_in_cosmos(task_id, {"status_message": "AIによるテキスト処理を実行中...", "progress_percent": 10})
            add_log_and_update("AIによるテキストデータの分類・スコアリングを開始します。")
            df = _process_text_columns_in_batch(df, text_processing_rules, batch_size=1000, log_callback=add_log_and_update)
            add_log_and_update("AIによるテキスト処理が完了しました。")

        update_task_in_cosmos(task_id, {"status_message": "AIが前処理コードを生成中...", "progress_percent": 15})
        add_log_and_update("AIに最適な前処理コードの生成を依頼しています。")
        
        # 初回のコード生成プロンプト
        preprocessing_prompt = f"""あなたは優秀なPythonデータサイエンティストです。以下のプロジェクト全体のコンテキストに基づき、データの前処理を行うPythonコードを生成してください。
{context_summary_for_ai}

### 【最重要】厳格な出力ルール
1.  あなたの応答は、**` ```python ... ``` `で囲まれたPythonコードブロックのみ**でなければなりません。
2.  解説など、コード以外のテキストは絶対に含めないでください。
3.  コードは既存の`df`変数を操作し、最終結果を`df_processed`に代入してください。
4.  **外部変数への依存禁止**: コード内で `user_preprocessing_rules` などの外部変数を参照してはいけません。提示されたルールの内容（欠損値処理の方法など）は、**具体的なPythonコード（fillna, dropnaなど）として書き下してください。**
5.  `import`や`pd.read_csv`は不要です。
6.  **【最重要】データ型のクリーニング**: 「9科目」「1,000円」のような**単位やカンマ、不要な文字列が含まれる数値列**に対して欠損値補完や統計量（中央値など）の計算を行う場合は、**必ず事前に正規表現等で不要な文字を除去し、`pd.to_numeric(..., errors='coerce')` で数値型に変換してから**計算処理を行ってください。文字列（Object型）のまま`median()`などを呼び出すとエラーになります。

### Pythonコードのみを出力してください:
"""
        # 初回のコード生成
        preprocessing_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=preprocessing_prompt)],feature_category="Insight Canvas")
        code_match = re.search(r'```python\n(.*?)```', preprocessing_code_response, re.DOTALL)
        if not code_match: 
             # 万が一コードブロックが見つからない場合は、生のテキストを試すか、エラーにする
             if "def " in preprocessing_code_response or "df" in preprocessing_code_response:
                 code_to_run_preprocessing = preprocessing_code_response.strip()
             else:
                 raise Exception(f"AIが有効な前処理コードブロックを生成できませんでした。")
        else:
             code_to_run_preprocessing = code_match.group(1).strip()
        
        add_log_and_update("生成された前処理コードを実行します。")
        update_task_in_cosmos(task_id, {"status_message": "前処理を実行中...", "progress_percent": 25})

        # --- 自動修復ループ (Retry Loop) ---
        max_retries = getattr(config, 'MAX_PYTHON_EXEC_RETRIES', 3)
        exec_result_preprocessing = {}
        
        for attempt in range(max_retries + 1):
            # コード実行
            exec_result_preprocessing = execute_python_safely(code_to_run_preprocessing, session_id=project_id, global_vars={'df': df.copy()})
            
            # 成功した場合
            if not exec_result_preprocessing.get('error'):
                logger.info(f"前処理コードの実行に成功しました (試行回数: {attempt+1})")
                add_log_and_update(f"前処理コードの実行が完了しました。")
                break
            
            # エラーが発生した場合
            error_msg = exec_result_preprocessing.get('error')
            logger.warning(f"前処理コード実行エラー (試行 {attempt+1}/{max_retries+1}): {error_msg}")
            
            # リトライ回数が残っている場合、AIに修正を依頼
            if attempt < max_retries:
                add_log_and_update(f"前処理でエラーが発生しました。AIが修正を試みます (試行 {attempt+1}/{max_retries})...")
                
                fix_prompt = f"""あなたはPythonコードのデバッグ専門家です。以下の前処理コードは実行時にエラーが発生しました。エラーメッセージを分析し、コードを修正してください。

### プロジェクトコンテキスト
{context_summary_for_ai}

### 失敗したコード
```python
{code_to_run_preprocessing}
```

### 発生したエラー
```
{error_msg}
```

### 修正のヒント
- `TypeError`や`ValueError`の場合、データの型が想定と異なっている可能性があります（例: '9科目'のような文字列が混入している）。`pd.to_numeric(..., errors='coerce')` や `.astype(str).str.replace(...)` を使って適切にクリーニングしてください。
- 外部変数は使用せず、すべてコード内で完結させてください。

### 【最重要】出力形式
修正後の**完全なPythonコード**のみを、**` ```python ... ``` `コードブロック**で囲んで出力してください。解説は不要です。
"""
                try:
                    fix_response, _ = generate_with_gemini(
                        model_name=config.DATA_ANALYSIS_MODEL_NAME, 
                        max_output_tokens=65535, 
                        contents=[types.Part.from_text(text=fix_prompt)],
                        feature_category="Insight Canvas"
                    )
                    fix_match = re.search(r'```python\n(.*?)```', fix_response, re.DOTALL)
                    if fix_match:
                        code_to_run_preprocessing = fix_match.group(1).strip()
                    else:
                        logger.warning("AIが修正コードの生成に失敗しました（フォーマット不正）。")
                        break # フォーマット不正ならリトライ中断
                except Exception as e:
                    logger.error(f"AIによるコード修正中にエラー: {e}")
                    break # AI呼び出しエラーなら中断
            else:
                # リトライ回数切れ
                add_log_and_update("前処理の自動修正に失敗しました。")
        
        # 前処理サマリーの作成
        preprocessing_summary = {
            "code": code_to_run_preprocessing,
            "stdout": "".join(item['content'] for item in exec_result_preprocessing.get('results', []) if item['type'] == 'stdout'),
            "error": exec_result_preprocessing.get('error')
        }

        update_analysis_project(project_id, user_id, {"preprocessingSummary": preprocessing_summary})
        
        if exec_result_preprocessing.get('error'): raise Exception(f"前処理でエラーが発生: {exec_result_preprocessing.get('error')}")
        
        df_processed = exec_result_preprocessing.get('final_scope', {}).get('df_processed')
        if df_processed is None:
             logger.warning("前処理コードで `df_processed` が定義されませんでした。元のDataFrameを使用します。")
             df_processed = df.copy()

        add_log_and_update("予測ターゲットの有無を確認しています...")
        target_variable, model_type = None, None
        primary_factors = []

        for col, role in column_roles_info.items():
            if role == '予測ターゲット（数値）':
                target_variable, model_type = col, 'regression'; break
            elif role == '予測ターゲット（カテゴリ）':
                target_variable, model_type = col, 'classification'; break
        
        if target_variable and model_type:
            update_task_in_cosmos(task_id, {"status_message": f"予測モデル({target_variable})を構築し、主要因を特定中...", "progress_percent": 35})
            add_log_and_update(f"予測ターゲット '{target_variable}' が設定されているため、{model_type} モデルの構築と主要因の特定を開始します。")
            
            try:
                prediction_result_meta, identified_factors = _run_prediction_task(df.copy(), target_variable, model_type, project_id)
                
                if identified_factors:
                    primary_factors = identified_factors
                    add_log_and_update(f"モデルから主要因として {', '.join(primary_factors)} を特定しました。")
                    context_summary_for_ai += f"\n- **システムが特定した重要因子**: {', '.join(primary_factors)}"
                else:
                    add_log_and_update("モデルから統計的に有意な主要因は特定できませんでした。")

                if prediction_result_meta and not prediction_result_meta.get('error'):
                    prediction_result_meta['userId'] = user_id
                    create_analysis_project(prediction_result_meta) 
                    add_log_and_update(f"予測モデルのメタデータをDBに保存しました (ID: {prediction_result_meta['id']})。")
                    update_analysis_project(project_id, user_id, {"regressionAnalysisResult": prediction_result_meta})
                    add_log_and_update("予測モデルの構築が完了しました。")
                else:
                    error_msg = prediction_result_meta.get('error', '不明なエラー') if prediction_result_meta else '結果がありません'
                    add_log_and_update(f"予測モデルの構築中にエラーが発生しました: {error_msg}")

            except Exception as e_pred:
                logger.error(f"予測モデルタスクの実行中に予期せぬエラーが発生しました: {e_pred}", exc_info=True)
                add_log_and_update(f"予測モデルの構築に失敗しました: {str(e_pred)}")
        else:
            add_log_and_update("予測ターゲットが設定されていないため、モデル構築とそれに伴う主要因特定はスキップします。")

        if isinstance(df_processed, pd.DataFrame):
            duplicated_cols = df_processed.columns.duplicated()
            if duplicated_cols.any():
                logger.warning(f"重複した列が検出されたため、自動的に削除します: {df_processed.columns[duplicated_cols].tolist()}")
                df_processed = df_processed.loc[:, ~duplicated_cols]
                    
        update_task_in_cosmos(task_id, {"status_message": "AIが分析計画を立案中...", "progress_percent": 40})
        
        df_processed_for_plan = df_processed
        df_info_buffer = io.StringIO(); df_processed_for_plan.info(buf=df_info_buffer)
        
        newly_created_columns = _get_newly_created_columns(df, df_processed)
        visualization_prompt_addition = ""
        if newly_created_columns:
            visualization_prompt_addition = f"""
### 【追加指示】可視化の必須要件
前処理の結果、以下の新しい列が生成されました: `{', '.join(newly_created_columns)}`
あなたの分析計画には、これらの**新しい列それぞれに対する可視化（カテゴリ列なら件数比較の棒グラフ、スコア列なら分布のヒストグラムなど）を必ず1つ以上含めてください。**
"""

        primary_factors_str = ', '.join([f"'{f}'" for f in primary_factors]) if primary_factors else "なし"
        target_variable_str = f"'{target_variable}'" if target_variable else "（目的変数が設定されていません）"

        timeseries_instructions = ""
        datetime_cols_for_prompt = df_processed.select_dtypes(include=['datetime64[ns]', 'datetime64[ns, UTC]']).columns.tolist()
        if datetime_cols_for_prompt:
            datetime_col = datetime_cols_for_prompt[0]
            ts_target_col = next(
                (col for col, role in column_roles_info.items() if role.startswith('予測ターゲット') and pd.api.types.is_numeric_dtype(df_processed[col])),
                None
            )
            if not ts_target_col:
                numeric_cols_for_ts = df_processed.select_dtypes(include=np.number).columns.tolist()
                if numeric_cols_for_ts: ts_target_col = numeric_cols_for_ts[0]
            
            if ts_target_col:
                timeseries_instructions = f"""
**2. 【特別指示】時系列・季節性分析 (以下の列を使用)**
   - **日付/時刻列**: `{datetime_col}`
   - **分析対象の数値列**: `{ts_target_col}`
   - 上記の2列が利用可能な場合、以下の分析を**必ず**含めてください:
     a. **統合的な季節性分析**: ...
     b. **【修正】自己相関分析**: 将来予測の可能性を探るため、以下の2つの分析を**それぞれ独立した分析計画として必ず生成してください**。
        - **自己相関係数(ACF)グラフ:** `plt.subplots()`で`fig, ax`を作成し、`plot_acf(データ, ax=ax)`を呼び出してください。
        - **偏自己相関係数(PACF)グラフ:** 同様に`plt.subplots()`で`fig, ax`を作成し、`plot_pacf(データ, ax=ax)`を呼び出してください。
        - **重要:** データは `df_processed.set_index('{datetime_col}')['{ts_target_col}'].dropna()` のように、**欠損を除外した1次元のPandas Series**として渡してください。
"""

        analysis_code_prompt = f"""あなたは、データからビジネスインサイトを抽出し、その結果を必ずグラフで可視化する、世界トップクラスのデータサイエンティストです。

### 【最重要】あなたの役割と厳格なルール
1.  **階層的な分析**: システムが特定した「主要因」を最優先で深掘りし、その後で他の変数を分析します。
2.  **原則、全分析をグラフ化**: あなたの仕事は、分析結果をテキスト（標準出力）で示すことではありません。**全ての分析計画は、必ず `matplotlib` や `seaborn` を使って結果をグラフとして描画する**Pythonコードを生成しなければなりません。**グラフを生成しない分析は禁止です。**
3.  **【超重要】分析データのJSON出力**: グラフを描画するために集計・加工したデータ（例: `groupby()`の結果など）は、**必ず`print(df_for_plot.to_json(orient='split', indent=2))`のように、JSON形式で標準出力してください。**
4.  **使用するDataFrame**: 分析には、**既にメモリ上に存在する `df_processed` という名前のpandas DataFrameのみを使用してください。**
5.  **【絶対禁止】データの再定義**: `df_processed = ...` や `pd.DataFrame(...)`、`pd.read_csv(...)` などでデータを再定義・上書き・読み込みすることは**固く禁じます**。これを行うと全データが失われます。必ず既存の `df_processed` をそのまま使ってください。
6.  **利用可能なライブラリ**: `pandas`, `numpy`, `matplotlib`, `seaborn`, `japanize_matplotlib`, `scipy.stats`, `statsmodels.api`, `statsmodels.graphics.tsaplots` のみ使用可能です。
7.  **グラフ描画の作法**:
    - コードの**先頭で必ず `japanize_matplotlib.japanize()` を呼び出し**、日本語の文字化けを防いでください。
    - グラフのタイトルや軸ラベルには、必ず日本語で分かりやすい名前を付けてください。
    - **絶対に `plt.show()` は呼び出さないでください。**
    - **seaborn警告回避**: `hue`引数を適切に使用してください。
8.  **JSONフォーマットの厳守**: あなたの応答は全体が単一のJSON配列でなければなりません。
9.  **統計的仮説検定の義務化**: グループ間の比較を行う場合は、**必ず適切な統計的仮説検定を実行してください。**
10. **自己完結したコード**: 各分析コードは、それ単体で完結するように記述してください。
11. **正しいメソッドの使用**: `df.sort_values(by='...')` など、正しい引数を使用してください。

### 依頼の概要
以下のコンテキストに従い、ユーザーの分析目的を達成するための**最低でも {insight_count} 個の多様な分析計画**をJSON形式で立案してください。

---
### 【プロジェクト全体のコンテキスト】
{context_summary_for_ai}
---
### 【前処理後のデータ情報】
{df_processed_for_plan.head(3).to_markdown()}
{df_info_buffer.getvalue()}
---
{visualization_prompt_addition}
---
### 【分析の進め方と特別指示】
**目的変数:** {target_variable_str}

**1. 最重要分析：主要因の深掘り (システム特定)**
   - **主要因リスト**: **[{primary_factors_str}]**
   - これらについて、多角的な分析を実行してください。

{timeseries_instructions}

**3. 補足分析：その他の変数の調査**

### 出力形式 (JSONのみ):
```json
[
  {{
    "title": "...",
    "description": "...",
    "code": "..."
  }}
]
```"""
        code_gen_response_str, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME,
            contents=[types.Part.from_text(text=analysis_code_prompt)],
            max_output_tokens=65535,
            temperature=0.2,
            generation_config_override={"response_mime_type": "application/json"},
            feature_category="Insight Canvas"
        )

        json_match = re.search(r'\[.*\]', code_gen_response_str, re.DOTALL)
        if json_match:
            json_string = json_match.group(0)
            corrected_json_string = re.sub(r'\\(?![/"\\bfnrtu])', r'\\\\', json_string)
            try:
                analysis_plan = json.loads(corrected_json_string)
            except json.JSONDecodeError as e:
                logger.error(f"JSONのパースに失敗しました。AIの応答(修正後): {corrected_json_string[:1000]}")
                raise Exception(f"AIが生成したJSONの解析に失敗しました。エラー: {e}") from e
        else:
            logger.error(f"AIが分析計画のJSONを生成できませんでした。AIの応答(先頭500文字): {code_gen_response_str[:500]}")
            raise Exception("AIが分析計画の生成に失敗しました。")
        
        update_task_in_cosmos(task_id, {"status_message": f"{len(analysis_plan)}個の分析を実行中...", "progress_percent": 50})
        
        update_analysis_project(project_id, user_id, {"insights": []})
        from concurrent.futures import ThreadPoolExecutor, as_completed

        insights = []
        with ThreadPoolExecutor(max_workers=max(4, os.cpu_count() or 1)) as executor:
            futures = {executor.submit(_execute_and_interpret_plan, plan, project_id, df_processed.copy()): plan for plan in analysis_plan}
            
            for i, future in enumerate(as_completed(futures)):
                try:
                    insight_result = future.result()
                    insights.append(insight_result)
                    update_task_in_cosmos(task_id, {
                        "status_message": f"分析を実行中... ({i + 1}/{len(analysis_plan)})",
                        "progress_percent": 50 + int(25 * (i + 1) / len(analysis_plan))
                    })
                except Exception as exc:
                    original_plan = futures[future]
                    logger.error(f'分析プラン "{original_plan.get("title", "無題")}" で例外が発生しました: {exc}', exc_info=True)

        # --- 時系列分析と異常検知の自動実行 ---
        update_task_in_cosmos(task_id, {"status_message": "時系列データの異常を検出中...", "progress_percent": 75})
        add_log_and_update("時系列データに異常な変動がないか確認しています...")
        anomaly_insight = _run_anomaly_detection(df_processed.copy(), project_id, column_roles_info, analysis_purpose)
        if anomaly_insight:
            insights.append(anomaly_insight)
            add_log_and_update("異常検知アラートを生成しました。")

        # 時系列予測
        update_task_in_cosmos(task_id, {"status_message": "時系列予測モデルを構築中...", "progress_percent": 77})
        add_log_and_update("時系列データが存在する場合、未来予測を実行します...")
        timeseries_model_meta = _run_timeseries_analysis(df_processed.copy(), project_id, user_id, project_data.get('projectName', ''))
        if timeseries_model_meta and not timeseries_model_meta.get('error'):
            create_analysis_project(timeseries_model_meta)
            add_log_and_update(f"時系列予測モデルのメタデータをDBに保存しました (ID: {timeseries_model_meta['id']})。")
            update_analysis_project(project_id, user_id, {"timeseriesForecastResult": timeseries_model_meta})
            add_log_and_update("時系列予測が完了しました。")
        else:
            add_log_and_update("時系列予測の対象となるデータがないか、モデル構築に失敗したため、スキップしました。")

        update_task_in_cosmos(task_id, {"status_message": "分析結果を整理・集約中...", "progress_percent": 80})
        
        def _categorize_insight(insight):
            title = insight.get('title', '').lower()
            if insight.get('is_anomaly_alert'): return 'anomaly_alerts'
            if 'category' in insight: return insight['category']
            if 'トレンド' in title or '推移' in title: return 'トレンド分析'
            if '季節' in title or '月別' in title: return '季節性分析'
            if 'acf' in title or 'pacf' in title or '自己相関' in title: return '自己相関分析'
            if '相関' in title or '関係' in title or '回帰' in title: return '相関・関係性分析'
            if '満足度' in title or '評価' in title or 'スコア' in title: return '評価・スコア分析'
            if 'クロス集計' in title or '内訳' in title: return 'クロス集計'
            if '比較' in title or '分布' in title: return '分布・比較分析'
            return 'その他のインサイト'

        primary_factor_insights = []
        remaining_insights = []
        if primary_factors:
            primary_insights_map = {factor: [] for factor in primary_factors}
            for insight in insights:
                found_in_primary = False
                for factor in primary_factors:
                    if factor in insight.get('title', '') or factor in insight.get('code', ''):
                        primary_insights_map[factor].append(insight)
                        found_in_primary = True
                        break
                if not found_in_primary:
                    remaining_insights.append(insight)
            
            for factor, factor_insights in primary_insights_map.items():
                if factor_insights:
                    primary_factor_insights.append({"factor": factor, "insights": factor_insights})
        else:
            remaining_insights = insights

        categorized_insights = {}
        for insight in remaining_insights:
            category = _categorize_insight(insight)
            if category not in categorized_insights:
                categorized_insights[category] = []
            categorized_insights[category].append(insight)
        
        anomaly_alerts = categorized_insights.pop('anomaly_alerts', [])

        final_insights_structure = {
            "primary_factor_insights": primary_factor_insights,
            "categorized_insights": [{"category": cat, "insights": ins_list} for cat, ins_list in categorized_insights.items()],
            "anomaly_alerts": anomaly_alerts
        }

        insights_bytes = json.dumps(final_insights_structure, ensure_ascii=False).encode('utf-8')
        insights_blob_name = f"results/{project_id}/insights.json"
        upload_bytes_to_blob(insights_bytes, insights_blob_name)
        update_analysis_project(project_id, user_id, {"insightsBlobName": insights_blob_name})

        update_task_in_cosmos(task_id, {"status_message": "クラスター分析を実行中...", "progress_percent": 80})
        cluster_analysis_result, df_clustered = _run_cluster_analysis(project_id, user_id, df_processed.copy(), context_summary_for_ai)

        if cluster_analysis_result:
            cluster_bytes = json.dumps(cluster_analysis_result, ensure_ascii=False).encode('utf-8')
            cluster_blob_name = f"results/{project_id}/cluster_analysis_result.json"
            upload_bytes_to_blob(cluster_bytes, cluster_blob_name)
            update_analysis_project(project_id, user_id, {"clusterAnalysisResultBlobName": cluster_blob_name})
        
        if df_clustered is not None:
            df_processed = df_clustered

        try:
            processed_data_buffer = io.BytesIO()
            df_processed.to_csv(processed_data_buffer, index=False, encoding='utf-8-sig')
            processed_data_buffer.seek(0)
            processed_blob_name = f"processed_data/{project_id}/processed_data.csv"
            upload_bytes_to_blob(processed_data_buffer.getvalue(), processed_blob_name)
            update_analysis_project(project_id, user_id, {"processedDataBlobName": processed_blob_name})
        except Exception as e_save:
            logger.error(f"分析後データの保存に失敗: {e_save}", exc_info=True)

        update_task_in_cosmos(task_id, {"status_message": "AIが分析ストーリーを生成中...", "progress_percent": 90})

        insights_summary_list = [{"title": i['title'], "description": i['description']} for i in insights if not i.get('error') and not i.get('is_anomaly_alert')]
        cluster_summary_for_ai = ""
        if cluster_analysis_result and not cluster_analysis_result.get('error'):
            cluster_summary_for_ai = "\n\n## クラスター分析による顧客セグメント分類\n"
            for cluster_info in cluster_analysis_result.get('summary', {}).get('clusters', []):
                cluster_summary_for_ai += f"- **{cluster_info.get('name')}**: {cluster_info.get('description')}\n"

        analysis_story_prompt = f"""あなたは、数多くの分析結果から本質を見抜き、経営層にも伝わるレポートを作成する、トップクラスのデータコンサルタントです。

### 分析の全体目的
{analysis_purpose}

### 発見されたインサイトのリスト
{json.dumps(insights_summary_list, indent=2, ensure_ascii=False)}

### クラスター分析の結果
{cluster_summary_for_ai}

### 【最重要】あなたのタスク
上記のすべての情報を統合し、以下の**3つのセクションから成る総括レポート**をMarkdown形式で生成してください。

---
### 1. エグゼクティブサマリー：最重要インサイト Top 3-5
まず、提供された全インサイトの中から、**ビジネスインパクトが最も大きい、あるいは最も想定外で注目すべき発見を3〜5個厳選**してください。

---
### 2. 分析のストーリー
次に、発見したインサイトを論理的に繋げ、データが語る物語を説明してください。
+ **【最重要指示】** ストーリーを記述する際は、**必ずインサイトリストから具体的なカテゴリ名、数値などを引用**してください。

---
### 3. 推奨される次のアクション
あなたは**事業企画部のマネージャー**です。上記の分析結果に基づき、明日から現場が実行できる具体的なアクションプランを**以下のMarkdownテーブル形式で**3つ提案してください。

**【重要】テーブル作成の厳守ルール**
1. **テーブルの前後には必ず空行を入れてください。**（空行がないと表として表示されません）
2. ヘッダー行の直下には、必ず**セパレータ行（`| :--- | :--- | ...`）**を含めてください。
3. カラム数は必ず以下の例と同じ5列にしてください。

(空行)
| アクション名 | 具体的なアクション内容 | 担当部署（案） | 主要評価指標 (KPI) | 期待される効果 |
| :--- | :--- | :--- | :--- | :--- |
| （例）若年層向けCP | 20代の購入率向上を狙い... | マーケティング部 | 新規獲得数 | 売上10%増 |
| （アクション2） | ... | ... | ... | ... |
| （アクション3） | ... | ... | ... | ... |
(空行)
"""

        final_report_md, _ = generate_with_gemini(
            model_name=config.DOC_GEN_MODEL_NAME, 
            max_output_tokens=65535, 
            contents=[types.Part.from_text(text=analysis_story_prompt)],
            feature_category="Insight Canvas"
        )
        if not final_report_md: raise Exception("AIによる分析ストーリーの生成に失敗しました。")
        final_report_md = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', final_report_md)

        update_task_in_cosmos(task_id, {"status_message": "分析結果のインフォグラフィックを生成中...", "progress_percent": 95})
        
        # 2つのBlob名を受け取る
        blob1, blob2 = _generate_and_save_infographic(final_report_md, project_id, user_id)
        
        project_updates = {
            "summaryReport": final_report_md, 
            "status": "analysis_completed"
        }
        
        if blob1: project_updates["infographicBlobName1"] = blob1
        if blob2: project_updates["infographicBlobName2"] = blob2
        
        update_analysis_project(project_id, user_id, project_updates)
        
        update_task_in_cosmos(task_id, {"status": "SUCCESS", "status_message": "分析が完了しました。", "progress_percent": 100})
    except Exception as e:
        logger.error(f"分析タスク {task_id} でエラー: {e}", exc_info=True)
        
        # エラー時もログ保存を試みる
        add_log_and_update(f"エラーが発生しました: {str(e)}")
        error_log_blob = _save_logs_to_blob()
        if error_log_blob:
             update_analysis_project(project_id, user_id, {"preprocessingLogBlobName": error_log_blob})

        update_task_in_cosmos(task_id, {"status": "FAILURE", "error": str(e), "status_message": "分析中にエラーが発生しました。"})


def _execute_and_interpret_plan(plan: dict, project_id: str, df: pd.DataFrame) -> dict:
    """単一の分析プランを実行し、信頼度を判定、結果を解釈する"""
    code_to_run = plan.get('code', '')
    original_description = plan.get('description', '説明なし')
    
    insight = {
        'id': f"insight_{uuid.uuid4().hex}",
        'title': plan.get('title', '無題の分析'),
        'description': original_description,
        'description_plain': None,
        'stdout': '',
        'error': None,
        'charts': [],
        'deep_dive_suggestions': [],
        'reliability_level': 'Reference', # デフォルトの信頼度レベル
        'code': code_to_run # 初期コードを設定
    }

    df_info_buffer = io.StringIO()
    df.info(buf=df_info_buffer)
    df_schema_for_correction = f"""
### 【重要】利用可能なDataFrame `df_processed` のスキーマ情報:
{df.head(3).to_markdown()}
{df_info_buffer.getvalue()}
"""

    for attempt in range(config.MAX_PYTHON_EXEC_RETRIES + 1):

        exec_result = execute_python_safely(code_to_run, session_id=project_id, global_vars={'df_processed': df.copy()})
        error = exec_result.get('error')
        results = exec_result.get('results', [])
        
        is_fatal_error = error and 'Traceback (most recent call last):' in error
        if not is_fatal_error:
            if error:
                logger.warning(f"実行中に警告または非致命的なエラーが出力されました: {error}")
            error = None
        
        if not error:
            # 実行結果から、統計テキストとJSONデータを分離する
            raw_stdout = "".join(item['content'] for item in results if item['type'] == 'stdout')
            analytics_data_json = None
            stats_text = raw_stdout
            
            match = re.search(r'--- BEGIN ANALYTICS DATA ---(.*?)--- END ANALYTICS DATA ---', raw_stdout, re.DOTALL)
            if match:
                json_str = match.group(1).strip()
                try:
                    # JSONデータをパースして保持
                    analytics_data_json = json.loads(json_str)
                    # 元のstdoutからJSON部分を削除して、統計テキストのみを残す
                    stats_text = raw_stdout.replace(match.group(0), "").strip()
                except json.JSONDecodeError:
                    logger.warning("Failed to parse analytics data JSON from stdout.")

            # insight辞書の更新
            insight['stdout'] = stats_text # JSONを除いたテキスト
            insight['charts'] = [item['content'] for item in results if item['type'] == 'chart_gcs_uri']
            # --- 変更点: 実際に成功したコードを保存 ---
            insight['code'] = code_to_run 
            # ---------------------------------------
            
            if not insight['charts'] and not insight['stdout'].strip():
                insight['error'] = "分析コードは実行されましたが、グラフやテキストなどの可視的な出力がありませんでした。"
                return insight

            # --- ステップ1: 信頼度レベルの判定 ---
            reliability_level = 'Reference'
            stdout_text = insight.get('stdout', '').lower()
            code_text = code_to_run.lower()
            
            # レベルA (確定的) の判定: p値が有意水準未満か
            p_value_match = re.search(r'(p-value|pvalue|p値)\s*[:=]\s*([\d.e-]+)', stdout_text)
            if p_value_match:
                try:
                    p_value = float(p_value_match.group(2))
                    if p_value < 0.05:
                        reliability_level = 'Confirmed'
                except (ValueError, IndexError):
                    pass # p値の抽出に失敗した場合は無視

            # レベルB (示唆的) の判定 (レベルAでない場合)
            if reliability_level == 'Reference':
                suggestive_keywords = ['corr', 'groupby', 'mean', 'median', 'ttest', 'f_oneway', 'boxplot', 'barplot', '比較', '相関']
                if any(keyword in code_text for keyword in suggestive_keywords):
                    reliability_level = 'Suggested'
            
            insight['reliability_level'] = reliability_level

            try:
                chart_info = "グラフは生成されませんでした。"
                if insight['charts']:
                    chart_info = f"はい、グラフが生成されました: {', '.join(insight['charts'])}"

                analytics_data_prompt_section = "分析データ(JSON): 分析コードからJSONデータは出力されませんでした。"
                if analytics_data_json:
                    analytics_data_prompt_section = f"### 分析データ (グラフの元データ):\n```json\n{json.dumps(analytics_data_json, indent=2, ensure_ascii=False)}\n```"

                # --- ステップ2: AIへの指示プロンプトを信頼度レベルで強化 ---
                interpretation_prompt = f"""あなたは、分析結果からビジネス上の結論を導き出す、データサイエンスに精通した専門家です。以下の情報から読み取れる**結論や発見**を、**地の文（レポートの本文）としてそのまま使える形式で**記述してください。

### 分析の目的: {original_description}
### 実行されたPythonコード:\n```python\n{code_to_run}\n```
{analytics_data_prompt_section}
### 実行結果のテキスト出力（統計量に注目）:\n```\n{stats_text[:2000]}\n```
### 生成されたグラフ: {chart_info}
### 【重要】分析結果の信頼度レベル: {reliability_level}

### 【最重要】厳格な出力ルール
1.  **【超重要】分析データ(JSON)を絶対的な正とする**: あなたの回答は、**必ず「分析データ(JSON)」に含まれる具体的なカテゴリ名、数値に基づいて**記述してください。JSONデータが存在しない場合にのみ、他の情報を参考にしてください。これにより、**存在しないカテゴリ名（例：「販路X」）を生成するハルシネーションを絶対に防いでください。**
2.  **信頼度レベルに応じた表現を厳守**:
    *   信頼度レベルが **'Confirmed'** の場合: **「～ということが統計的に確認されました」「～という有意な差があります」** のように、**断定的な表現**を使用してください。
    *   信頼度レベルが **'Suggested'** の場合: **「～という傾向が見られます」「～の可能性が示唆されます」** といった表現を用い、**「ただし、これが偶然でないことを証明するには、追加の統計的検証が推奨されます」** という旨の**注意喚起を必ず文末に含めてください。**
    *   信頼度レベルが **'Reference'** の場合: **「データによると～となっています」「～の分布はこのようになっています」** のように、事実を客観的に記述するに留め、**断定的な結論や因果関係には言及しないでください。**
3.  **具体的かつ定量的に記述**: 「分析データ(JSON)」を基に、**最も値が高い（または低い）カテゴリ名と、その具体的な数値を必ず引用**してください。「差がある」だけでなく、「**『Web』の受注率が58.2%で最も高く、『テレアポ』の15.3%を大幅に上回っています**」のように、比較対象を明確にして記述します。
4.  **結論のみを出力**: あなたの応答は、分析から導かれる**結論そのもの**でなければなりません。
5.  **統計的根拠を必ず含める**: 実行結果のテキスト出力に含まれる**p値や相関係数などの統計量**を可能な限り引用し、その発見が統計的にどの程度信頼できるかに言及してください。
6.  **無価値な回答の禁止**: **「分かりません」「断定できません」といった回答は絶対に禁止**します。統計的に有意でない場合でも、グラフから読み取れる視覚的な傾向を述べ、その上で「統計的な裏付けは弱い」または「サンプル数が少ないため参考値です」といった形で指摘してください。
7.  **前置きや接頭辞は厳禁**: 「はい、承知しました。」「インサイト：」といった前置きは**一切含めないでください。**
8.  **Markdownの活用**: 強調したいキーワードは `**` で囲んでください。

### あなたのタスク: 上記のルールを厳守し、**JSONデータを基にした、ハルシネーションのない、具体的で定量的な結論となる文章だけ**を生成してください。
"""
                expert_description, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, temperature=0.4, contents=[types.Part.from_text(text=interpretation_prompt)],feature_category="Insight Canvas")
                if expert_description and not expert_description.startswith("[エラー"):
                    insight['description'] = expert_description.strip()

                plain_language_prompt = f"""あなたは、複雑なデータ分析の結果を、専門知識のないビジネス担当者にも分かるように簡単な言葉で要約する専門家です。
### 専門家による分析結果
{insight['description']}
### あなたのタスク: 上記の専門的な解説を、以下のルールに従って**平易な言葉で書き直してください。**
### 【最重要】厳格なルール
1.  **結論ファースト**: 最も重要な発見や結論を、単刀直入に述べてください。
2.  **専門用語を避ける**: 「p値」「統計的に有意」などの専門用語は使わず、「偶然とは考えにくい」「明らかな差がある」といった表現に置き換えてください。
3.  **簡潔に**: 全体を1〜3文の短い文章で要約してください。
4.  **キーワードを強調**: 分析結果の中で、ユーザーに最も注目してほしい**キーワード**や**具体的な数値**を`**`で囲んで強調してください。
5.  **前置きは不要**: 「はい、承知しました」「つまり、」などの応答は含めず、書き直した文章だけを出力してください。
### あなたのタスク: 上記のルールを厳守し、**平易な言葉で書き直した文章だけ**を生成してください。
"""
                plain_description, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, temperature=0.5, contents=[types.Part.from_text(text=plain_language_prompt)],feature_category="Insight Canvas")
                if plain_description and not plain_description.startswith("[エラー"):
                    insight['description_plain'] = plain_description.strip()
                else:
                    insight['description_plain'] = insight['description']

                deep_dive_prompt = f"""あなたは、データアナリストの優秀なアシスタントです。ある分析によって、以下のインサイトが得られました。この発見を基に、さらに一歩踏み込んだ分析を行うための「次の質問」を提案してください。
### 元となったインサイト
- **タイトル:** {insight['title']}
- **結論:** {insight['description']}
### あなたのタスク: このインサイトから、ビジネスユーザーが次に関心を持ちそうな、具体的で分析可能な質問を3つ提案してください。平易な言葉で、簡潔な指示文の形でお願いします。
### 【最重要】出力形式 (解説を含まない、単一のJSON配列のみ):
```json
[
  "この顧客セグメントについて、年齢層別の内訳も表示して。",
  "この傾向は、直近3ヶ月のデータに絞っても見られますか？"
]
```"""
                suggestions_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, temperature=0.5, contents=[types.Part.from_text(text=deep_dive_prompt)], generation_config_override={"response_mime_type": "application/json"},feature_category="Insight Canvas")
                if suggestions_response and not suggestions_response.strip().startswith("[エラー"):
                    insight['deep_dive_suggestions'] = json.loads(suggestions_response)
            except Exception as e:
                logger.warning(f"インサイトの解釈中にエラーが発生しました: {e}", exc_info=True)
                if not insight.get('description_plain'):
                    insight['description_plain'] = insight.get('description', original_description)

            insight['error'] = None
            return insight

        # --- 変更点: 失敗した場合も実行したコードを保存 ---
        insight['error'] = error
        insight['code'] = code_to_run
        # ------------------------------------------------
        
        logger.warning(f"分析「{insight['title']}」の実行に失敗 (試行 {attempt + 1}/{config.MAX_PYTHON_EXEC_RETRIES + 1})。エラー: {error}")
        
        if attempt < config.MAX_PYTHON_EXEC_RETRIES:
            try:
                fix_prompt = f"""あなたはPythonコードのデバッグの専門家です。以下のコードは実行時にエラーを発生させました。エラーメッセージと正しいデータスキーマを参考にコードを修正してください。
{df_schema_for_correction}
### 目的: {original_description}
### 失敗したコード:\n```python\n{code_to_run}\n```
### エラーメッセージ:\n```\n{error}\n```
### 【最重要】厳格な出力ルール
1.  あなたの応答は、修正後の完全なPythonコードを含む**` ```python ... ``` `コードブロックのみ**でなければなりません。
2.  **必ず上記のスキーマ情報に存在する正しいカラム名のみを使用してください。**
3.  解説や謝罪など、コード以外のテキストは**絶対に含めないでください。**
### 修正後のPythonコード:"""

                fixed_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=fix_prompt)],feature_category="Insight Canvas")
                
                code_match = re.search(r'```python\n(.*?)```', fixed_code_response, re.DOTALL)
                if code_match:
                    code_to_run = code_match.group(1).strip()
                    logger.info(f"AIがコードの修正案を生成しました。リトライします。")
                    continue
                else:
                    logger.warning("AIは修正コードを生成できませんでした。リトライを中断します。")
                    break
            except Exception as e:
                logger.error(f"コード修正のAI呼び出し中にエラー: {e}", exc_info=True)
                break
        
    return insight

@data_analyzer_bp.route('/api/insights/<project_id>', methods=['GET'])
def get_insights_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    project_data = get_analysis_project(project_id, user_id)
    if not project_data: return jsonify({"error": "プロジェクトが見つかりません。"}), 404
    
    insights_data = {}
    cluster_analysis_result = None

    # インサイトをBlobから読み込む
    if insights_blob_name := project_data.get("insightsBlobName"):
        try:
            insights_bytes = download_blob_to_bytes(insights_blob_name)
            if insights_bytes:
                insights_data = json.loads(insights_bytes.decode('utf-8'))
        except Exception as e:
            logger.error(f"Blob '{insights_blob_name}' からのインサイト読み込みに失敗: {e}", exc_info=True)
    # 後方互換性
    elif "insights" in project_data:
        insights_data = project_data.get("insights", {})

    # クラスター分析結果をBlobから読み込む
    if cluster_blob_name := project_data.get("clusterAnalysisResultBlobName"):
        try:
            cluster_bytes = download_blob_to_bytes(cluster_blob_name)
            if cluster_bytes:
                cluster_analysis_result = json.loads(cluster_bytes.decode('utf-8'))
        except Exception as e:
            logger.error(f"Blob '{cluster_blob_name}' からのクラスター分析結果読み込みに失敗: {e}", exc_info=True)
    # 後方互換性
    elif "clusterAnalysisResult" in project_data:
        cluster_analysis_result = project_data.get("clusterAnalysisResult")
    
    # --- 【修正箇所】前処理ログをBlobから読み込む ---
    preprocessing_log = []
    if log_blob_name := project_data.get("preprocessingLogBlobName"):
        try:
            log_bytes = download_blob_to_bytes(log_blob_name)
            if log_bytes:
                preprocessing_log = json.loads(log_bytes.decode('utf-8'))
        except Exception as e:
            logger.error(f"Blob '{log_blob_name}' からの前処理ログ読み込みに失敗: {e}", exc_info=True)
    # 後方互換性: BlobになくCosmosに直接ある場合
    elif "preprocessingLog" in project_data:
        preprocessing_log = project_data.get("preprocessingLog", [])
    # ---------------------------------------------

    summary_report = project_data.get("summaryReport", "")
    preprocessing_summary = project_data.get("preprocessingSummary", {})
    regression_analysis_result = project_data.get("regressionAnalysisResult")
    timeseries_forecast_result = project_data.get("timeseriesForecastResult")
    
    infographic_urls = []
    if blob1 := project_data.get("infographicBlobName1"):
        if url1 := get_blob_sas_url(blob1): infographic_urls.append(url1)
    # 後方互換: 古いキーがあればそれもチェック
    elif blob_old := project_data.get("infographicBlobName"):
        if url_old := get_blob_sas_url(blob_old): infographic_urls.append(url_old)

    if blob2 := project_data.get("infographicBlobName2"):
        if url2 := get_blob_sas_url(blob2): infographic_urls.append(url2)

    gcs_bucket_name = os.environ.get("GOOGLE_CLOUD_STORAGE_BUCKET_NAME")

    def _process_insight_list(insight_list):
        """インサイトのリストをループし、各インサイトに署名付きURLを追加するヘルパー関数"""
        if not isinstance(insight_list, list):
            return
            
        for insight in insight_list:
            chart_urls = []
            gcs_blob_names = insight.get('charts', [])
            for blob_name in gcs_blob_names:
                # GCSの場合
                if gcs_bucket_name and blob_name.startswith(f'gs://{gcs_bucket_name}/'):
                    blob_path_only = blob_name.replace(f'gs://{gcs_bucket_name}/', '')
                    signed_url = get_gcs_signed_url(blob_path_only, bucket_name=gcs_bucket_name)
                    if signed_url:
                        chart_urls.append(signed_url)
                # Azure Blob Storageの場合 (httpを含まないパス形式を想定)
                elif blob_name and not blob_name.startswith('http') and not blob_name.startswith('gs://'):
                     sas_url = get_blob_sas_url(blob_name)
                     if sas_url:
                         chart_urls.append(sas_url)
                # 既にhttp URLの場合はそのまま
                elif blob_name and blob_name.startswith('http'):
                     chart_urls.append(blob_name)

            insight['chart_urls'] = chart_urls

    # 新しい階層化されたデータ構造に対応
    if isinstance(insights_data, dict):
        # 1. 主要因インサイトの処理
        for factor_group in insights_data.get("primary_factor_insights", []):
            _process_insight_list(factor_group.get("insights", []))
        
        # 2. その他のインサイトの処理
        for category_group in insights_data.get("categorized_insights", []):
            _process_insight_list(category_group.get("insights", []))
        
        # 3. 異常検知アラートの処理
        _process_insight_list(insights_data.get("anomaly_alerts", []))
    
    # 後方互換性のため、古いフラットなリスト形式も念のため処理
    elif isinstance(insights_data, list):
        _process_insight_list(insights_data)
    
    if cluster_analysis_result and isinstance(cluster_analysis_result, dict):
        # elbow_chart_urlとcluster_chart_urlの両方を署名付きURLに変換する
        for key in ['elbow_chart_url', 'cluster_chart_url']:
            if uri := cluster_analysis_result.get(key):
                # GCS URI
                if gcs_bucket_name and isinstance(uri, str) and uri.startswith(f'gs://{gcs_bucket_name}/'):
                    blob_path_only = uri.replace(f'gs://{gcs_bucket_name}/', '')
                    signed_url = get_gcs_signed_url(blob_path_only, bucket_name=gcs_bucket_name)
                    if signed_url:
                        cluster_analysis_result[key] = signed_url
                # Azure Blob
                elif isinstance(uri, str) and not uri.startswith('http') and not uri.startswith('gs://'):
                    sas_url = get_blob_sas_url(uri)
                    if sas_url:
                        cluster_analysis_result[key] = sas_url
    
    if regression_analysis_result and isinstance(regression_analysis_result, dict):
        # 1. Azure Blob Storage のパスがある場合 (新ロジック)
        if blob_name := regression_analysis_result.get('feature_importance_chart_blob_name'):
            sas_url = get_blob_sas_url(blob_name)
            if sas_url:
                regression_analysis_result['feature_importance_chart_url'] = sas_url
        
        # 2. GCS URIがある場合 (旧ロジック互換)
        elif gcs_bucket_name:
            gcs_uri = regression_analysis_result.get('feature_importance_chart_url_gcs')
            if gcs_uri and isinstance(gcs_uri, str) and gcs_uri.startswith(f'gs://{gcs_bucket_name}/'):
                blob_path_only = gcs_uri.replace(f'gs://{gcs_bucket_name}/', '')
                signed_url = get_gcs_signed_url(blob_path_only, bucket_name=gcs_bucket_name)
                if signed_url:
                    regression_analysis_result['feature_importance_chart_url'] = signed_url

    # 時系列予測モデルのグラフURL処理
    if timeseries_forecast_result and isinstance(timeseries_forecast_result, dict):
        # GCS
        if gcs_bucket_name:
            gcs_uri = timeseries_forecast_result.get('feature_importance_chart_url_gcs')
            if gcs_uri and isinstance(gcs_uri, str) and gcs_uri.startswith(f'gs://{gcs_bucket_name}/'):
                blob_path_only = gcs_uri.replace(f'gs://{gcs_bucket_name}/', '')
                signed_url = get_gcs_signed_url(blob_path_only, bucket_name=gcs_bucket_name)
                if signed_url:
                    timeseries_forecast_result['forecast_chart_url'] = signed_url
                    timeseries_forecast_result['feature_importance_chart_url'] = signed_url

    return jsonify({
        "preprocessing_summary": preprocessing_summary,
        "preprocessingLog": preprocessing_log, # Blobから読み込んだログを返す
        "insights": insights_data,
        "summary_report": summary_report,
        "infographic_urls": infographic_urls,
        "clusterAnalysisResult": cluster_analysis_result,
        "regressionAnalysisResult": regression_analysis_result,
        "timeseriesForecastResult": timeseries_forecast_result
    })


@data_analyzer_bp.route('/api/insight/<project_id>/<insight_id>', methods=['PATCH'])
def update_insight_api(project_id, insight_id):
    """
    【改訂版 v3.1】インサイトカードのメモや採用ステータスを更新する。
    Blob Storageへのオフロードに対応しつつ、後方互換性も維持する。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            return jsonify({"error": "プロジェクトが見つかりません。"}), 404
        
        # --- 1. Blob StorageまたはCosmos DBからインサイトデータを読み込む ---
        insights_data = None
        insights_blob_name = project_data.get("insightsBlobName")
        
        if insights_blob_name:
            # Blob Storageから読み込む (新しい形式)
            try:
                insights_bytes = download_blob_to_bytes(insights_blob_name)
                if insights_bytes:
                    insights_data = json.loads(insights_bytes.decode('utf-8'))
                else:
                    raise Exception("Insights blob data is empty or could not be downloaded.")
            except Exception as e:
                logger.error(f"Blob '{insights_blob_name}' からのインサイト読み込みに失敗: {e}", exc_info=True)
                return jsonify({"error": f"インサイトファイルの読み込みに失敗しました: {e}"}), 500
        elif "insights" in project_data:
            # 後方互換性: Cosmos DBのフィールドから直接読み込む (古い形式)
            insights_data = project_data.get('insights', {})

        if not insights_data:
             return jsonify({"error": "インサイトデータが見つかりません。"}), 404

        # --- 2. 更新対象のインサイトを検索する ---
        target_insight = None
        
        # 新しい階層化データ構造の場合 (dict)
        if isinstance(insights_data, dict):
            all_insight_lists = (
                [group.get("insights", []) for group in insights_data.get("primary_factor_insights", [])] +
                [group.get("insights", []) for group in insights_data.get("categorized_insights", [])] +
                [insights_data.get("anomaly_alerts", [])]
            )
            
            for insight_list in all_insight_lists:
                if not isinstance(insight_list, list): continue
                
                target_insight = next((i for i in insight_list if i and i.get('id') == insight_id), None)
                if target_insight:
                    break
        
        # 古いフラットなリスト構造の場合 (list)
        elif isinstance(insights_data, list):
            target_insight = next((i for i in insights_data if isinstance(i, dict) and i.get('id') == insight_id), None)
        
        if not target_insight:
            return jsonify({"error": "指定されたインサイトが見つかりません。"}), 404
        
        # --- 3. メモリ上でインサイトデータを更新 ---
        if 'is_adopted' in data:
            target_insight['is_adopted'] = data.get('is_adopted')
        if 'memo' in data:
            target_insight['memo'] = data.get('memo')
        
        # --- 4. 更新後のデータを適切な場所に書き戻す ---
        if insights_blob_name:
            # Blob Storageに上書き保存 (新しい形式)
            updated_bytes = json.dumps(insights_data, ensure_ascii=False).encode('utf-8')
            upload_bytes_to_blob(updated_bytes, insights_blob_name)
        else:
            # 後方互換性: Cosmos DBフィールドを更新 (古い形式)
            update_analysis_project(project_id, user_id, {"insights": insights_data})
        
        return jsonify({"success": True, "insight": target_insight})
        
    except Exception as e:
        logger.error(f"インサイト更新中にエラー (プロジェクトID: {project_id}, インサイトID: {insight_id}): {e}", exc_info=True)
        return jsonify({"error": "インサイトの更新中にサーバーで予期せぬエラーが発生しました。"}), 500

@data_analyzer_bp.route('/api/chat', methods=['POST'])
def analysis_chat_api():
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    project_id, user_message, context_insight, ab_test_prompt = data.get('project_id'), data.get('message'), data.get('context_insight'), data.get('ab_test_prompt')
    
    project_data = get_analysis_project(project_id, user_id)
    if not project_data: return jsonify({"error": "プロジェクトが見つかりません。"}), 404
    
    chat_history = project_data.get('analysisChatHistory', [])
    chat_history.append({"role": "user", "content": user_message})
    
    intent_prompt = f"ユーザーのメッセージ: \"{user_message}\"\n\nこのメッセージは「Pythonコードの生成・実行」を要求していますか、「自然言語での対話」を要求していますか？ `code_generation` または `natural_language` の一言で答えてください。"
    intent_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=500, contents=[types.Part.from_text(text=intent_prompt)],feature_category="Insight Canvas")
    intent = intent_response.strip()

    ai_response = {}
    if "code_generation" in intent:
        processed_blob_name = project_data.get("processedDataBlobName")

        if not processed_blob_name:
            return jsonify({"type": "text", "reply": "チャットでの分析を実行する前に、メインの分析を完了してください。前処理済みデータが見つかりません。"}), 200

        data_bytes = download_blob_to_bytes(processed_blob_name)
        
        try:
            df_processed = pd.read_csv(io.BytesIO(data_bytes))
        except Exception as e:
            return jsonify({"type": "text", "reply": f"分析用データの読み込みに失敗しました: {e}"}), 200

        if df_processed is None: return jsonify({"type": "text", "reply": "チャットでの分析中にデータファイルの読み込みに失敗しました。"}), 200
        
        df_info_buffer = io.StringIO()
        df_processed.info(buf=df_info_buffer)
        df_schema = f"""
### 利用可能なDataFrame `df_processed` のスキーマ情報:
{df_processed.head(3).to_markdown()}
{df_info_buffer.getvalue()}
"""

        # --- 1. Initial code generation ---
        code_gen_prompt = f"""あなたはPythonデータ分析専門家です。会話履歴と最新の指示に基づき、`df_processed`というDataFrameを操作する分析コードを生成してください。
{df_schema}
### 会話履歴 (参考): {json.dumps(chat_history[-5:-1], ensure_ascii=False)}
### 最新の指示: {user_message}
### 【最重要】厳格な出力ルール
1.  あなたの応答は、**` ```python ... ``` `で囲まれたPythonコードブロックのみ**でなければなりません。
2.  挨拶、解説、謝罪など、**コード以外のテキストは絶対に含めないでください。**
3.  `import`や`pd.read_csv`は不要です。常に`df_processed`を操作してください。
4.  **必ず上記のスキーマ情報に存在するカラム名のみを使用してください。**
### Pythonコードのみを出力してください:"""
        
        code_str_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=code_gen_prompt)],feature_category="Insight Canvas")
        
        code_to_run = None
        code_match = re.search(r'```python\n(.*?)```', code_str_response, re.DOTALL)
        if code_match:
            code_to_run = code_match.group(1).strip()
        else:
            ai_response = {"type": "text", "reply": f"申し訳ありません、コードの生成に失敗しました。AIの応答: {code_str_response}"}

        # --- 2. Execution and Retry Loop ---
        if code_to_run:
            final_exec_result = None
            max_retries = config.MAX_PYTHON_EXEC_RETRIES if hasattr(config, 'MAX_PYTHON_EXEC_RETRIES') else 2
            for attempt in range(max_retries):
                exec_result = execute_python_safely(code_to_run, session_id=project_id, global_vars={'df_processed': df_processed.copy()})
                error = exec_result.get('error')
                final_exec_result = exec_result # Store the latest result

                if not error:
                    logger.info(f"チャットコードの実行に成功 (試行 {attempt + 1})。")
                    break # Success, exit loop
                
                logger.warning(f"チャットコードの実行に失敗 (試行 {attempt + 1}/{max_retries})。エラー: {error}")
                if attempt < max_retries - 1:
                    fix_prompt = f"""あなたはPythonコードのデバッグの専門家です。以下のコードは実行時にエラーを発生させました。エラーメッセージと正しいデータスキーマを参考にコードを修正してください。
{df_schema}
### ユーザーの目的: {user_message}
### 失敗したコード:
```python
{code_to_run}
```
### エラーメッセージ:
```
{error}
```
### 【最重要】厳格な出力ルール
1.  あなたの応答は、修正後の完全なPythonコードを含む**` ```python ... ``` `コードブロックのみ**でなければなりません。
2.  解説や謝罪など、コード以外のテキストは**絶対に含めないでください。**
3.  **必ず上記のスキーマ情報に存在するカラム名のみを使用してください。**
### 修正後のPythonコード:"""
                    
                    fixed_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=fix_prompt)],feature_category="Insight Canvas")
                    fix_match = re.search(r'```python\n(.*?)```', fixed_code_response, re.DOTALL)
                    if fix_match:
                        code_to_run = fix_match.group(1).strip()
                        logger.info("AIが修正コードを生成しました。リトライします。")
                        continue
                    else:
                        logger.warning("AIは修正コードを生成できませんでした。リトライを中断します。")
                        break
            
            # --- 3. Process the final result (success or last failure) ---
            processed_results = []
            gcs_uris = []
            for result_item in final_exec_result.get('results', []):
                if result_item.get('type') == 'chart_gcs_uri':
                    gcs_uris.append(result_item['content'])
                    bucket_name = os.environ.get("GOOGLE_CLOUD_STORAGE_BUCKET_NAME", "")
                    blob_name = result_item['content'].replace(f'gs://{bucket_name}/', '')
                    sas_url = get_gcs_signed_url(blob_name, bucket_name=bucket_name)
                    if sas_url: processed_results.append({"type": "chart_sas_url", "content": sas_url})
                else:
                    processed_results.append(result_item)
            ai_response = {"type": "code_result", "results": processed_results, "error": final_exec_result.get('error')}

            if not final_exec_result.get('error'):
                summary_json = None
                try:
                    insight_gen_prompt = f"""あなたはデータ分析結果を要約する専門家です。以下の情報を元に、この分析の「タイトル」と「インサイト」を生成してください。
### ユーザーの指示: {user_message}
### 実行されたコード:\n```python\n{code_to_run}\n```
### 実行結果のテキスト:\n{json.dumps(final_exec_result.get('results', []), ensure_ascii=False, indent=2)}
### タスク: 上記を基に、以下のキーを持つJSONオブジェクトを生成してください。
- `title`: 分析内容を簡潔に表す日本語のタイトル。
- `insight`: ビジネスユーザー向けの簡潔な発見や結論（日本語、1〜2文）。
### 【最重要】出力形式: あなたの応答は、上記のキーを持つ単一のJSONオブジェクトの文字列**のみ**でなければなりません。解説やマークダウンの囲みは**絶対に含めないでください**。
"""
                    summary_response, _ = generate_with_gemini(
                        model_name=config.DATA_ANALYSIS_MODEL_NAME,
                        max_output_tokens=65535,
                        contents=[types.Part.from_text(text=insight_gen_prompt)], 
                        generation_config_override={"response_mime_type": "application/json"},
                        feature_category="Insight Canvas"
                    )
                    
                    json_match = re.search(r'\{.*\}', summary_response, re.DOTALL)
                    if json_match:
                        summary_json = json.loads(json_match.group(0))
                    else:
                        raise ValueError("AIの応答に有効なJSONオブジェクトが含まれていません。")

                except (json.JSONDecodeError, ValueError) as e:
                    logger.warning(f"チャットからのインサイト要約生成（JSON解析）に失敗: {e}. AI Raw Response: {summary_response[:500]}")
                    summary_json = {
                        "title": "チャットからの追加分析",
                        "insight": "チャットでの対話を通じて追加された分析です（AIによる自動要約に失敗しました）。"
                    }
                
                new_insight = {
                    'id': f"insight_{uuid.uuid4().hex}",
                    'title': summary_json.get('title', 'チャットからの追加分析'),
                    'description': summary_json.get('insight', 'チャットでの対話を通じて追加された分析です。'),
                    'stdout': "".join(item['content'] for item in final_exec_result.get('results', []) if item['type'] == 'stdout'),
                    'error': None, 'charts': gcs_uris, 'memo': ''
                }

                # 常に階層化された辞書を維持するように修正
                insights_data = project_data.get('insights', {})

                # insights_dataが辞書でない（古いリスト形式など）場合は、新しい構造に初期化
                if not isinstance(insights_data, dict):
                    insights_data = {
                        "primary_factor_insights": [],
                        "other_insights": [],
                        "anomaly_alerts": []
                    }
                
                # project_dataのinsightsを更新# 常に階層化された辞書を維持するように修正
                insights_data = project_data.get('insights', {})
                if not isinstance(insights_data, dict):
                    insights_data = {"categorized_insights": []} # 構造を初期化

                # 「categorized_insights」リストが存在しない場合は作成
                if 'categorized_insights' not in insights_data:
                    insights_data['categorized_insights'] = []

                # 「その他のインサイト」カテゴリを検索、なければ作成
                other_insights_category = next((cat for cat in insights_data['categorized_insights'] if cat.get('category') == 'その他のインサイト'), None)

                if not other_insights_category:
                    other_insights_category = {"category": "その他のインサイト", "insights": []}
                    insights_data['categorized_insights'].append(other_insights_category)

                # 新しいインサイトを「その他のインサイト」カテゴリに追加
                other_insights_category['insights'].append(new_insight)
                
                update_analysis_project(project_id, user_id, {"insights": insights_data})
                
                gcs_bucket_name = os.environ.get("GOOGLE_CLOUD_STORAGE_BUCKET_NAME")
                if gcs_bucket_name:
                    chart_urls = [get_gcs_signed_url(uri.replace(f'gs://{gcs_bucket_name}/', ''), bucket_name=gcs_bucket_name) for uri in new_insight.get('charts', [])]
                    new_insight['chart_urls'] = [url for url in chart_urls if url]
                
                ai_response['new_insight'] = new_insight

    else:
        chat_prompt = f"""あなたはデータ分析の専門家です。以下の文脈を元に、ユーザーの質問に回答してください。
- **プロジェクト全体の目的**: {project_data.get('analysisPurpose','')}
- **現在注目しているインサイト**: {context_insight or '(なし)'}
- **これまでの会話履歴**: {json.dumps(chat_history[:-1], ensure_ascii=False)}
- **ユーザーの最新の質問**: {user_message}"""
        
        if ab_test_prompt:
            chat_prompt += f"""

### 【特別指示】A/Bテストの設計
上記の文脈、特に注目しているインサイトに基づき、**{ab_test_prompt}**
この効果をビジネス成果に繋げるための具体的なA/Bテストを設計してください。
応答は、必ず以下の5項目を含む構成で記述してください。
1.  **仮説 (Hypothesis)**: このテストで何を検証しようとしているのか。
2.  **対象 (Target Audience)**: 誰に対してテストを行うのか。
3.  **テスト内容 (Test Details)**: 具体的に何を、どのように変更するのか（A案 vs B案）。
4.  **評価指標 (Key Metrics)**: テストの成否をどの数値で判断するのか（主要指標と副次指標）。
5.  **推奨期間 (Duration)**: どのくらいの期間テストを実施すべきか。
"""
        reply_text, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=chat_prompt)],feature_category="Insight Canvas")
        ai_response = {"type": "text", "reply": reply_text}
    
    chat_history.append({"role": "assistant", "content": ai_response})
    update_analysis_project(project_id, user_id, {"analysisChatHistory": chat_history})
    return jsonify(ai_response)

@data_analyzer_bp.route('/api/profile/<project_id>/apply_preprocessing', methods=['PATCH'])
def apply_preprocessing_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    code_to_run = data.get('code')

    if not code_to_run:
        return jsonify({"error": "実行するコードがありません。"}), 400

    try:
        # 1. プロジェクト情報を取得してデータファイルの場所を特定
        project_data = get_analysis_project(project_id, user_id)
        if not project_data or 'dataFile' not in project_data:
            return jsonify({"error": "プロジェクトまたはデータファイルが見つかりません。"}), 404
        
        data_file_info = project_data['dataFile']
        blob_name = data_file_info['blob_name']
        
        # 2. 現在のデータをダウンロードしてDataFrameに読み込み
        data_bytes = download_blob_to_bytes(blob_name)
        df = _load_df_from_bytes_with_temp_file(data_bytes, data_file_info['filename'])
        if df is None:
            raise Exception("データファイルの読み込みに失敗しました。")

        # 3. 受け取ったコードを安全なサンドボックスで実行
        # 修正: AIが生成するコードは 'df' を操作するため、変数名を 'df' に統一
        exec_result = execute_python_safely(code_to_run, session_id=project_id, global_vars={'df': df.copy()})

        if exec_result.get('error'):
            raise Exception(f"コードの実行中にエラーが発生しました: {exec_result.get('error')}")

        # 4. 変更後のDataFrameを取得
        df_processed = exec_result.get('final_scope', {}).get('df')
        if df_processed is None:
            raise Exception("コード実行後、処理済みのDataFrame 'df' が見つかりませんでした。")
        
        # 5. 変更後のDataFrameをCSV形式のバイトデータに変換
        output = io.StringIO()
        df_processed.to_csv(output, index=False)
        updated_data_bytes = output.getvalue().encode('utf-8')

        # 6. 同じBLOB名でファイルを上書きアップロード
        upload_bytes_to_blob(updated_data_bytes, blob_name)

        logger.info(f"プロジェクト {project_id} のデータファイルが前処理コードによって更新されました。")
        return jsonify({"success": True, "message": "データの前処理が適用されました。"})

    except Exception as e:
        logger.error(f"前処理の適用中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": f"前処理の適用に失敗しました: {str(e)}"}), 500


@data_analyzer_bp.route('/api/profile/<project_id>/download', methods=['GET'])
def download_data_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data or 'dataFile' not in project_data:
            return "プロジェクトまたはデータファイルが見つかりません。", 404
            
        data_file_info = project_data['dataFile']
        blob_name = data_file_info['blob_name']
        original_filename = data_file_info['filename']

        # BLOBからデータをバイトとしてダウンロード
        data_bytes = download_blob_to_bytes(blob_name)
        if data_bytes is None:
            return "ファイルのダウンロードに失敗しました。", 404

        # バイトデータを send_file を使ってレスポンスとして返す
        return send_file(
            io.BytesIO(data_bytes),
            mimetype='text/csv', # or application/vnd.ms-excel etc.
            as_attachment=True,
            download_name=original_filename
        )

    except Exception as e:
        logger.error(f"データダウンロード中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return "サーバーエラーが発生しました。", 500
    
@data_analyzer_bp.route('/api/profile/<project_id>/rerun', methods=['POST'])
def rerun_profiling_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    try:
        # 新しいタスクIDを生成
        task_id = f"task_{project_id}_profiling"
        # 既存のタスクがあれば更新、なければ作成
        update_task_in_cosmos(task_id, {"status": "PENDING", "task_type": "data_profiling", "project_id": project_id, "user_id": user_id})
        
        # プロファイリングタスクを非同期で実行
        executor.submit(_run_profiling_task, task_id, project_id, user_id)
        
        return jsonify({"success": True, "message": "プロファイリングの再実行を開始しました。", "task_id": task_id}), 202

    except Exception as e:
        logger.error(f"プロファイリング再実行の起動中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "再実行の開始に失敗しました。"}), 500
    
def _generate_code_from_rules(rules: dict) -> str:
    """保存されたルール(JSON)を、実行可能なPython(pandas)コードに変換する。"""
    code_lines = []
    
    # 1. 欠損値処理
    missing_rules = rules.get('missing_values', {})
    for col, method in missing_rules.items():
        code_lines.append(f"# [{col}]列の欠損値処理: {method}")
        if method == 'drop':
            code_lines.append(f"df.dropna(subset=['{col}'], inplace=True)")
        elif method == 'mean':
            code_lines.append(f"if pd.api.types.is_numeric_dtype(df['{col}']): df['{col}'].fillna(df['{col}'].mean(), inplace=True)")
        elif method == 'median':
            code_lines.append(f"if pd.api.types.is_numeric_dtype(df['{col}']): df['{col}'].fillna(df['{col}'].median(), inplace=True)")
        elif method == 'mode':
            code_lines.append(f"if not df['{col}'].mode().empty: df['{col}'].fillna(df['{col}'].mode().iloc[0], inplace=True)")
        elif method == 'zero':
            code_lines.append(f"if pd.api.types.is_numeric_dtype(df['{col}']): df['{col}'].fillna(0, inplace=True)")
        elif method == 'unknown':
            code_lines.append(f"df['{col}'].fillna('不明', inplace=True)")
    
    # 2. 外れ値処理
    outlier_rules = rules.get('outliers', {})
    for col, config in outlier_rules.items():
        method = config.get('method')
        threshold = config.get('threshold', 1.0) / 100.0
        code_lines.append(f"# [{col}]列の外れ値処理: {method} (上位/下位 {threshold*100}%)")
        if method == 'clip':
            code_lines.append(f"if pd.api.types.is_numeric_dtype(df['{col}']):")
            code_lines.append(f"    lower_bound = df['{col}'].quantile({threshold})")
            code_lines.append(f"    upper_bound = df['{col}'].quantile(1 - {threshold})")
            code_lines.append(f"    df['{col}'] = df['{col}'].clip(lower=lower_bound, upper=upper_bound)")
        elif method == 'remove':
            code_lines.append(f"if pd.api.types.is_numeric_dtype(df['{col}']):")
            code_lines.append(f"    lower_bound = df['{col}'].quantile({threshold})")
            code_lines.append(f"    upper_bound = df['{col}'].quantile(1 - {threshold})")
            code_lines.append(f"    df = df[(df['{col}'] >= lower_bound) & (df['{col}'] <= upper_bound)]")
    
    # 3. 低頻度カテゴリ統合
    low_freq_rules = rules.get('low_frequency', {})
    for col, threshold_pct in low_freq_rules.items():
        threshold = threshold_pct / 100.0
        code_lines.append(f"# [{col}]列の低頻度カテゴリ統合 (しきい値: {threshold_pct}%)")
        code_lines.append(f"if not pd.api.types.is_numeric_dtype(df['{col}']):")
        code_lines.append(f"    value_counts = df['{col}'].value_counts(normalize=True)")
        code_lines.append(f"    to_replace = value_counts[value_counts < {threshold}].index")
        code_lines.append(f"    df['{col}'] = df['{col}'].replace(to_replace, 'その他')")
        
    return "\n".join(code_lines)

@data_analyzer_bp.route('/api/profile/<project_id>/download_processed', methods=['GET'])
def download_processed_data_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data or 'dataFile' not in project_data:
            return jsonify({"error": "プロジェクトまたはデータファイルが見つかりません。"}), 404
        
        rules = project_data.get("preprocessingRules", {})
        code_to_run = _generate_code_from_rules(rules)
        
        data_file_info = project_data['dataFile']
        data_bytes = download_blob_to_bytes(data_file_info['blob_name'])
        df = _load_df_from_bytes_with_temp_file(data_bytes, data_file_info['filename'])
        if df is None:
            raise Exception("データファイルの読み込みに失敗しました。")

        if code_to_run:
            exec_result = execute_python_safely(code_to_run, session_id=f"{project_id}_download", global_vars={'df': df.copy()})
            if exec_result.get('error'):
                raise Exception(f"前処理コードの実行中にエラーが発生しました: {exec_result.get('error')}")
            
            df_processed = exec_result.get('final_scope', {}).get('df')
            if df_processed is None:
                df_processed = df # エラーにはせず、元のdfを返す
        else:
            df_processed = df

        output = io.BytesIO()
        df_processed.to_csv(output, index=False, encoding='utf-8-sig')
        output.seek(0)
        
        original_filename = data_file_info['filename']
        base, _ = os.path.splitext(original_filename)
        download_name = f"{base}_processed.csv"

        return send_file(
            output,
            mimetype='text/csv',
            as_attachment=True,
            download_name=download_name
        )

    except Exception as e:
        logger.error(f"加工済みデータのダウンロード中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": f"ダウンロードに失敗しました: {str(e)}"}), 500
    
@data_analyzer_bp.route('/api/analysis/<project_id>/download_processed', methods=['GET'])
def download_analysis_processed_data_api(project_id):
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        project_data = get_analysis_project(project_id, user_id)
        processed_blob_name = project_data.get("processedDataBlobName")

        if not processed_blob_name:
            return jsonify({"error": "分析後のデータが見つかりません。分析が完了しているか確認してください。"}), 404
            
        data_bytes = download_blob_to_bytes(processed_blob_name)
        if data_bytes is None:
            return jsonify({"error": "ファイルのダウンロードに失敗しました。"}), 404

        original_filename = project_data.get('dataFile', {}).get('filename', 'data.xlsx')
        base, _ = os.path.splitext(original_filename)
        if not base:
            base = project_id
        download_name = f"{base}_analyzed.csv"

        return send_file(
            io.BytesIO(data_bytes),
            mimetype='text/csv',
            as_attachment=True,
            download_name=download_name
        )

    except Exception as e:
        logger.error(f"分析後データのダウンロード中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": f"サーバーエラーが発生しました: {str(e)}"}), 500

def _run_cluster_analysis(project_id: str, user_id: str, df: pd.DataFrame, context_summary: str) -> dict | None:
    """
    【改訂版 v4】データに対してクラスター分析を実行し、AIによる命名と可視化を行う。
    'cluster_name'列の追加漏れを修正し、関数内のデータフレームの一貫性を担保。
    """
    # --- 1. 関数のセットアップとインポート ---
    from sklearn.decomposition import PCA
    from sklearn.preprocessing import StandardScaler
    from google.genai import types
    import pandas as pd
    import numpy as np
    import json
    import re
    import config

    logger.info(f"プロジェクト {project_id} のクラスター分析を開始します。")
    
    MAX_RETRIES = config.MAX_PYTHON_EXEC_RETRIES if hasattr(config, 'MAX_PYTHON_EXEC_RETRIES') else 2

    try:
        # --- 2. 日付変換とデータ準備 ---
        df_local = df.copy() # この関数内で使用するローカルコピーを作成
        for col in df_local.select_dtypes(include=['object']).columns:
            try:
                # 文字列型の列を日付型に変換しようと試みる
                converted_series = pd.to_datetime(df_local[col], errors='coerce')
                # 元々NULLでなかった値の半分以上が正常に日付に変換できた場合、その列をdatetime型とみなす
                original_not_null = df_local[col].notna().sum()
                if original_not_null > 0 and converted_series.notna().sum() / original_not_null > 0.5:
                    df_local[col] = converted_series
                    logger.info(f"[{project_id}] クラスタ分析のため、日付/時刻列 '{col}' を再変換しました。")
            except (ValueError, TypeError):
                # 変換に失敗した列はそのまま無視
                continue

        logger.info(f"[{project_id}] クラスター分析: データ前処理を開始。")
        # 再変換した df_local を使って列を定義する
        date_cols = df_local.select_dtypes(include=['datetime64[ns]', 'datetime64[ns, UTC]']).columns.tolist()
        numeric_cols = df_local.select_dtypes(include=np.number).columns.tolist()
        categorical_cols = df_local.select_dtypes(include=['object', 'category']).columns.tolist()

        if not numeric_cols:
            logger.warning(f"プロジェクト {project_id} には分析可能な数値データがありません。")
            return {"error": "クラスター分析には少なくとも1つ以上の数値データが必要です。", "skipped": True}

        df_for_clustering = df_local.copy()

        # 1. 数値データの標準化
        # StandardScalerは分散0の列などでNaNを生む可能性があるため、この段階ではNaNが残る可能性があります
        df_for_clustering[numeric_cols] = StandardScaler().fit_transform(df_for_clustering[numeric_cols])
        
        # 2. 強制的な欠損値処理 (Imputation)
        # KMeansはNaNを許容しないため、標準化後に残ったNaNやInfを確実に0（平均的な値）で埋めます
        if numeric_cols:
            imputer = SimpleImputer(strategy='constant', fill_value=0)
            df_for_clustering[numeric_cols] = imputer.fit_transform(df_for_clustering[numeric_cols])

        # 3. カテゴリデータのダミー変数化
        if categorical_cols:
            df_for_clustering = pd.get_dummies(df_for_clustering, columns=categorical_cols, dummy_na=False)
            
            # ダミー変数化で生成された列も含めて、再度NaNチェックと補完を行う（念のため）
            # 新しく生成された列（bool型やint型）をfloatに変換して欠損処理を統一
            df_for_clustering = df_for_clustering.astype(float)
            df_for_clustering.fillna(0, inplace=True)

        # 分析に使用する最終的な列リストを定義
        final_cols_for_clustering = [col for col in df_for_clustering.columns if col in numeric_cols or any(cat_col in col for cat_col in categorical_cols) or col in df_for_clustering.columns] 
        # (注: get_dummiesで列名が変わるため、columns全体から数値型になったものを対象にするのが安全です)
        final_cols_for_clustering = df_for_clustering.select_dtypes(include=np.number).columns.tolist()

        if not final_cols_for_clustering:
             raise Exception("前処理の結果、分析に使用できる列がなくなりました。")
    
        full_code_log = "# --- クラスター分析 実行コードログ ---\n"

        # --- 3. エルボー法による最適クラスター数推定（リトライ機能付き）---
        logger.info(f"[{project_id}] クラスター分析: エルボー法のコード生成・実行を開始。")
        elbow_code_prompt = f"""あなたはデータサイエンティストです。
メモリ上には既に `df_for_clustering` という前処理済みDataFrame（行数: {len(df_for_clustering)}）が存在します。
これを使ってKMeansクラスタリングを実行し、エルボー法で最適クラスター数を推定するPythonコードを生成してください。

### タスク
1.  `sklearn.cluster.KMeans` を使い、クラスター数を1から10まで変化させたときのSSE（イナーシャ）を計算してください。
2.  結果を折れ線グラフで描画し、タイトルに「エルボー法による最適クラスター数の推定」と付けてください。

### 【最重要】禁止事項とルール
1.  **`df_for_clustering` を再定義（`pd.DataFrame(...)`）、上書き、読み込み（`read_csv`）は絶対に行わないでください。** 必ずメモリ上の既存変数を使用してください。
2.  **ダミーデータの作成は厳禁です。**「例としてダミーデータを作成します」といったコードは不要です。
3.  **欠損値処理（fillna, SimpleImputer）は不要です。** データは既に前処理済みです。
4.  出力は、` ```python ... ``` `で囲まれたPythonコードブロックのみとします。
5.  `plt.show()`は呼び出さないでください。
"""
        elbow_code = ""
        elbow_exec_result = None
        for attempt in range(MAX_RETRIES):
            if not elbow_code:
                # コード生成
                elbow_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=elbow_code_prompt)],feature_category="Insight Canvas")
                code_match = re.search(r'```(?:python\n)?(.*?)\n?```', elbow_code_response, re.DOTALL)
                if not code_match:
                    logger.warning(f"[{project_id}] エルボー法コード生成に失敗 (試行 {attempt + 1})。AI応答形式が不正です。")
                    elbow_code_prompt += "\n\n前回の応答は形式が不適切でした。必ずPythonコードのみを` ```python ... ``` `で囲んで出力してください。"
                    continue
                elbow_code = code_match.group(1).strip()
            
            # コード実行
            elbow_exec_result = execute_python_safely(elbow_code, session_id=project_id, global_vars={'df_for_clustering': df_for_clustering[final_cols_for_clustering]})
            
            if not elbow_exec_result.get('error'):
                logger.info(f"[{project_id}] エルボー法コードの実行に成功。")
                break # 成功
            
            # 失敗時の処理
            error = elbow_exec_result.get('error')
            logger.warning(f"[{project_id}] エルボー法コードの実行に失敗 (試行 {attempt + 1})。エラー: {error}")
            
            if attempt < MAX_RETRIES - 1:
                fix_prompt = f"""Pythonコードのデバッグ専門家として、以下のコードのエラーを修正してください。
### 失敗したコード:
```python
{elbow_code}
```
### エラーメッセージ:```
{error}
```
### 【最重要】修正後の完全なPythonコードのみを` ```python ... ``` `で囲んで出力してください。"""
                fix_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=fix_prompt)],feature_category="Insight Canvas")
                fix_match = re.search(r'```(?:python\n)?(.*?)\n?```', fix_response, re.DOTALL)
                if fix_match:
                    elbow_code = fix_match.group(1).strip()
                    logger.info(f"[{project_id}] AIがエルボー法コードの修正案を生成。リトライします。")
                else:
                    logger.warning(f"[{project_id}] AIが修正コードを生成できませんでした。リトライを中断します。")
                    break
        
        if elbow_exec_result is None:
            raise Exception("AIによるエルボー法コードの生成が複数回失敗しました。AIの応答形式が不正です。")
        if elbow_exec_result.get('error'): 
            raise Exception(f"エルボー法コードの実行に失敗しました: {elbow_exec_result['error']}")
        
        elbow_chart_uri = next((item['content'] for item in elbow_exec_result.get('results', []) if item['type'] == 'chart_gcs_uri'), None)

        if 'elbow_code' in locals():
            full_code_log += f"\n# [Step 1] エルボー法による最適クラスター数推定\n{elbow_code}\n"

        # --- 4. AIによる最適クラスター数の決定 ---
        logger.info(f"[{project_id}] クラスター分析: 最適クラスター数をAIが決定中。")
        k_decision_prompt = f"""あなたはデータサイエンティストです。一般的に、エルボー法ではグラフの「肘」のように見える点が最適なクラスター数とされます。データセットのカラム数（{len(df.columns)}）やデータの性質を考慮し、2から6の範囲で最も可能性の高い最適クラスター数を1つだけ、**数字のみで**回答してください。"""
        optimal_k_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=k_decision_prompt)],feature_category="Insight Canvas")
        k_match = re.search(r'\d+', optimal_k_response)
        optimal_k = int(k_match.group(0)) if k_match else 3
        if not k_match: logger.warning(f"[{project_id}] AIから最適なkを取得できませんでした。デフォルト値k={optimal_k}を使用します。")
        

        # --- 5. KMeansクラスタリングの実行（リトライ機能付き）---
        logger.info(f"[{project_id}] クラスター分析: k={optimal_k}でのクラスタリングコード生成・実行を開始。")
        clustering_code_prompt = f"""あなたはデータサイエンティストです。
メモリ上には既に `df_for_clustering` という前処理済みDataFrame（行数: {len(df_for_clustering)}）が存在します。
これを使ってKMeansクラスタリングを実行してください。

### タスク
1.  `sklearn.cluster.KMeans` を使い、既存の `df_for_clustering` に対して `n_clusters={optimal_k}` でクラスタリングを実行してください。
2.  結果のラベル配列（`labels_`）を、`cluster_labels` という変数に代入してください。

### 【最重要】禁止事項とルール
1.  **`df_for_clustering` を再定義（`pd.DataFrame(...)`）、上書き、読み込みは絶対に行わないでください。** メモリ上の変数をそのまま使用してください。
2.  **行数の削減は厳禁です。** `df_for_clustering.head()` やスライス `[:5]` などは絶対に使用しないでください。入力データの行数と、出力される `cluster_labels` の要素数は完全に一致する必要があります。
3.  欠損値処理は完了しているため、`SimpleImputer` は不要です。
4.  出力は、` ```python ... ``` `で囲まれたPythonコードブロックのみとします。
"""
        clustering_code = ""
        clustering_exec_result = None
        for attempt in range(MAX_RETRIES):
            if not clustering_code:
                # コード生成
                clustering_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=clustering_code_prompt)],feature_category="Insight Canvas")
                code_match = re.search(r'```(?:python\n)?(.*?)\n?```', clustering_code_response, re.DOTALL)
                if not code_match:
                    logger.warning(f"[{project_id}] クラスタリングコード生成に失敗 (試行 {attempt + 1})。")
                    continue
                clustering_code = code_match.group(1).strip()
            
            # コード実行
            clustering_exec_result = execute_python_safely(clustering_code, session_id=project_id, global_vars={'df_for_clustering': df_for_clustering[final_cols_for_clustering], 'df_original': df_local.copy()})
            if not clustering_exec_result.get('error'):
                logger.info(f"[{project_id}] クラスタリングコードの実行に成功。")
                break
            
            # 失敗時の処理
            error = clustering_exec_result.get('error')
            logger.warning(f"[{project_id}] クラスタリングコード実行に失敗 (試行 {attempt + 1})。エラー: {error}")
            if attempt < MAX_RETRIES - 1:
                fix_prompt = f"""Pythonコードのデバッグ専門家として、以下のコードのエラーを修正してください。
### 失敗したコード:
```python
{clustering_code}
```
### エラーメッセージ:
```
{error}
```
### 修正のルール
- `ValueError: Input X contains NaN` の場合でも、データ削減（`head()`など）は行わず、`sklearn.impute.SimpleImputer` 等で値を補完してください。
- **【最重要】`df_for_clustering` の行数を減らすことは絶対に禁止です。** 結果のラベル配列の長さが変わるとシステムエラーになります。

### 【最重要】修正後の完全なPythonコードのみを` ```python ... ``` `で囲んで出力してください。"""
                fix_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=fix_prompt)],feature_category="Insight Canvas")
                fix_match = re.search(r'```(?:python\n)?(.*?)\n?```', fix_response, re.DOTALL)
                if fix_match:
                    clustering_code = fix_match.group(1).strip()
                    logger.info(f"[{project_id}] AIがクラスタリングコードの修正案を生成。リトライします。")
                else:
                    break
        
        if clustering_exec_result is None:
            raise Exception("AIによるクラスタリングコードの生成が複数回失敗しました。")
        if clustering_exec_result.get('error'): 
            raise Exception(f"クラスタリングコードの実行に失敗しました: {clustering_exec_result['error']}")

        cluster_labels = clustering_exec_result.get('final_scope', {}).get('cluster_labels')
        if cluster_labels is None:
            raise Exception("`cluster_labels` がAIのコード内で定義されませんでした。")

        df_final_clustered = df_local.copy()
        
        if len(df_final_clustered) != len(cluster_labels):
            raise Exception(f"元データとクラスタリング結果の行数が一致しません。元データ: {len(df_final_clustered)}行, 結果: {len(cluster_labels)}ラベル")

        df_final_clustered['cluster'] = cluster_labels

        if 'clustering_code' in locals():
            full_code_log += f"\n# [Step 2] KMeansクラスタリング (k={optimal_k})\n{clustering_code}\n"

        # --- 6. AIによる各クラスターの命名 ---
        logger.info(f"[{project_id}] クラスター分析: 各クラスターの命名をAIが実行中。")
        cluster_summaries = {
            f"クラスター{i}": {
                "numeric_stats": df_final_clustered[df_final_clustered['cluster'] == i][numeric_cols].describe().transpose().to_dict('index'),
                "count": len(df_final_clustered[df_final_clustered['cluster'] == i])
            }
            for i in sorted(df_final_clustered['cluster'].unique())
        }
        naming_prompt = f"""あなたは優れたマーケティングアナリストです。以下の各クラスターの統計的特徴を分析し、それぞれのクラスターに最も的確な「分類名」と「短い解説文」を日本語で付けてください。
### データコンテキスト
{context_summary}
### 各クラスターの統計情報```json
{json.dumps(cluster_summaries, indent=2, ensure_ascii=False)}
```
### 出力形式 (JSONのみ)
```json
{{
  "clusters": [
    {{ "cluster_id": 0, "name": "（例）ロイヤル顧客層", "description": "（例）購入頻度と単価が最も高く..." }},
    {{ "cluster_id": 1, "name": "（例）新規お試し層", "description": "（例）最近利用を開始し..." }}
  ]
}}```"""
        naming_response_str, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=naming_prompt)], generation_config_override={"response_mime_type": "application/json"},feature_category="Insight Canvas")
        cluster_names_info = json.loads(naming_response_str)
        cluster_name_map = {c['cluster_id']: c['name'] for c in cluster_names_info['clusters']}
        
        df_final_clustered['cluster_name'] = df_final_clustered['cluster'].map(cluster_name_map)
        
        
        # --- 7. 条件に応じた可視化 (日付列の有無を最優先) ---
        logger.info(f"[{project_id}] クラスター分析: 可視化コード生成・実行を開始。")
        
        is_timeseries = bool(date_cols)
        cluster_chart_uri = None
        df_for_plot = None
        visualization_prompt = ""        

        # 7-A. 日付列が存在する場合: 時系列プロット (特徴量が1つでもこちらを優先)
        if is_timeseries and date_cols and numeric_cols:
            logger.info(f"[{project_id}] 時系列データとして検出。時系列プロットによる可視化を実行します。")
            
            value_col_name = numeric_cols[0]
            date_col_name = date_cols[0]
            
            visualization_prompt = f"""あなたはPythonのデータ可視化専門家です。
メモリ上には既に `df_plot` というDataFrame（列: `{date_col_name}`, `{value_col_name}`, `cluster_name`）が存在します。
これを使って可視化コードを生成してください。

### タスク
1.  `seaborn`と`matplotlib`を使い、時系列データをクラスタごとに色分けして可視化してください。
2.  背景に全体推移（グレー）、手前にクラスタ別散布図を描画してください。

### 【最重要】禁止事項とルール
1.  **`df_plot` を再定義（`data = {{...}}`など）、上書き、読み込みは絶対に行わないでください。** 必ずメモリ上の既存変数を使用してください。
2.  **ダミーデータの作成は厳禁です。コメントアウトした状態であっても、ダミーデータの例は一切記述しないでください。**
3.  コードの先頭で `import japanize_matplotlib; japanize_matplotlib.japanize()` を呼び出してください。
4.  出力は、` ```python ... ``` `で囲まれたPythonコードブロックのみとします。
"""

            if 'full_code_log' in locals():
                full_code_log += f"""
# [Step 3-Prep] 可視化用データの準備 (時系列)
# クラスタリング結果を含んだデータフレームから必要な列を抽出
df_plot = df_final_clustered[['{date_col_name}', '{value_col_name}', 'cluster_name']].copy()
# df_plot が可視化の入力となります
"""

            df_for_plot = df_final_clustered[[date_col_name, value_col_name, 'cluster_name']].copy()

        # 7-B. 日付列が存在しない、かつ特徴量が2つ以上の場合: PCAプロット
        elif not is_timeseries and df_for_clustering[final_cols_for_clustering].shape[1] >= 2:
            logger.info(f"[{project_id}] 非時系列データとして検出。PCAによる可視化を実行します。")
            pca = PCA(n_components=2)
            principal_components = pca.fit_transform(df_for_clustering[final_cols_for_clustering])
            pca_df = pd.DataFrame(data=principal_components, columns=['PC1', 'PC2'])
            pca_df.index = df_final_clustered.index

            if 'full_code_log' in locals():
                full_code_log += f"""
# [Step 3-Prep] 可視化用データの準備 (次元圧縮)
from sklearn.decomposition import PCA

# df_for_clustering は Step 1, 2 で使用した前処理済みデータ
pca = PCA(n_components=2)
principal_components = pca.fit_transform(df_for_clustering)

df_plot = pd.DataFrame(data=principal_components, columns=['PC1', 'PC2'])
df_plot['cluster_name'] = df_final_clustered['cluster_name']
# df_plot が可視化の入力となります
"""

            df_for_plot = pd.concat([pca_df, df_final_clustered['cluster_name']], axis=1)
            explained_variance_ratio = pca.explained_variance_ratio_

            visualization_prompt = f"""あなたはPythonのデータ可視化専門家です。
メモリ上には既に `df_plot` というDataFrame（列: `PC1`, `PC2`, `cluster_name`）が存在します。
これを使って散布図を作成してください。

### タスク
1.  `seaborn`を使い、`cluster_name`で色分けした散布図を描画してください。
2.  タイトルは「PCAによるクラスター分布」としてください。

### 【最重要】禁止事項とルール
1.  **`df_plot` を再定義（`data = {{...}}`など）、上書き、読み込みは絶対に行わないでください。** 必ずメモリ上の既存変数を使用してください。
2.  **ダミーデータの作成は厳禁です。コメントアウトした状態であっても、ダミーデータの例は一切記述しないでください。**
3.  コードの先頭で `import japanize_matplotlib; japanize_matplotlib.japanize()` を呼び出してください。
4.  出力は、` ```python ... ``` `で囲まれたPythonコードブロックのみとします。
"""

        # 7-C. 上記のいずれにも当てはまらない場合 (可視化スキップ)
        else:
            logger.warning(f"[{project_id}] クラスター分析の可視化をスキップします。理由: 可視化に適した列の組み合わせ（日付列 or 2つ以上の特徴量）が見つかりませんでした。")  

        # 7-D. 可視化コードの実行 (共通ロジック)
        if df_for_plot is not None and visualization_prompt:
            viz_code = ""
            viz_exec_result = None
            for attempt in range(MAX_RETRIES):
                if not viz_code:
                    viz_code_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=visualization_prompt)],feature_category="Insight Canvas")
                    code_match = re.search(r'```(?:python\n)?(.*?)\n?```', viz_code_response, re.DOTALL)
                    if not code_match:
                        logger.warning(f"[{project_id}] 可視化コード生成に失敗 (試行 {attempt + 1})。")
                        continue
                    viz_code = code_match.group(1).strip()
                
                viz_exec_result = execute_python_safely(viz_code, session_id=project_id, global_vars={'df_plot': df_for_plot.copy()})
                if not viz_exec_result.get('error'):
                    logger.info(f"[{project_id}] 可視化コードの実行に成功。")
                    break
                
                error = viz_exec_result.get('error')
                logger.warning(f"[{project_id}] 可視化コード実行に失敗 (試行 {attempt + 1})。エラー: {error}")
                if attempt < MAX_RETRIES - 1:
                    fix_prompt = f"""Pythonコードのデバッグ専門家として、以下のコードのエラーを修正してください。
### 失敗したコード:
```python
{viz_code}```
### エラーメッセージ:
```
{error}```
### 【最重要】修正後の完全なPythonコードのみを` ```python ... ``` `で囲んで出力してください。"""
                    fix_response, _ = generate_with_gemini(model_name=config.DATA_ANALYSIS_MODEL_NAME, max_output_tokens=65535, contents=[types.Part.from_text(text=fix_prompt)],feature_category="Insight Canvas")
                    fix_match = re.search(r'```(?:python\n)?(.*?)\n?```', fix_response, re.DOTALL)
                    if fix_match:
                        viz_code = fix_match.group(1).strip()
                        logger.info(f"[{project_id}] AIが可視化コードの修正案を生成。リトライします。")
                    else:
                        break
            
            if viz_exec_result and not viz_exec_result.get('error'):
                cluster_chart_uri = next((item['content'] for item in viz_exec_result.get('results', []) if item['type'] == 'chart_gcs_uri'), None)

        if 'viz_code' in locals():
            full_code_log += f"\n# [Step 3] クラスタリング結果の可視化\n{viz_code}\n"

        # --- 8. 最終結果の返却 ---
        logger.info(f"[{project_id}] クラスター分析の全工程が正常に完了。")
        return {
            "summary": cluster_names_info,
            "elbow_chart_url": elbow_chart_uri,
            "cluster_chart_url": cluster_chart_uri,
            "skipped": False,
            "error": None,
            "code": full_code_log
        }, df_final_clustered

    except Exception as e:
        logger.error(f"クラスター分析タスク {project_id} でエラー: {e}", exc_info=True)
        return {"error": str(e), "skipped": False}, None


def _run_anomaly_detection(df: pd.DataFrame, project_id: str, column_roles: dict, analysis_purpose: str) -> dict | None:
    """
    【改善版】DataFrame内の時系列データを分析し、異常値を検出する。
    検出された場合、結果を可視化したグラフを含むインサイト形式の辞書で返す。
    """
    logger.info(f"プロジェクト {project_id} の異常検知を開始。")
    
    datetime_cols = df.select_dtypes(include=['datetime64[ns]', 'datetime64[ns, UTC]']).columns.tolist()
    if not datetime_cols:
        logger.info("分析可能な日付/時刻型の列がないため、異常検知をスキップします。")
        return None
        
    date_col = datetime_cols[0]
    numeric_cols = [col for col in df.select_dtypes(include=np.number).columns if col != date_col]
    if not numeric_cols:
        logger.info(f"日付列 '{date_col}' に対応する数値列がないため、異常検知をスキップします。")
        return None

    value_col = None

    try:
        logger.info("AIに時系列分析に最適な列の選択を依頼します...")
        prompt = f"""あなたはデータ分析の専門家です。ユーザーの分析目的とデータカラムの情報を基に、時系列分析（特に異常検知）に最も適した**単一の数値列**を選択してください。

### ユーザーの分析目的
{analysis_purpose}

### 各列の役割
{json.dumps(column_roles, indent=2, ensure_ascii=False)}

### 分析対象となる数値列の候補と基本統計量
{df[numeric_cols].describe().to_json(indent=2)}

### あなたのタスク
上記の情報を総合的に判断し、分析目的の達成に最も貢献する可能性が高い列を1つだけ選んでください。
例えば、「受注の失注分析」が目的なら、「受注金額」や「受注件数」といった列が適しています。「目標シェア率」のような結果指標や、目的と直接関係ない列は避けるべきです。

### 【最重要】出力形式
選択した列名を、以下のJSON形式で**JSONオブジェクトのみ**を返却してください。解説やその他のテキストは一切含めないでください。
```json
{{
  "selected_column": "（ここに選択した列名）"
}}
```"""

        ai_response_str, _ = generate_with_gemini(
            model_name=config.DATA_ANALYSIS_MODEL_NAME,
            contents=[types.Part.from_text(text=prompt)],
            max_output_tokens=500,
            temperature=0.0,
            generation_config_override={"response_mime_type": "application/json"},
            feature_category="Insight Canvas"
        )
        
        response_json = json.loads(ai_response_str)
        selected_col_by_ai = response_json.get("selected_column")

        if selected_col_by_ai and selected_col_by_ai in numeric_cols:
            value_col = selected_col_by_ai
            logger.info(f"AIが異常検知の対象列として '{value_col}' を選択しました。")
        else:
            logger.warning(f"AIは有効な列を選択できませんでした。AIの応答: '{selected_col_by_ai}'")

    except Exception as e:
        logger.error(f"AIによる列選択中にエラーが発生しました: {e}。フォールバックロジックを実行します。")

    if not value_col:
        logger.info("AIによる列選択に失敗したため、ルールベースのフォールバックロジックを実行します。")
        target_cols = [col for col, role in column_roles.items() if role == '予測ターゲット（数値）' and col in numeric_cols]
        if target_cols:
            value_col = target_cols[0]
            logger.info(f"フォールバック: 役割が'予測ターゲット（数値）'の '{value_col}' を選択しました。")
        
        if not value_col:
            business_keywords = ["売上", "受注", "金額", "数量", "件数", "コスト", "利益", "price", "sales", "amount", "quantity", "count", "cost", "profit"]
            for col in numeric_cols:
                if any(keyword in col.lower() for keyword in business_keywords):
                    value_col = col
                    logger.info(f"フォールバック: ビジネスキーワードを含む '{value_col}' を選択しました。")
                    break
        
        if not value_col:
            value_col = numeric_cols[0]
            logger.info(f"フォールバック: 最初の数値列である '{value_col}' を選択しました。")

    if not value_col:
        logger.info("分析対象の適切な数値列が見つからないため、異常検知をスキップします。")
        return None

    logger.info(f"異常検知対象: 日付='{date_col}', 値='{value_col}'")

    code_to_run = f"""
import pandas as pd
import numpy as np
from statsmodels.tsa.seasonal import seasonal_decompose
import matplotlib.pyplot as plt
import japanize_matplotlib

japanize_matplotlib.japanize()

# 日付でソートし、インデックスに設定
df_ad = df_processed[['{date_col}', '{value_col}']].copy()
df_ad = df_ad.dropna().sort_values(by='{date_col}').set_index('{date_col}')

# 日次などでリサンプリングして集計（欠損を埋めるため）
df_resampled = df_ad['{value_col}'].resample('D').mean().interpolate()

anomalies_result = None

# データが2シーズン未満の場合は実行不可
if len(df_resampled) >= 365 * 2:
    # 季節性分解を実行
    decomposition = seasonal_decompose(df_resampled, model='additive', period=365)
    
    residual = decomposition.resid.dropna()
    mean = residual.mean()
    std = residual.std()
    cut_off = std * 3
    lower_bound, upper_bound = mean - cut_off, mean + cut_off
    
    anomalies_indices = residual[(residual < lower_bound) | (residual > upper_bound)].index
    
    if not anomalies_indices.empty:
        anomalies_result = {{
            "target_column": "{value_col}",
            "dates": anomalies_indices.strftime('%Y-%m-%d').tolist()
        }}
        
        # --- グラフ描画処理 ---
        plt.figure(figsize=(12, 6))
        plt.plot(df_resampled.index, df_resampled.values, label='観測値', color='gray', alpha=0.7)
        plt.plot(decomposition.trend.index, decomposition.trend.values, label='トレンド', color='cornflowerblue')
        
        # 異常点を抽出してプロット
        anomalies_values = df_resampled[df_resampled.index.isin(anomalies_indices)]
        plt.scatter(anomalies_values.index, anomalies_values.values, color='red', s=100, zorder=5, label='検出された異常値')
        
        plt.title('時系列データの異常検知結果 ({value_col})')
        plt.xlabel('日付')
        plt.ylabel('{value_col}')
        plt.legend()
        plt.grid(True, linestyle='--', alpha=0.6)
        plt.tight_layout()
else:
    print("データ期間が短すぎるため、季節分解による異常検知はスキップします。")
"""
    # 3. コードを実行
    exec_result = execute_python_safely(code_to_run, session_id=f"{project_id}_anomaly", global_vars={'df_processed': df.copy()})

    if exec_result.get('error'):
        logger.error(f"異常検知コードの実行中にエラー: {exec_result.get('error')}")
        return None

    # 4. 結果をインサイト形式に整形
    anomalies_result = exec_result.get('final_scope', {}).get('anomalies_result')
    if anomalies_result:
        logger.info(f"異常を {len(anomalies_result['dates'])} 件検出しました。")
        
        charts = [item['content'] for item in exec_result.get('results', []) if item['type'] == 'chart_gcs_uri']
        
        insight = {
            'id': f"insight_{uuid.uuid4().hex}",
            'title': f"【異常検知アラート】'{anomalies_result['target_column']}' の時系列データに特異な変動を検出",
            'description': f"以下の日付で、通常の変動範囲を大きく超える動きが検出されました。特定のイベントやデータ入力ミスの可能性があります。\n- 検出日: {', '.join(anomalies_result['dates'])}",
            'stdout': "".join(item['content'] for item in exec_result.get('results', []) if item['type'] == 'stdout'),
            'error': None,
            'charts': charts,
            'is_anomaly_alert': True,
            'code': code_to_run
        }
        return insight
    
    logger.info("異常は検出されませんでした。")
    return None

def _find_layout_by_name(prs, name):
    """指定された名前のレイアウトを検索する。見つからない場合はNoneを返す。"""
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    logger.warning(f"レイアウト '{name}' がテンプレートに見つかりませんでした。")
    return None

def _parse_markdown_table(md_text):
    """Markdownテーブル文字列をヘッダーと行のリストにパースする。テーブルでなければNoneを返す。"""
    lines = [line.strip() for line in md_text.strip().split('\n') if line.strip()]
    # テーブルの簡易判定: 2行以上で、2行目がセパレータであること
    if len(lines) < 2 or not re.match(r'\|?(\s*:?-+:?\s*\|)+', lines[1]):
        return None, None

    header = [h.strip() for h in lines[0].strip('|').split('|')]
    rows = []
    for line in lines[2:]:
        rows.append([cell.strip() for cell in line.strip('|').split('|')])
        
    # ヘッダーと最初のデータ行の列数が一致するか確認
    if rows and len(header) != len(rows[0]):
        logger.warning("Markdownテーブルのヘッダーとデータ行の列数が一致しません。")
        return None, None

    return header, rows

def _add_table_to_placeholder(slide, placeholder, headers, rows):
    """スライドのプレースホルダーの位置にテーブルを追加し、元のプレースホルダーを削除する。"""
    if not headers or not rows:
        return
    try:
        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
        table = table_shape.table

        # ヘッダーの設定
        for i, header_text in enumerate(headers):
            table.cell(0, i).text = header_text
            # ヘッダーセルのスタイル調整（例）
            table.cell(0, i).text_frame.paragraphs[0].font.bold = True

        # データ行の設定
        for r_idx, row_data in enumerate(rows):
            for c_idx, cell_data in enumerate(row_data):
                if c_idx < len(headers):
                    table.cell(r_idx + 1, c_idx).text = cell_data
        
        # 元のプレースホルダーを削除
        sp = placeholder.element
        sp.getparent().remove(sp)
        logger.info("Markdownテーブルをスライドに描画しました。")
    except Exception as e:
        logger.error(f"テーブルの描画中にエラーが発生しました: {e}", exc_info=True)

def _identify_primary_factors(df: pd.DataFrame, target_variable: str, n_factors: int = 5) -> list:
    """
    RandomForestRegressorを用いて、目的変数に対する重要度が高い特徴量を特定する。
    カテゴリ変数は自動的にダミー変数化され、欠損値も補完される。

    Args:
        df (pd.DataFrame): 分析対象のデータフレーム。
        target_variable (str): 目的変数のカラム名。
        n_factors (int, optional): 抽出する上位の特徴量の数。デフォルトは5。

    Returns:
        list: 重要度が高い特徴量のカラム名のリスト。
    """
    logger.info(f"重要因子の特定を開始します。ターゲット: '{target_variable}', 上位{n_factors}件")
    
    if target_variable not in df.columns:
        logger.error(f"指定されたターゲット変数 '{target_variable}' がデータフレームに存在しません。")
        return []

    try:
        df_copy = df.copy()

        # 目的変数が数値でない場合、回帰モデルは使用できないため処理を中断
        if not pd.api.types.is_numeric_dtype(df_copy[target_variable]):
            logger.warning(f"ターゲット変数 '{target_variable}' は数値ではないため、重要度分析をスキップします。")
            return []

        # ターゲットと特徴量に分割
        y = df_copy[target_variable]
        X = df_copy.drop(columns=[target_variable])

        # ターゲットに欠損値がある行は、特徴量とターゲットの両方から除外
        if y.isnull().any():
            valid_indices = y.dropna().index
            X = X.loc[valid_indices].reset_index(drop=True)
            y = y.loc[valid_indices].reset_index(drop=True)

        # 特徴量の欠損値処理
        # 数値列は中央値で、カテゴリ列は最頻値で補完
        numeric_cols = X.select_dtypes(include=np.number).columns
        for col in numeric_cols:
            if X[col].isnull().any():
                median_val = X[col].median()
                # FutureWarningを回避するため、こちらの書き方に修正
                X[col] = X[col].fillna(median_val)

        categorical_cols = X.select_dtypes(include=['object', 'category']).columns
        for col in categorical_cols:
            if X[col].isnull().any():
                mode_val = X[col].mode().iloc[0] if not X[col].mode().empty else 'Unknown'
                # こちらも同様に修正
                X[col] = X[col].fillna(mode_val)

        # カテゴリ変数をダミー変数化
        X_processed = pd.get_dummies(X, columns=categorical_cols, dummy_na=False)
        
        if X_processed.empty:
            logger.warning("前処理の結果、分析に使用できる特徴量がありませんでした。")
            return []

        # RandomForest回帰モデルの学習 (n_jobs=1 に変更してデッドロックを回避)
        model = RandomForestRegressor(n_estimators=100, random_state=42, n_jobs=1)
        model.fit(X_processed, y)

        # 特徴量の重要度を抽出し、上位N件を取得
        importances = pd.Series(model.feature_importances_, index=X_processed.columns)
        top_factors = importances.sort_values(ascending=False).head(n_factors).index.tolist()

        logger.info(f"特定された重要因子: {top_factors}")
        return top_factors

    except Exception as e:
        logger.error(f"重要因子の特定中に予期せぬエラーが発生しました: {e}", exc_info=True)
        return []



def _add_markdown_text_to_shape(text_frame, markdown_text):
    """
    【改訂版】Markdownテキストを解析し、太字や箇条書きを維持しながら
    python-pptxのTextFrameにRunとして追加する。
    """
    text_frame.clear()
    text_frame.word_wrap = True
    
    # MarkdownをHTMLに変換（tables拡張機能も有効化）
    html = markdown.markdown(markdown_text, extensions=['markdown.extensions.tables'])
    soup = BeautifulSoup(f"<div>{html}</div>", 'html.parser')

    def _apply_nodes_recursive(node, current_paragraph, level=0, is_bold=False, is_italic=False):
        if isinstance(node, NavigableString):
            if str(node).strip():
                run = current_paragraph.add_run()
                run.text = str(node)
                run.font.bold = is_bold
                run.font.italic = is_italic
        elif node.name:
            # ブロックレベル要素の場合、新しい段落を開始
            if node.name in ['p', 'h1', 'h2', 'h3', 'h4', 'ul', 'ol', 'li']:
                # 現在の段落に既にテキストがあれば、新しい段落を作成
                if current_paragraph.text or len(current_paragraph.runs) > 0:
                    current_paragraph = text_frame.add_paragraph()
                
                current_paragraph.level = level
                # リスト項目の場合はインデントと箇条書き記号を設定
                if node.name == 'li':
                    # 新しい段落が作成されたので、テキストをクリアして箇条書き記号を有効にする
                    current_paragraph.text = "" 
                
            new_is_bold = is_bold or node.name in ['strong', 'b']
            new_is_italic = is_italic or node.name in ['em', 'i']
            new_level = level
            if node.name in ['ul', 'ol']:
                new_level += 1

            for child in node.children:
                # テーブルはここで処理しない
                if child.name == 'table':
                    continue
                current_paragraph = _apply_nodes_recursive(child, current_paragraph, new_level, new_is_bold, new_is_italic)
        return current_paragraph

    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    p.text = ""
    p.level = 0
    
    last_paragraph = p
    for content_node in soup.find('div').children:
        # テーブルタグはスキップ（別途処理するため）
        if content_node.name == 'table':
            continue
        last_paragraph = _apply_nodes_recursive(content_node, last_paragraph)

    # 空の段落が残っていれば削除
    if not text_frame.paragraphs[-1].text.strip() and len(text_frame.paragraphs) > 1:
        text_frame._txBody.remove(text_frame.paragraphs[-1]._p)

def find_body_placeholder(slide):
    """
    【改訂版 v3.2 - エラー修正】本文用のプレースホルダーを堅牢な方法で検索する。
    has_title属性の誤用を修正。
    """
    # 1. BODYまたはOBJECTタイプを最優先で検索
    for shape in slide.placeholders:
        if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
            try:
                if shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    logger.info(f"  プレースホルダー検索: BODY/OBJECTタイプ '{shape.name}' を発見。")
                    return shape
            except ValueError:
                continue

    # 2. 名前に 'Content' or '本文' を含むものを検索 (タイトルは除外)
    title_shape = slide.shapes.title if slide.shapes.title else None # 'has_title' ではなく 'title' の存在チェック
    for shape in slide.placeholders:
        if shape != title_shape and shape.has_text_frame:
            if 'Content' in shape.name or '本文' in shape.name:
                logger.info(f"  プレースホルダー検索: 名前 '{shape.name}' を発見。")
                return shape

    # 3. タイトル以外の、テキストフレームを持つ最初のプレースホルダーをフォールバックとして使用
    for shape in slide.placeholders:
        if shape != title_shape and shape.has_text_frame:
            logger.info(f"  プレースホルダー検索: フォールバックとして '{shape.name}' を発見。")
            return shape
    
    logger.warning("スライドに本文用プレースホルダーが見つかりません。")
    return None


def _add_long_text_to_slides(prs, layout, title, long_text):
    """
    長いMarkdownテキストを意味的な単位（見出しなど）で分割し、複数のスライドにわたって追加する。
    """
    text_content = (long_text or "").strip()
    if not text_content:
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = title
        return

    # Markdownの見出し(###)や番号付きリスト(1.)でセクションを分割
    # re.splitはキャプチャグループを含むと区切り文字もリストに残す
    sections = re.split(r'(^###\s.*|^\d+\.\s.*)', text_content, flags=re.MULTILINE)
    
    content_blocks = []
    # 最初のセクション（最初の見出しより前の部分）
    if sections[0].strip():
        content_blocks.append((title, sections[0].strip()))

    # 見出しと内容のペアを作成
    for i in range(1, len(sections), 2):
        heading = sections[i].strip('# ').strip()
        content = sections[i+1].strip()
        content_blocks.append((heading, content))

    # 分割されなかった場合、全体を一つのブロックとする
    if not content_blocks:
        content_blocks.append((title, text_content))

    for heading, content in content_blocks:
        if not content.strip():
            continue

        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = heading
        
        body_ph = find_body_placeholder(slide)
        if body_ph:
            headers, rows = _parse_markdown_table(content)
            if headers and rows:
                _add_table_to_placeholder(slide, body_ph, headers, rows)
            else:
                _add_markdown_text_to_shape(body_ph.text_frame, content)

def _parse_summary_report(report_text: str) -> dict:
    """
    【改訂版】サマリーレポートを、AIの出力揺らぎに強い堅牢な方法で3つのセクションに分割する。
    """
    report_text = report_text.strip()
    
    # セクション名をキーとする辞書を初期化
    parsed_content = {
        "エグゼクティブサマリー": "",
        "分析のストーリー": "",
        "推奨される次のアクション": ""
    }
    
    section_keywords = list(parsed_content.keys())
    
    # キーワードの出現位置を検索
    positions = {}
    for keyword in section_keywords:
        try:
            # "### 1. エグゼクティブサマリー" のような見出し全体ではなく、キーワード自体で検索
            pos = report_text.lower().index(keyword.lower())
            positions[keyword] = pos
        except ValueError:
            continue
            
    # 見つかったキーワードを出現位置でソート
    sorted_keywords = sorted(positions.keys(), key=lambda k: positions[k])
    
    if not sorted_keywords:
        # キーワードが一つも見つからなければ、全文を「分析のストーリー」に入れる
        parsed_content["分析のストーリー"] = report_text
        return parsed_content
        
    # 各セクションの内容をスライスして抽出
    for i, keyword in enumerate(sorted_keywords):
        start_index = positions[keyword]
        
        # 次のキーワードの開始位置を検索、なければテキストの最後まで
        end_index = len(report_text)
        if i + 1 < len(sorted_keywords):
            end_index = positions[sorted_keywords[i+1]]
            
        # キーワード自体の部分（見出し）も含めてセクションの内容を切り出す
        section_full_text = report_text[start_index:end_index].strip()
        
        # 見出し行を除いた本文を取得する
        lines = section_full_text.split('\n')
        # 最初の行が見出しであると仮定して、それ以降の内容をコンテンツとする
        content_text = '\n'.join(lines[1:]).strip()
        
        # もしコンテンツが空であれば、見出し行がなかったと見なして全文を入れる
        if not content_text:
            content_text = section_full_text

        parsed_content[keyword] = content_text
            
    return parsed_content

@data_analyzer_bp.route('/api/export/<project_id>/pptx', methods=['GET'])
def export_pptx_api(project_id):
    """
    【改修版】分析結果をPPTX形式でエクスポートする。
    クラスター分析、予測モデル、時系列予測の結果もスライドに含めるように修正。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    try:
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
             return jsonify({"error": "エクスポート対象のプロジェクトデータが見つかりません。"}), 404

        # プロファイリング結果をBlob Storageまたは古いフィールドから読み込む
        profiling_result = None
        if blob_name := project_data.get("profilingResultBlobName"):
            try:
                result_bytes = download_blob_to_bytes(blob_name)
                if result_bytes:
                    profiling_result = json.loads(result_bytes.decode('utf-8'))
            except Exception as e:
                logger.error(f"PPTXエクスポート中にプロファイリング結果の読み込みに失敗: {e}", exc_info=True)
        
        # 後方互換性のため、古いフィールドもチェック
        if profiling_result is None and "profilingResult" in project_data:
            profiling_result = project_data.get("profilingResult")

        # インサイトをBlobから読み込む
        insights_data = {}
        if insights_blob_name := project_data.get("insightsBlobName"):
            try:
                insights_bytes = download_blob_to_bytes(insights_blob_name)
                if insights_bytes:
                    insights_data = json.loads(insights_bytes.decode('utf-8'))
            except Exception as e:
                 logger.error(f"PPTXエクスポート中にインサイトの読み込みに失敗: {e}", exc_info=True)
                 # エラーでも処理を続行、インサイトなしでPPTXが生成される
        # 後方互換性
        elif "insights" in project_data:
            insights_data = project_data.get("insights", {})

        # クラスター分析結果をBlobから読み込む
        cluster_result = None
        if cluster_blob_name := project_data.get("clusterAnalysisResultBlobName"):
            try:
                cluster_bytes = download_blob_to_bytes(cluster_blob_name)
                if cluster_bytes:
                    cluster_result = json.loads(cluster_bytes.decode('utf-8'))
            except Exception as e:
                logger.error(f"PPTXエクスポート中にクラスター分析結果の読み込みに失敗: {e}", exc_info=True)
        # 後方互換性
        elif "clusterAnalysisResult" in project_data:
            cluster_result = project_data.get("clusterAnalysisResult")

        if not insights_data and not project_data.get("summaryReport"):
             return jsonify({"error": "エクスポート対象の分析結果が見つかりません。"}), 404

        # --- 1. プレゼンテーションとレイアウトの準備 ---
        prs = None
        template_path = getattr(config, 'PPTX_GENERATOR_DEFAULT_TEMPLATE', 'default_template.pptx')
        if os.path.exists(template_path):
            try:
                prs = Presentation(template_path)
                logger.info(f"テンプレート '{template_path}' を読み込みました。")
            except Exception as e:
                prs = Presentation()
                logger.error(f"テンプレート読み込み失敗: {e}。空のプレゼンテーションを作成します。")
        else:
            prs = Presentation()
            logger.warning(f"テンプレート '{template_path}' が見つかりません。空のプレゼンテーションを作成します。")
            
        layout_title_content = _find_layout_by_name(prs, 'Layout_1_Title_And_BulletedText') or prs.slide_layouts[1]
        layout_title_two_content = _find_layout_by_name(prs, 'Layout_6_Title_And_Two_Contents') or prs.slide_layouts[3] if len(prs.slide_layouts) > 5 else prs.slide_layouts[5]

        # --- 2. 1枚目: プロジェクト概要スライド ---
        slide = prs.slides.add_slide(layout_title_content)
        if slide.shapes.title:
            slide.shapes.title.text = "プロジェクト概要"
        
        body_shape = find_body_placeholder(slide)
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            
            p1 = tf.add_paragraph(); p1.text = "分析目的"; p1.font.bold = True; p1.level = 0
            p2 = tf.add_paragraph(); p2.text = project_data.get('analysisPurpose', 'N/A'); p2.level = 1
            tf.add_paragraph() # 空行
            
            p3 = tf.add_paragraph(); p3.text = "参考資料"; p3.font.bold = True; p3.level = 0
            ref_files = project_data.get('referenceFiles', [])
            ref_text = ", ".join([f.get('filename', '不明') for f in ref_files]) if ref_files else "なし"
            p4 = tf.add_paragraph(); p4.text = ref_text; p4.level = 1
            tf.add_paragraph() # 空行

            p5 = tf.add_paragraph(); p5.text = "データ概要"; p5.font.bold = True; p5.level = 0
            profiling_summary = (profiling_result or {}).get('summary', {})
            row_count = profiling_summary.get('rowCount', 'N/A')
            col_count = profiling_summary.get('columnCount', 'N/A')
            p6 = tf.add_paragraph(); p6.text = f"行数: {row_count}, 列数: {col_count}"; p6.level = 1
            
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # --- 3. サマリーレポートの分割とスライド作成 ---
        summary_report = project_data.get("summaryReport", "")
        if summary_report:
            parsed_summary = _parse_summary_report(summary_report)
            
            for title, content in parsed_summary.items():
                if not content.strip(): continue
                
                # Markdownテーブルの存在をチェック
                table_match = re.search(r'(\|.*\|[\r\n]+)\|.*\|', content)

                # テーブルが見つかった場合、2コンテンツレイアウトを使用
                if title == "推奨される次のアクション" and table_match:
                    slide = prs.slides.add_slide(layout_title_two_content)
                    if slide.shapes.title:
                        slide.shapes.title.text = title

                    # テーブルより前のテキストと、テーブル自体のテキストに分割
                    text_before_table = content[:table_match.start()].strip()
                    table_markdown = content[table_match.start():].strip()

                    # 左と右のプレースホルダーを特定
                    text_ph, table_ph = None, None
                    title_shape = slide.shapes.title
                    content_phs = sorted([p for p in slide.placeholders if p != title_shape], key=lambda p: p.left)
                    if len(content_phs) >= 2:
                        text_ph, table_ph = content_phs[0], content_phs[1]
                    
                    # 左側にテキスト部分を描画
                    if text_ph and text_before_table:
                        _add_markdown_text_to_shape(text_ph.text_frame, text_before_table)

                    # 右側にテーブル部分を描画
                    if table_ph and table_markdown:
                        headers, rows = _parse_markdown_table(table_markdown)
                        if headers and rows:
                            _add_table_to_placeholder(slide, table_ph, headers, rows)
                        else: # テーブル解析失敗時はテキストとして描画
                            _add_markdown_text_to_shape(table_ph.text_frame, table_markdown)

                # 「分析のストーリー」の場合、複数スライドに分割
                elif title == "分析のストーリー":
                    _add_long_text_to_slides(prs, layout_title_content, title, content)

                # それ以外のセクションは従来通り1枚のスライドに出力
                else:
                    slide = prs.slides.add_slide(layout_title_content)
                    if slide.shapes.title:
                        slide.shapes.title.text = title
                    
                    body_ph = find_body_placeholder(slide)
                    if body_ph:
                        _add_markdown_text_to_shape(body_ph.text_frame, content)

        # プロジェクトデータからインフォグラフィックのBlob名を取得
        infographic_blobs = []
        if b1 := project_data.get("infographicBlobName1"): infographic_blobs.append({"name": b1, "title": "分析サマリー・インフォグラフィック (概況)"})
        # 後方互換
        elif b_old := project_data.get("infographicBlobName"): infographic_blobs.append({"name": b_old, "title": "分析サマリー・インフォグラフィック"})
        
        if b2 := project_data.get("infographicBlobName2"): infographic_blobs.append({"name": b2, "title": "推奨される次のアクション・インフォグラフィック"})

        for info in infographic_blobs:
            blob_name = info["name"]
            slide_title = info["title"]
            try:
                infographic_bytes = download_blob_to_bytes(blob_name)
                if infographic_bytes:
                    slide = prs.slides.add_slide(layout_title_two_content)
                    if slide.shapes.title: slide.shapes.title.text = slide_title
                    
                    # 画像配置ロジック (既存と同じ)
                    placeholders = [ph for ph in slide.placeholders if ph.element.tag.endswith('body')]
                    if placeholders:
                        target_ph = placeholders[0]
                        left, top, width, height = target_ph.left, target_ph.top, target_ph.width, target_ph.height
                        slide.shapes.add_picture(io.BytesIO(infographic_bytes), left, top, width=width, height=height)
                        sp = target_ph.element; sp.getparent().remove(sp)
                        if len(placeholders) > 1: sp2 = placeholders[1].element; sp2.getparent().remove(sp2)
                    else:
                        margin = Inches(0.5)
                        pic_width = prs.slide_width - (margin * 2)
                        pic_height = prs.slide_height - (margin * 2)
                        top = margin
                        if slide.shapes.title: top = slide.shapes.title.top + slide.shapes.title.height + Inches(0.2); pic_height = prs.slide_height - top - margin
                        slide.shapes.add_picture(io.BytesIO(infographic_bytes), margin, top, width=pic_width, height=pic_height)
            except Exception as e:
                logger.error(f"インフォグラフィック({blob_name})のPPTX追加中にエラー: {e}", exc_info=True)

        # --- 4. クラスター分析結果のスライド ---
        cluster_result = None
        if cluster_blob_name := project_data.get("clusterAnalysisResultBlobName"):
            try:
                cluster_bytes = download_blob_to_bytes(cluster_blob_name)
                if cluster_bytes:
                    cluster_result = json.loads(cluster_bytes.decode('utf-8'))
            except Exception as e:
                logger.error(f"PPTXエクスポート中にクラスター分析結果の読み込みに失敗: {e}", exc_info=True)
        # 後方互換性のため、古いフィールドもチェック
        elif "clusterAnalysisResult" in project_data:
            cluster_result = project_data.get("clusterAnalysisResult")

        if cluster_result and not cluster_result.get('error') and cluster_result.get('summary'):
            slide = prs.slides.add_slide(layout_title_two_content)
            if slide.shapes.title:
                slide.shapes.title.text = "クラスター分析：顧客セグメントの発見"
            
            # プレースホルダーを安全に取得
            text_ph, chart_ph = None, None
            title_shape = slide.shapes.title
            content_phs = sorted([p for p in slide.placeholders if p != title_shape], key=lambda p: p.left)
            if len(content_phs) >= 2:
                text_ph, chart_ph = content_phs[0], content_phs[1]

            if text_ph:
                cluster_descriptions = ""
                for cluster in cluster_result.get('summary', {}).get('clusters', []):
                    name = cluster.get('name', f"クラスター {cluster.get('cluster_id', '')}")
                    description = cluster.get('description', '解説はありません。')
                    cluster_descriptions += f"**{name}**\n\n{description}\n\n"
                _add_markdown_text_to_shape(text_ph.text_frame, cluster_descriptions)

            if chart_ph:
                gcs_chart_uri = cluster_result.get('cluster_chart_url')
                if gcs_chart_uri:
                    try:
                        image_bytes = download_gcs_blob_to_bytes(gcs_chart_uri)
                        if image_bytes:
                            slide.shapes.add_picture(io.BytesIO(image_bytes), chart_ph.left, chart_ph.top, chart_ph.width, chart_ph.height)
                            sp = chart_ph.element; sp.getparent().remove(sp)
                    except Exception as e:
                        logger.error(f"クラスター分析グラフの挿入に失敗: {e}", exc_info=True)
                        if chart_ph.has_text_frame: chart_ph.text_frame.text = "グラフ画像の表示に失敗しました。"
        
        # --- 5. 予測モデル結果のスライド (回帰・分類) ---
        regression_result = project_data.get("regressionAnalysisResult")
        if regression_result and not regression_result.get('error'):
            slide = prs.slides.add_slide(layout_title_two_content)
            if slide.shapes.title:
                slide.shapes.title.text = f"予測モデル結果: {regression_result.get('target_variable', '')}"

            text_ph, chart_ph = None, None
            title_shape = slide.shapes.title
            content_phs = sorted([p for p in slide.placeholders if p != title_shape], key=lambda p: p.left)
            if len(content_phs) >= 2:
                text_ph, chart_ph = content_phs[0], content_phs[1]

            if text_ph:
                accuracy_label = '精度 (R²)' if regression_result.get('model_type') == 'regression' else '精度 (Accuracy)'
                model_summary = (
                    f"**予測ターゲット:** {regression_result.get('target_variable', 'N/A')}\n\n"
                    f"**モデルタイプ:** {regression_result.get('model_type', 'N/A')}\n\n"
                    f"**アルゴリズム:** {regression_result.get('algorithm', 'N/A')}\n\n"
                    f"**{accuracy_label}:** {regression_result.get('accuracy_score', 0):.3f}\n\n"
                    f"### AIによる解説\n\n{regression_result.get('commentary', '')}"
                )
                _add_markdown_text_to_shape(text_ph.text_frame, model_summary)
            
            if chart_ph:

                chart_blob_name = regression_result.get('feature_importance_chart_blob_name')
                image_bytes = None

                if chart_blob_name:
                    try:
                        image_bytes = download_blob_to_bytes(chart_blob_name)
                    except Exception as e:
                        logger.error(f"予測モデルグラフ(Blob)のダウンロードに失敗: {e}")

                # フォールバック: GCS URI (旧データ互換)
                if not image_bytes:
                    gcs_chart_uri = regression_result.get('feature_importance_chart_url_gcs')
                    if gcs_chart_uri:
                        try:
                            image_bytes = download_gcs_blob_to_bytes(gcs_chart_uri)
                        except Exception as e:
                            logger.error(f"予測モデルグラフ(GCS)のダウンロードに失敗: {e}")

                if image_bytes:
                    try:
                        slide.shapes.add_picture(io.BytesIO(image_bytes), chart_ph.left, chart_ph.top, chart_ph.width, chart_ph.height)
                        sp = chart_ph.element; sp.getparent().remove(sp)
                    except Exception as e:
                         logger.error(f"予測モデルグラフの挿入に失敗: {e}", exc_info=True)
                         if chart_ph.has_text_frame: chart_ph.text_frame.text = "グラフ画像の表示に失敗しました。"

        # --- 6. 時系列予測結果のスライド ---
        timeseries_result = project_data.get("timeseriesForecastResult")
        if timeseries_result and not timeseries_result.get('error'):
            slide = prs.slides.add_slide(layout_title_two_content)
            if slide.shapes.title:
                slide.shapes.title.text = f"時系列予測結果: {timeseries_result.get('target_variable', '')}"

            text_ph, chart_ph = None, None
            title_shape = slide.shapes.title
            content_phs = sorted([p for p in slide.placeholders if p != title_shape], key=lambda p: p.left)
            if len(content_phs) >= 2:
                text_ph, chart_ph = content_phs[0], content_phs[1]
            
            if text_ph:
                summary_text = (
                    f"**予測ターゲット:** {timeseries_result.get('target_variable', 'N/A')}\n\n"
                    f"**モデルタイプ:** 時系列予測\n\n"
                    f"**アルゴリズム:** {timeseries_result.get('algorithm', 'N/A')}\n\n"
                    f"**{timeseries_result.get('accuracy_score_label', '精度')}:** {timeseries_result.get('accuracy_score', 0) * 100:.1f}%\n\n"
                    f"### AIによる解説\n\n{timeseries_result.get('commentary', '')}"
                )
                _add_markdown_text_to_shape(text_ph.text_frame, summary_text)

            if chart_ph:
                image_bytes = None
                
                # 1. GCS URIがあればそれを使う
                gcs_chart_uri = timeseries_result.get('feature_importance_chart_url_gcs')
                if gcs_chart_uri:
                     image_bytes = download_gcs_blob_to_bytes(gcs_chart_uri)
                
                # 2. なければ SAS URLからダウンロードを試みる (requestsを使用)
                if not image_bytes:
                    sas_url = timeseries_result.get('feature_importance_chart_url')
                    if sas_url:
                        try:
                            import requests
                            resp = requests.get(sas_url)
                            if resp.status_code == 200:
                                image_bytes = resp.content
                        except Exception as e:
                            logger.error(f"時系列グラフ(SAS URL)のダウンロード失敗: {e}")

                if image_bytes:
                    try:
                        slide.shapes.add_picture(io.BytesIO(image_bytes), chart_ph.left, chart_ph.top, chart_ph.width, chart_ph.height)
                        sp = chart_ph.element; sp.getparent().remove(sp)
                    except Exception as e:
                        logger.error(f"時系列予測グラフの挿入に失敗: {e}", exc_info=True)

        # --- 7. 全インサイトをリストに集約 & スライド作成 ---
        all_insights = []
        #insights_data = project_data.get("insights", {})
        if isinstance(insights_data, dict):
            for factor_group in insights_data.get("primary_factor_insights", []):
                all_insights.extend(factor_group.get("insights", []))
            for category_group in insights_data.get("categorized_insights", []):
                all_insights.extend(category_group.get("insights", []))
            all_insights.extend(insights_data.get("anomaly_alerts", []))
        elif isinstance(insights_data, list):
            all_insights = insights_data # 後方互換性
        
        for insight in all_insights:
            if not insight or insight.get('error'): continue
            
            title = insight.get('title', '無題のインサイト')
            description = insight.get('description', '')
            gcs_chart_uris = insight.get('charts', [])
            has_chart = bool(gcs_chart_uris)
            
            headers, rows = _parse_markdown_table(description)
            is_table = bool(headers and rows)

            current_layout = layout_title_two_content if has_chart else layout_title_content
            slide = prs.slides.add_slide(current_layout)
            if slide.shapes.title:
                slide.shapes.title.text = title

            if has_chart:
                text_ph, chart_ph = None, None
                title_shape = slide.shapes.title
                content_phs = sorted([p for p in slide.placeholders if p != title_shape], key=lambda p: p.left)
                if len(content_phs) >= 2:
                    text_ph, chart_ph = content_phs[0], content_phs[1]
                elif len(content_phs) == 1:
                    text_ph = content_phs[0]

                if text_ph:
                    if is_table:
                        _add_table_to_placeholder(slide, text_ph, headers, rows)
                    else:
                        _add_markdown_text_to_shape(text_ph.text_frame, description)
                
                if chart_ph:
                    try:
                        image_bytes = download_gcs_blob_to_bytes(gcs_chart_uris[0])
                        if image_bytes:
                            slide.shapes.add_picture(io.BytesIO(image_bytes), chart_ph.left, chart_ph.top, chart_ph.width, chart_ph.height)
                            sp = chart_ph.element; sp.getparent().remove(sp)
                    except Exception as e:
                        logger.error(f"インサイトグラフの挿入に失敗: {e}", exc_info=True)
                        if chart_ph.has_text_frame: chart_ph.text_frame.text = "グラフ画像の表示に失敗しました。"
            else:
                body_ph = find_body_placeholder(slide)
                if body_ph:
                    if is_table:
                        _add_table_to_placeholder(slide, body_ph, headers, rows)
                    else:
                        _add_markdown_text_to_shape(body_ph.text_frame, description)

        # --- 8. PPTXファイルをメモリ上で保存して送信 ---
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)

        base_name = os.path.splitext(project_data.get('projectName', 'report'))[0]
        if not base_name:
            base_name = project_id
        download_name = f"{base_name}_analysis_report.pptx"

        return send_file(
            file_stream, as_attachment=True, download_name=download_name,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        logger.error(f"PPTXエクスポート中にエラー (プロジェクトID: {project_id}): {e}", exc_info=True)
        return jsonify({"error": "PPTXファイルのエクスポート中にサーバーでエラーが発生しました。"}), 500

def _run_prediction_task(df: pd.DataFrame, target_variable: str, model_type: str, project_id: str, n_factors: int = 5) -> tuple[dict | None, list]:
    """
    データフレーム、目的変数、モデルタイプを受け取り、予測モデルを構築・評価する。
    【改訂版】
    - One-Hotエンコーディング等の前処理をパイプラインに組み込む。
    - UIでのWhat-if分析用に、'変換前の' 特徴量名と詳細情報をメタデータとして保存する。
    """

    # 再現用のPythonコード文字列を作成
    reproduction_code = f"""# --- 予測モデル構築コード（再現用） ---
# 注意: このコードはシステム内部で実行された処理を再現するためのものです。
# 実際にはサーバー上の環境で実行されました。

import pandas as pd
from sklearn.model_selection import train_test_split
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.ensemble import RandomForest{'Regressor' if model_type=='regression' else 'Classifier'}
from sklearn.metrics import {'mean_absolute_percentage_error' if model_type=='regression' else 'f1_score'}

# 1. データの準備
target = '{target_variable}'
X = df.drop(columns=[target])
y = df[target]

# 数値・カテゴリ列の定義
numeric_features = X.select_dtypes(include=['number']).columns.tolist()
categorical_features = X.select_dtypes(include=['object', 'category']).columns.tolist()

# 2. 前処理パイプライン
numeric_transformer = Pipeline(steps=[
    ('imputer', SimpleImputer(strategy='median')),
    ('scaler', StandardScaler())
])

categorical_transformer = Pipeline(steps=[
    ('imputer', SimpleImputer(strategy='most_frequent')),
    ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
])

preprocessor = ColumnTransformer(
    transformers=[
        ('num', numeric_transformer, numeric_features),
        ('cat', categorical_transformer, categorical_features)
    ]
)

# 3. モデル定義
model = RandomForest{'Regressor' if model_type=='regression' else 'Classifier'}(random_state=42, n_jobs=-1)

clf = Pipeline(steps=[('preprocessor', preprocessor), ('classifier', model)])

# 4. 学習と評価
X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
clf.fit(X_train, y_train)

print(f"Model Score: {{clf.score(X_test, y_test):.4f}}")
"""

    logger.info(f"プロジェクト {project_id} のための予測モデル構築タスクを開始します。")
    logger.info(f"  - 目的変数: {target_variable}")
    logger.info(f"  - モデルタイプ: {model_type}")
    logger.info(f"  - 主要因の抽出数: {n_factors}")

    try:
        # --- 1. データ準備 ---
        if target_variable not in df.columns:
            raise ValueError(f"目的変数 '{target_variable}' がデータに存在しません。")

        # 日付列を除外
        df_for_prediction = df.select_dtypes(exclude=['datetime64[ns]', 'datetime64[ns, UTC]'])
        # 目的変数の欠損を除外
        df_for_prediction = df_for_prediction.dropna(subset=[target_variable])
        
        if df_for_prediction.empty:
            raise ValueError("目的変数の欠損値を除外した結果、データが0件になりました。")
        
        if target_variable not in df_for_prediction.columns:
            if pd.api.types.is_datetime64_any_dtype(df[target_variable].dtype):
                 raise ValueError(f"目的変数 '{target_variable}' が日付/時刻型のため、この予測モデルではターゲットとして使用できません。")
            else:
                 raise ValueError(f"目的変数 '{target_variable}' がデータから失われました。")

        X = df_for_prediction.drop(columns=[target_variable])
        y = df_for_prediction[target_variable]

        numeric_features = X.select_dtypes(include=[np.number, 'bool']).columns.tolist()
        categorical_features = X.select_dtypes(include=['object', 'category']).columns.tolist()
        
        # ★★★ 修正箇所: What-if分析用に、変換前の元のカラムリストを保持する ★★★
        model_original_features = X.columns.tolist()

        # --- 2. 前処理パイプラインの構築 ---
        numeric_transformer = Pipeline(steps=[
            ('imputer', SimpleImputer(strategy='median')),
            ('scaler', StandardScaler())
        ])
        categorical_transformer = Pipeline(steps=[
            ('imputer', SimpleImputer(strategy='most_frequent')),
            # handle_unknown='ignore'により、学習時にないカテゴリが来てもエラーにせず無視する
            ('onehot', OneHotEncoder(handle_unknown='ignore', sparse_output=False))
        ])
        
        preprocessor = ColumnTransformer(
            transformers=[
                ('num', numeric_transformer, numeric_features),
                ('cat', categorical_transformer, categorical_features)
            ],
            remainder='drop' 
        )

        # --- 3. モデルの選択と完全なパイプラインの作成 ---
        if model_type == 'regression':
            model = RandomForestRegressor(random_state=42, n_jobs=1)
            algorithm_name = "Random Forest Regressor"
        else:
            model = RandomForestClassifier(random_state=42, n_jobs=1)
            algorithm_name = "Random Forest Classifier"
        
        # 前処理とモデルを結合したパイプライン
        pipeline = Pipeline(steps=[('preprocessor', preprocessor), ('classifier', model)])

        # --- 4. 学習の実行 (テストデータ分割あり) ---
        X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42)
        
        # ★★★ 修正箇所: UIフォーム生成用の特徴量詳細情報を、変換前のデータから作成する ★★★
        feature_details = {}
        MAX_CATEGORIES_FOR_SELECT = 50
        
        for feature in model_original_features:
            if feature not in X_train.columns:
                continue

            col_type = X_train[feature].dtype
            if pd.api.types.is_numeric_dtype(col_type):
                # 数値型: 最小・最大・中央値を保存（スライダーの範囲に使用）
                feature_details[feature] = {
                    "type": "numeric", 
                    "min": float(X_train[feature].min()),
                    "max": float(X_train[feature].max()), 
                    "median": float(X_train[feature].median())
                }
            else:
                # カテゴリ型: ユニークな値をリストとして保存（プルダウンの選択肢に使用）
                unique_values = X_train[feature].dropna().unique().tolist()
                mode_value = X_train[feature].mode().iloc[0] if not X_train[feature].mode().empty else ""
                
                if len(unique_values) > 50: 
                    feature_details[feature] = {"type": "text", "mode": str(mode_value)}
                else:
                    feature_details[feature] = {
                        "type": "categorical", 
                        "categories": sorted([str(v) for v in unique_values]),
                        "mode": str(mode_value)
                    }
        
        # パイプラインで学習
        pipeline.fit(X_train, y_train)

        # --- 5. 評価指標の計算 ---
        y_pred = pipeline.predict(X_test)
        if model_type == 'regression':
            accuracy_score_val = pipeline.score(X_test, y_test)
            error_metric_val = mean_absolute_percentage_error(y_test, y_pred)
        else:
            accuracy_score_val = pipeline.score(X_test, y_test)
            error_metric_val = f1_score(y_test, y_pred, average='weighted')

        # --- 6. 特徴量重要度の計算と主要因の特定 ---
        top_factors = []
        original_feature_importances = {}
        try:
            # OneHot変換後の特徴量名を取得
            transformed_feature_names = pipeline.named_steps['preprocessor'].get_feature_names_out()
            importances = pipeline.named_steps['classifier'].feature_importances_
            
            # OneHotで分割された特徴量の重要度を、元のカラム名に集約する
            for name, imp in zip(transformed_feature_names, importances):
                # 'num__', 'cat__' などのプレフィックスを除去
                original_name = re.sub(r'^(num__|cat__|remainder__)', '', name)
                
                # カテゴリ変数の場合、OneHotにより 'ColName_Value' のようになるため、元の 'ColName' を復元
                if any(cat_feat in name for cat_feat in categorical_features):
                    # アンダースコア区切りの最後の要素（値部分）を除去して元のカラム名とみなす簡易ロジック
                    # (カラム名自体にアンダースコアが含まれる場合への完全な対応は複雑だが、ここでは近似的に処理)
                    original_name = re.sub(r'_[^_]*$', '', original_name)
                
                original_feature_importances.setdefault(original_name, 0)
                original_feature_importances[original_name] += imp

            # 重要度の高い順にソート
            sorted_original_factors = sorted(original_feature_importances.items(), key=lambda item: item[1], reverse=True)
            top_factors = [name for name, _ in sorted_original_factors[:n_factors]]
            logger.info(f"モデルから特定された主要因: {top_factors}")

        except Exception as e_feat:
            logger.error(f"特徴量重要度の計算または主要因の特定に失敗: {e_feat}", exc_info=True)

        # --- 7. 特徴量重要度の可視化 ---
        feature_importance_chart_url, chart_blob_name = None, None
        if original_feature_importances:
            try:
                df_for_plot = pd.DataFrame(
                    list(original_feature_importances.items()),
                    columns=['feature', 'importance']
                ).sort_values('importance', ascending=False).head(20)

                plt.figure(figsize=(10, 8))
                sns.barplot(x='importance', y='feature', data=df_for_plot)
                plt.title(f'Feature Importance for {target_variable} (Top 20)')
                plt.tight_layout()
                
                buf = io.BytesIO()
                plt.savefig(buf, format='png')
                plt.close()
                buf.seek(0)
                
                chart_blob_name = f"charts/{project_id}/feature_importance_{uuid.uuid4().hex}.png"
                upload_bytes_to_blob(buf.getvalue(), chart_blob_name)
            except Exception as e_chart:
                logger.error(f"特徴量重要度グラフの生成またはアップロードに失敗: {e_chart}", exc_info=True)

        # --- 8. AIによる解説生成 ---
        top_feature_name = top_factors[0] if top_factors else "N/A"
        if model_type == 'regression':
            commentary = f"AIによる自動生成された解説: 回帰モデルを構築し、テストデータで評価した結果、精度(R²)は **{accuracy_score_val:.2f}** となりました。また、予測値が実際の値から平均して **{error_metric_val:.2%}** 程度乖離していることを示すMAPE（平均絶対パーセント誤差）が算出されました。予測に最も影響を与えているのは「**{top_feature_name}**」です。"
        else:
            commentary = f"AIによる自動生成された解説: 分類モデルを構築し、テストデータで評価した結果、精度(正解率)は **{accuracy_score_val:.2f}** となりました。予測に最も影響を与えているのは「**{top_feature_name}**」です。"

        # --- 9. パイプラインの保存と結果の整形 ---
        pipeline_blob_name = None
        try:
            pipeline_buffer = io.BytesIO()
            joblib.dump(pipeline, pipeline_buffer)
            pipeline_buffer.seek(0)
            pipeline_blob_name = f"models/{project_id}/pipeline_v1_{uuid.uuid4().hex[:8]}.joblib"
            upload_bytes_to_blob(pipeline_buffer.getvalue(), pipeline_blob_name)
        except Exception as e_save:
            logger.error(f"パイプラインの保存に失敗: {e_save}", exc_info=True)

        model_id = f"model_{project_id}_{uuid.uuid4().hex[:8]}"
        result = {
            "id": model_id, 
            "model_id": model_id, 
            "project_id": project_id,
            "projectName": df.attrs.get('name', project_id), 
            "target_variable": target_variable,
            "model_type": model_type, 
            "algorithm": algorithm_name, 
            "accuracy_score": float(accuracy_score_val),
            "error_metric": float(error_metric_val), 
            
            # ★★★ 修正箇所: 変換前の特徴量名を保存し、UIがこれを利用できるようにする ★★★
            "feature_names": model_original_features,
            "feature_details": feature_details, 
            
            "commentary": commentary,
            "feature_importance_chart_blob_name": chart_blob_name,
            "pipeline_blob_name": pipeline_blob_name,
            "createdAt": datetime.now(timezone.utc).isoformat(), 
            "version": 1, 
            "error": None,
            "code": reproduction_code,
            "type": "PredictionModel",
        }
        
        return result, top_factors

    except Exception as e:
        logger.error(f"予測モデルタスク {project_id} でエラー: {e}", exc_info=True)
        return {"error": str(e)}, []


@data_analyzer_bp.route('/api/models', methods=['GET'])
def get_user_models_api():
    """
    現在のユーザーに紐づく、過去に作成された予測モデルの一覧を返す。
    InMemoryContainerとCosmos DBの両方に対応。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    
    try:
        container = get_analyzer_container()
        if not container:
            raise Exception("データベースコンテナに接続できません。")

        # --- ここからが改修箇所 ---

        # InMemoryContainerの場合、手動でフィルタリングとソートを行う
        if isinstance(container, InMemoryContainer):
            with container.lock:
                all_items = list(container.store.values())
            
            user_models = [
                item for item in all_items 
                if item.get("userId") == user_id and item.get("type") == "PredictionModel"
            ]
            # createdAtで降順ソート
            user_models.sort(key=lambda x: x.get('createdAt', '1970-01-01T00:00:00.000Z'), reverse=True)
            
            logger.info(f"インメモリDBからユーザー {user_id} のために {len(user_models)} 件の予測モデルを取得しました。")
            return jsonify(user_models)

        # Cosmos DB を使用する場合 (本番環境)
        query = "SELECT * FROM c WHERE c.userId = @userId AND c.type = 'PredictionModel' ORDER BY c.createdAt DESC"
        parameters = [{"name": "@userId", "value": user_id}]
        
        models = list(container.query_items(
            query=query,
            parameters=parameters,
            partition_key=user_id # パーティションキーを指定して効率化
        ))
        
        logger.info(f"Cosmos DBからユーザー {user_id} のために {len(models)} 件の予測モデルを取得しました。")
        return jsonify(models)

    except Exception as e:
        logger.error(f"ユーザー {user_id} のモデル一覧取得中にエラー: {e}", exc_info=True)
        return jsonify({"error": "モデル一覧の取得に失敗しました。"}), 500

@data_analyzer_bp.route('/api/models/<model_id>', methods=['GET'])
def get_model_detail_api(model_id):
    """
    指定されたモデルIDのメタデータを返す。
    （過去の安定バージョンに基づき修正）
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    if not model_id:
        return jsonify({"error": "モデルIDが指定されていません。"}), 400

    try:
        model_data = get_analysis_project(project_id=model_id, user_id=user_id)
        
        if model_data and model_data.get('type') == 'PredictionModel':
            logger.info(f"モデル詳細を取得しました: {model_id}")

            # ★★★ 修正箇所 ★★★
            # パイプラインを動的に読み込んで feature_names を上書きする処理を削除。
            # _run_prediction_task で保存された `feature_names` と `feature_details` を信頼し、そのまま使用する。
            
            # グラフのBlob名から、毎回新しい有効なSAS URLを生成して返す
            if blob_name := model_data.get('feature_importance_chart_blob_name'):
                sas_url = get_blob_sas_url(blob_name)
                if sas_url:
                    # 'feature_importance_chart_url' キーを動的に追加または上書きする
                    model_data['feature_importance_chart_url'] = sas_url

            # データベースから取得した `model_data` をそのまま返す
            return jsonify(model_data)
        else:
            logger.warning(f"モデルが見つからないか、アクセス権がありません: {model_id}")
            return jsonify({"error": "モデルが見つからないか、アクセス権がありません。"}), 404

    except Exception as e:
        logger.error(f"モデル詳細 {model_id} の取得中にエラー: {e}", exc_info=True)
        return jsonify({"error": "モデル詳細の取得に失敗しました。"}), 500

@data_analyzer_bp.route('/api/models/<model_id>/simulate', methods=['POST'])
def simulate_model_api(model_id):
    """
    指定された特徴量の値に基づき、What-ifシミュレーションを実行する。
    【改訂版】
    - 分類モデルにおけるSHAP値の多次元構造（クラスごとの寄与度）を正しく処理し、
      waterfallプロットに単一の説明オブジェクトを渡すように修正。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")
    data = request.get_json()
    features = data.get('features') 

    if not features:
        return jsonify({"error": "シミュレーション対象の特徴量がありません。"}), 400

    try:
        # 1. モデルとパイプラインをロード
        model_data = get_analysis_project(project_id=model_id, user_id=user_id)
        if not model_data or not model_data.get('pipeline_blob_name'):
            return jsonify({"error": "モデルが見つからないか、パイプラインが保存されていません。"}), 404

        pipeline_blob_name = model_data['pipeline_blob_name']
        pipeline_bytes = download_blob_to_bytes(pipeline_blob_name)
        if not pipeline_bytes:
            return jsonify({"error": "パイプラインファイルの読み込みに失敗しました。"}), 500
        
        pipeline = joblib.load(io.BytesIO(pipeline_bytes))

        # 2. 入力データをモデルが期待するカラム構成のDataFrameに変換
        if hasattr(pipeline.named_steps['preprocessor'], 'feature_names_in_'):
            model_feature_names = pipeline.named_steps['preprocessor'].feature_names_in_
        else:
            model_feature_names = model_data.get('feature_names', [])
        
        input_df_raw = pd.DataFrame([features])
        input_df = input_df_raw.reindex(columns=model_feature_names)

        # 3. 予測を実行
        prediction_array = pipeline.predict(input_df)
        prediction = prediction_array[0] if isinstance(prediction_array, (np.ndarray, list)) else prediction_array

        model_type = model_data.get('model_type', 'regression')
        if model_type == 'regression':
            final_prediction = float(prediction)
        else: # classification
            final_prediction = str(prediction)

        # 4. SHAP値を計算し、プロットを生成
        shap_plot_url = None
        try:
            preprocessor = pipeline.named_steps['preprocessor']
            model = pipeline.named_steps['classifier']
            
            # 前処理ステップのみを適用
            X_transformed = preprocessor.transform(input_df)

            # SHAP用に数値変換
            try:
                if hasattr(X_transformed, 'toarray'):
                    X_transformed_for_shap = X_transformed.toarray()
                else:
                    X_transformed_for_shap = X_transformed
                
                X_transformed_for_shap = pd.DataFrame(X_transformed_for_shap).apply(pd.to_numeric, errors='coerce').fillna(0).values
            except (ValueError, TypeError):
                logger.warning("SHAP計算用のデータ数値変換に失敗しました。")
                X_transformed_for_shap = X_transformed

            # TreeExplainerを使用してSHAP値を計算
            explainer = shap.TreeExplainer(model)
            shap_values = explainer.shap_values(X_transformed_for_shap)
            base_values = explainer.expected_value

            # ★★★ 修正箇所: 分類モデルのSHAP値の構造（クラス次元）を適切に処理 ★★★
            # shap_values がリストの場合（各クラスごとの配列のリスト）
            if isinstance(shap_values, list):
                # 2値分類ならインデックス1（陽性クラス）、それ以外なら0を使用
                target_class_index = 1 if len(shap_values) > 1 else 0
                shap_values_for_plot = shap_values[target_class_index]
                
                if isinstance(base_values, list) or isinstance(base_values, np.ndarray):
                     base_values_for_plot = base_values[target_class_index]
                else:
                     base_values_for_plot = base_values

            # shap_values がndarrayの場合（[samples, features, classes] または [samples, features]）
            elif isinstance(shap_values, np.ndarray):
                # 3次元配列の場合（クラス次元がある場合）
                if shap_values.ndim == 3:
                    target_class_index = 1 if shap_values.shape[2] > 1 else 0
                    shap_values_for_plot = shap_values[:, :, target_class_index]
                    
                    if isinstance(base_values, (list, np.ndarray)) and len(base_values) > 1:
                        base_values_for_plot = base_values[target_class_index]
                    else:
                         base_values_for_plot = base_values
                else:
                    # 回帰や2値分類でクラス次元が省略されている場合
                    shap_values_for_plot = shap_values
                    base_values_for_plot = base_values
            else:
                # Explanationオブジェクトなどが返ってきた場合のフォールバック
                shap_values_for_plot = shap_values
                base_values_for_plot = base_values

            # 特徴量名を取得
            feature_names_out = preprocessor.get_feature_names_out()

            # データの次元を確認し、1サンプルの場合はフラットにする
            if len(shap_values_for_plot.shape) > 1:
                 shap_values_for_plot = shap_values_for_plot[0]
            
            if len(X_transformed_for_shap.shape) > 1:
                 data_for_plot = X_transformed_for_shap[0]
            else:
                 data_for_plot = X_transformed_for_shap

            # SHAP Explanationオブジェクトの作成
            shap_explanation = shap.Explanation(
                values=shap_values_for_plot,
                base_values=base_values_for_plot,
                data=data_for_plot,
                feature_names=feature_names_out
            )

            # ウォーターフォールプロットの生成
            plt.figure()
            shap.plots.waterfall(shap_explanation, max_display=15, show=False)
            plt.tight_layout()
            
            buf = io.BytesIO()
            plt.savefig(buf, format='png', bbox_inches='tight')
            plt.close()
            buf.seek(0)
            
            chart_blob_name = f"charts/{model_id}/shap_{uuid.uuid4().hex}.png"
            upload_bytes_to_blob(buf.getvalue(), chart_blob_name)
            shap_plot_url = get_blob_sas_url(chart_blob_name)

        except Exception as e_shap:
            logger.error(f"SHAPプロットの生成に失敗: {e_shap}", exc_info=True)

        return jsonify({
            "prediction": final_prediction,
            "shap_plot_url": shap_plot_url
        })

    except Exception as e:
        logger.error(f"モデルシミュレーション {model_id} 中にエラー: {e}", exc_info=True)
        return jsonify({"error": f"シミュレーションの実行中にサーバーエラーが発生しました: {str(e)}"}), 500


@data_analyzer_bp.route('/api/project/<project_id>', methods=['DELETE'])
def delete_project_api(project_id):
    """分析プロジェクトと、それに紐づくファイル、予測モデルをすべて削除します。"""
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    logger.info(f"ユーザー '{user_id}' がプロジェクト '{project_id}' の削除を要求しました。")

    try:
        # 1. 削除対象のプロジェクト情報を取得
        project_data = get_analysis_project(project_id, user_id)
        if not project_data:
            logger.warning(f"削除対象のプロジェクト '{project_id}' が見つからないか、アクセス権がありません。")
            return jsonify({"error": "削除対象のプロジェクトが見つからないか、アクセス権がありません。"}), 404

        # 2. 関連する予測モデルを検索・削除
        related_models = get_models_by_project_id(project_id, user_id)
        for model in related_models:
            model_id = model.get('id')
            if not model_id: continue
            
            logger.info(f"関連モデル '{model_id}' を削除します。")
            # モデルの成果物（パイプライン、グラフ）をストレージから削除
            if pipeline_blob := model.get('pipeline_blob_name'):
                delete_blob(pipeline_blob)
            
            if chart_gcs_uri := model.get('feature_importance_chart_url_gcs'):
                # URIを直接渡すだけでOK
                _delete_from_gcs(chart_gcs_uri)
            
            # モデル本体をDBから削除
            delete_design_project(model_id, user_id)

        # 3. プロジェクトに紐づくストレージ上のファイルを削除
        gcs_bucket_name = os.environ.get("GOOGLE_CLOUD_STORAGE_BUCKET_NAME")
        
        # insights内のグラフ (GCS)
        insights_data = project_data.get('insights', {})
        if isinstance(insights_data, dict):
            # 新しい階層化データ構造に対応
            all_insight_lists = [
                group.get("insights", []) for group in insights_data.get("primary_factor_insights", [])
            ] + [
                group.get("insights", []) for group in insights_data.get("categorized_insights", [])
            ] + [
                insights_data.get("anomaly_alerts", [])
            ]
            for insight_list in all_insight_lists:
                if isinstance(insight_list, list):
                    for insight in insight_list:
                        if insight and isinstance(insight.get('charts'), list):
                            for chart_uri in insight['charts']:
                                # URIを直接渡すだけでOK
                                _delete_from_gcs(chart_uri)

        # クラスター分析のグラフ (GCS)
        cluster_result = project_data.get('clusterAnalysisResult', {})
        if isinstance(cluster_result, dict):
            for key in ['elbow_chart_url', 'cluster_chart_url']:
                if gcs_uri := cluster_result.get(key):
                    # GCS URIであることを確認
                    if gcs_uri and 'blob.core.windows.net' not in gcs_uri:
                       # URIを直接渡すだけでOK
                       _delete_from_gcs(gcs_uri)

        # 各種データファイル (Blob Storage)
        if data_file := project_data.get('dataFile', {}):
            if blob_name := data_file.get('blob_name'): delete_blob(blob_name)
        for ref_file in project_data.get('referenceFiles', []):
            if blob_name := ref_file.get('blob_name'): delete_blob(blob_name)
        if processed_blob := project_data.get('processedDataBlobName'):
            delete_blob(processed_blob)

        # 4. プロジェクト本体をDBから削除
        if delete_design_project(project_id, user_id):
            logger.info(f"プロジェクト '{project_id}' と関連リソースの削除が完了しました。")
            return jsonify({"success": True, "message": "プロジェクトが正常に削除されました。"}), 200
        else:
            raise Exception("データベースからのプロジェクト本体の削除に失敗しました。")

    except Exception as e:
        logger.error(f"プロジェクト {project_id} の削除中にエラーが発生しました: {e}", exc_info=True)
        return jsonify({"error": "プロジェクトの削除中にサーバーエラーが発生しました。"}), 500

def _run_timeseries_analysis(df: pd.DataFrame, project_id: str, user_id: str, project_name: str) -> dict | None:
    """
    DataFrameに時系列データが含まれている場合、SARIMAモデルを構築・保存し、
    モデル管理用のメタデータ辞書を返す。(改訂版 v2)
    エラーハンドリングと変数のスコープを修正。
    """
    logger.info(f"プロジェクト {project_id} の時系列予測モデル構築を開始。")

    try:
        # 1. 分析対象となるカラムを特定
        datetime_cols = df.select_dtypes(include=['datetime64[ns]', 'datetime64[ns, UTC]']).columns.tolist()
        if not datetime_cols:
            logger.info("分析可能な日付/時刻型の列がないため、時系列予測をスキップします。")
            return None
        
        date_col = datetime_cols[0]
        numeric_cols = [col for col in df.select_dtypes(include=np.number).columns if col != date_col]
        if not numeric_cols:
            logger.info(f"日付列 '{date_col}' に対応する数値列がないため、時系列予測をスキップします。")
            return None
        value_col = numeric_cols[0]

        logger.info(f"時系列予測対象: 日付='{date_col}', 値='{value_col}'")

        # 2. データの前処理
        ts_df = df[[date_col, value_col]].copy()
        ts_df = ts_df.dropna().sort_values(by=date_col)
        ts_df = ts_df.set_index(date_col)

        ts_daily = ts_df[value_col].resample('D').mean().interpolate()
        
        if len(ts_daily) < 14:
            logger.warning("時系列データが短すぎるため、予測をスキップします。")
            return None

        # 3. 訓練データとテストデータに分割
        train_size = int(len(ts_daily) * 0.8)
        train, test = ts_daily[0:train_size], ts_daily[train_size:]

        # 4. SARIMAモデルの学習
        model = SARIMAX(train, order=(1, 1, 1), seasonal_order=(1, 1, 1, 7))
        model_fit = model.fit(disp=False)

        # 5. 予測の実行
        forecast_steps = len(test) + 30
        forecast = model_fit.get_forecast(steps=forecast_steps)
        forecast_values = forecast.predicted_mean
        forecast_ci = forecast.conf_int()

        # 6. 評価指標の計算
        mape = mean_absolute_percentage_error(test, forecast_values[:len(test)])

        # 7. グラフの生成とアップロード
        plt.figure(figsize=(12, 6))
        plt.plot(ts_daily.index, ts_daily, label='実績値', color='gray', alpha=0.7)
        plt.plot(forecast_values.index, forecast_values, label='予測値', color='blue')
        plt.fill_between(forecast_ci.index, forecast_ci.iloc[:, 0], forecast_ci.iloc[:, 1], color='blue', alpha=0.1, label='95%信頼区間')
        plt.title(f'時系列予測: {value_col}')
        plt.xlabel(date_col)
        plt.ylabel(value_col)
        plt.legend()
        plt.grid(True, linestyle='--', alpha=0.6)
        plt.tight_layout()
        
        buf = io.BytesIO()
        plt.savefig(buf, format='png')
        plt.close()
        buf.seek(0)
        
        chart_filename = f"charts/{project_id}/timeseries_forecast_{uuid.uuid4().hex}.png"
        gcs_bucket_name = os.environ.get("GOOGLE_CLOUD_STORAGE_BUCKET_NAME")
        gcs_uri, forecast_chart_url = None, None
        
        # GCSへのアップロード (GCSが有効な場合)
        if gcs_bucket_name:
            gcs_uri = _upload_to_gcs(buf, chart_filename, gcs_bucket_name)
            if gcs_uri:
                forecast_chart_url = get_gcs_signed_url(chart_filename, bucket_name=gcs_bucket_name)
        
        # Azure Blob Storageへのアップロード (フォールバックまたはメイン)
        # GCSが無効、または失敗した場合はAzure Blobを使用
        if not forecast_chart_url:
             upload_bytes_to_blob(buf.getvalue(), chart_filename)
             forecast_chart_url = get_blob_sas_url(chart_filename)

        # 8. モデルオブジェクトの保存
        pipeline_blob_name = None
        try:
            pipeline_buffer = io.BytesIO()
            joblib.dump(model_fit, pipeline_buffer)
            pipeline_buffer.seek(0)
            pipeline_blob_name = f"models/{project_id}/sarima_v1_{uuid.uuid4().hex[:8]}.joblib"
            upload_bytes_to_blob(pipeline_buffer.getvalue(), pipeline_blob_name)
            logger.info(f"学習済みSARIMAモデルをBLOBに保存しました: {pipeline_blob_name}")
        except Exception as e_save:
            logger.error(f"SARIMAモデルの保存に失敗: {e_save}", exc_info=True)

        # 9. AIによる解説生成
        commentary = f"AIによる自動生成された解説: 過去のデータに基づき、今後30日間の「{value_col}」を予測しました。モデルは週単位の季節性を考慮したSARIMAモデルです。テストデータに対する予測精度（MAPE）は **{mape:.1%}** であり、平均的に約{mape:.1%}の誤差で予測できています。"

        # 10. 再現用コードの生成
        # ★重要: date_col, value_col が確実に存在し、処理が成功したこのタイミングでのみ生成する
        reproduction_code = f"""# --- 時系列予測モデル (SARIMA) 構築コード（再現用） ---
import pandas as pd
from statsmodels.tsa.statespace.sarimax import SARIMAX
import matplotlib.pyplot as plt

# 設定
date_col = '{date_col}'
value_col = '{value_col}'

# 1. データ準備
# 注意: df は分析に使用された元のデータフレームです
df_ts = df[[date_col, value_col]].dropna().sort_values(by=date_col).set_index(date_col)

# 日次リサンプリング・補完
# ※データの頻度に応じて 'D' (日次), 'W' (週次), 'M' (月次) などを調整してください
ts_resampled = df_ts[value_col].resample('D').mean().interpolate()

# 2. 学習データ分割 (直近20%をテストとする)
train_size = int(len(ts_resampled) * 0.8)
train, test = ts_resampled[0:train_size], ts_resampled[train_size:]

# 3. SARIMAモデル学習
# パラメータ order=(p,d,q), seasonal_order=(P,D,Q,s) は自動設定されたものです
model = SARIMAX(train, order=(1, 1, 1), seasonal_order=(1, 1, 1, 7))
model_fit = model.fit(disp=False)

# 4. 予測
forecast = model_fit.get_forecast(steps=len(test) + 30)
pred_mean = forecast.predicted_mean
conf_int = forecast.conf_int()

# 5. 可視化
plt.figure(figsize=(12, 6))
plt.plot(ts_resampled.index, ts_resampled, label='Observed')
plt.plot(pred_mean.index, pred_mean, label='Forecast', color='red')
plt.fill_between(conf_int.index, conf_int.iloc[:, 0], conf_int.iloc[:, 1], color='pink', alpha=0.3)
plt.title(f'Forecast for {{value_col}}')
plt.legend()
plt.show()
"""

        # 11. モデル管理用の完全なメタデータを作成
        model_id = f"model_{project_id}_{uuid.uuid4().hex[:8]}"
        result = {
            "id": model_id,
            "model_id": model_id,
            "project_id": project_id,
            "userId": user_id,
            "projectName": project_name,
            "target_variable": value_col,
            "model_type": "timeseries",
            "algorithm": "SARIMA",
            "accuracy_score_label": "予測精度 (MAPE)",
            "accuracy_score": mape,
            "commentary": commentary,
            "feature_importance_chart_url": forecast_chart_url,
            "feature_importance_chart_url_gcs": gcs_uri,
            "pipeline_blob_name": pipeline_blob_name,
            "createdAt": datetime.now(timezone.utc).isoformat(),
            "version": 1,
            "error": None,
            "code": reproduction_code, # ★コードを含める
            "type": "PredictionModel",
        }
        return result

    except Exception as e:
        logger.error(f"時系列予測タスク {project_id} でエラー: {e}", exc_info=True)
        return None


@data_analyzer_bp.route('/api/models/<model_id>/predict_batch', methods=['POST'])
def predict_batch_api(model_id):
    """
    アップロードされたファイルデータに対し、指定されたモデルで一括予測を実行し、
    予測結果を付与したCSVファイルを返す。
    One-Hotエンコーディングされるカテゴリ変数も元の形式で受け付けるように改訂。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    if 'predict_file' not in request.files:
        return jsonify({"error": "予測対象のデータファイルがありません。"}), 400
    
    predict_file = request.files['predict_file']

    try:
        # 1. モデルのメタデータを取得
        model_data = get_analysis_project(project_id=model_id, user_id=user_id)
        if not model_data or not model_data.get('pipeline_blob_name'):
            return jsonify({"error": "モデルが見つからないか、学習済みパイプラインが保存されていません。"}), 404

        # 2. 学習済みパイプラインをストレージからロード
        pipeline_blob_name = model_data['pipeline_blob_name']
        pipeline_bytes = download_blob_to_bytes(pipeline_blob_name)
        if not pipeline_bytes:
            return jsonify({"error": "パイプラインファイルの読み込みに失敗しました。"}), 500
        
        pipeline = joblib.load(io.BytesIO(pipeline_bytes))

        # 3. アップロードされた予測対象ファイルをDataFrameとして読み込み
        predict_data_bytes = predict_file.read()
        df_to_predict_original = _load_df_from_bytes_with_temp_file(predict_data_bytes, predict_file.filename)
        if df_to_predict_original is None:
            return jsonify({"error": "予測対象ファイルの読み込みに失敗しました。対応形式（CSV, Excel）か確認してください。"}), 400

        # 4. 予測用データをモデルが期待するカラム構成に整形
        # パイプラインから学習時の特徴量名リストを取得
        model_feature_names = pipeline.named_steps['preprocessor'].feature_names_in_
        
        # 学習時のカラム順序と構成に合わせる
        # これにより、アップロードされたファイルのカラム順が違っても、モデルに必要なカラムがなくてもエラーなく処理できます
        df_for_prediction = df_to_predict_original.reindex(columns=model_feature_names)
        
        # 5. 予測を実行
        # パイプラインが内部でOneHotEncodingや欠損値補完を自動的に適用します
        predictions = pipeline.predict(df_for_prediction)

        # 6. 【改訂】元のデータに予測結果の列を追加
        # カラム構成を変更する前のオリジナルのDataFrameに結果を結合します
        target_variable_name = model_data.get('target_variable', 'prediction')
        df_result = df_to_predict_original.copy()
        df_result[f"predicted_{target_variable_name}"] = predictions

        # 7. 結果をCSV形式でメモリ上のバッファに書き出す
        output_buffer = io.BytesIO()
        df_result.to_csv(output_buffer, index=False, encoding='utf-8-sig')
        output_buffer.seek(0)
        
        # 8. ダウンロード用のファイル名を生成
        base, ext = os.path.splitext(predict_file.filename)
        download_name = f"{base}_predicted.csv"

        # 9. CSVファイルをレスポンスとして返す
        return send_file(
            output_buffer,
            mimetype='text/csv',
            as_attachment=True,
            download_name=download_name
        )

    except Exception as e:
        logger.error(f"モデル {model_id} による一括予測中にエラー: {e}", exc_info=True)
        error_message = f"予測の実行中にサーバーエラーが発生しました: {str(e)}"
        return jsonify({"error": error_message}), 500

@data_analyzer_bp.route('/api/models/<model_id>/retrain', methods=['POST'])
def retrain_model_api(model_id):
    """
    指定されたモデルを、新しくアップロードされたデータで再学習する非同期タスクを開始する。
    """
    user_info = get_entra_user_info(request)
    user_id = user_info.get("userId", "local-user")

    if 'retrain_data_file' not in request.files:
        return jsonify({"error": "再学習用のデータファイルがありません。"}), 400
    
    retrain_file = request.files['retrain_data_file']
    
    try:
        # 1. 再学習用データを一時BLOBにアップロード
        retrain_data_bytes = retrain_file.read()
        temp_blob_name = f"temp_retrain_data/{model_id}/{uuid.uuid4().hex}/{retrain_file.filename}"
        upload_bytes_to_blob(retrain_data_bytes, temp_blob_name)

        # 2. 再学習タスクを作成・登録
        task_id = f"task_retrain_{model_id}_{uuid.uuid4().hex[:8]}"
        task_data = {
            "id": task_id,
            "task_id": task_id,
            "status": "PENDING",
            "task_type": "model_retraining",
            "model_id": model_id,
            "user_id": user_id,
            "new_data_blob_name": temp_blob_name,
            "new_data_filename": retrain_file.filename
        }
        create_task_in_cosmos(task_data)

        # 3. 非同期でタスクを実行
        executor.submit(_run_retraining_task, task_id, model_id, user_id, temp_blob_name, retrain_file.filename)

        return jsonify({"success": True, "message": "モデルの更新（再学習）タスクを開始しました。", "task_id": task_id}), 202

    except Exception as e:
        logger.error(f"モデル {model_id} の再学習タスク開始中にエラー: {e}", exc_info=True)
        return jsonify({"error": "再学習タスクの開始に失敗しました。"}), 500

def _run_retraining_task(task_id, model_id, user_id, new_data_blob_name, new_data_filename):
    """
    モデルの再学習を実際に行う非同期タスク。
    """
    try:
        update_task_in_cosmos(task_id, {"status": "PROCESSING", "status_message": "既存モデルの情報を読み込み中...", "progress_percent": 10})

        # 1. 更新対象の既存モデル情報を取得
        original_model_data_raw = get_analysis_project(project_id=model_id, user_id=user_id)
        if not original_model_data_raw:
            raise Exception(f"更新対象のモデル(ID: {model_id})が見つかりませんでした。")
        original_model_data = copy.deepcopy(original_model_data_raw)
        
        target_variable = original_model_data.get('target_variable')
        model_type = original_model_data.get('model_type')
        project_id = original_model_data.get('project_id')
        
        if not all([target_variable, model_type, project_id]):
            raise Exception("モデル情報（目的変数、モデルタイプ、プロジェクトID）が不完全です。")

        update_task_in_cosmos(task_id, {"status_message": "新しい学習データを読み込み中...", "progress_percent": 25})

        # 2. 新しい学習データを読み込み
        data_bytes = download_blob_to_bytes(new_data_blob_name)
        df_new = _load_df_from_bytes_with_temp_file(data_bytes, new_data_filename)
        if df_new is None:
            raise Exception("新しい学習データファイルの読み込みに失敗しました。")
        
        df_new.attrs['name'] = original_model_data.get('projectName', project_id)

        update_task_in_cosmos(task_id, {"status_message": "新しいデータでモデルを再学習中...", "progress_percent": 50})

        # 3. 新しいモデルを学習
        new_model_meta, _ = _run_prediction_task(df_new, target_variable, model_type, project_id)
        if not new_model_meta or new_model_meta.get('error'):
            raise Exception(f"モデルの再学習に失敗しました: {new_model_meta.get('error', '不明なエラー')}")

        update_task_in_cosmos(task_id, {"status_message": "モデル情報を更新中...", "progress_percent": 80})

        # 4. 既存のモデル情報を更新
        updates_for_db = {
            "accuracy_score": new_model_meta.get("accuracy_score"), "error_metric": new_model_meta.get("error_metric"),
            "feature_details": new_model_meta.get("feature_details"),
            "commentary": f"【更新】{datetime.now(timezone.utc).strftime('%Y-%m-%d')}に新しいデータで再学習。{new_model_meta.get('commentary', '')}",
            "feature_importance_chart_url": new_model_meta.get("feature_importance_chart_url"),
            "feature_importance_chart_url_gcs": new_model_meta.get("feature_importance_chart_url_gcs"),
            "pipeline_blob_name": new_model_meta.get("pipeline_blob_name"),
            "version": original_model_data.get("version", 1) + 1,
            "updatedAt": datetime.now(timezone.utc).isoformat()
        }
        update_analysis_project(model_id, user_id, updates_for_db)

        # 5. 古い成果物をクリーンアップ
        if old_pipeline_blob := original_model_data.get('pipeline_blob_name'):
            delete_blob(old_pipeline_blob)
        if old_chart_gcs := original_model_data.get('feature_importance_chart_url_gcs'):
            _delete_from_gcs(old_chart_gcs)
        
        update_task_in_cosmos(task_id, {"status": "SUCCESS", "status_message": "モデルの更新が完了しました。", "progress_percent": 100})

    except Exception as e:
        logger.error(f"モデル再学習タスク {task_id} でエラー: {e}", exc_info=True)
        update_task_in_cosmos(task_id, {"status": "FAILURE", "error": str(e), "status_message": "モデルの更新中にエラーが発生しました。"})
    
    finally:
        try:
            delete_blob(new_data_blob_name)
        except Exception as e_del:
            logger.error(f"一時ファイル {new_data_blob_name} の削除に失敗: {e_del}", exc_info=True)


